from __future__ import annotations

import hashlib
import importlib.util
import json
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SCHEMA_VERSION = "6.3"
RUNTIME_REVISION = "6.3.1"
GEOMETRY_TOLERANCE_IN = 0.03


def _load(name: str):
    path = Path(__file__).with_name(f"{name}.py")
    spec = importlib.util.spec_from_file_location(f"v63_acceptance_{name}", path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def _read_json(path: Path) -> dict[str, Any] | list[Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _issue(code: str, slide_id: str, message: str) -> dict[str, str]:
    return {
        "code": code,
        "severity": "blocker",
        "stage": "v63_pptx_acceptance",
        "slide_id": slide_id,
        "message": message,
    }


def _write_json_atomic(path: Path, value: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_suffix(path.suffix + ".tmp")
    temporary.write_text(
        json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    temporary.replace(path)


def _bbox_in(shape) -> list[float]:
    return [
        float(shape.left) / 914400.0,
        float(shape.top) / 914400.0,
        float(shape.width) / 914400.0,
        float(shape.height) / 914400.0,
    ]


def _geometry_matches(actual: list[float], expected: list[float]) -> bool:
    return all(
        abs(float(actual_value) - float(expected_value)) <= GEOMETRY_TOLERANCE_IN
        for actual_value, expected_value in zip(actual, expected)
    )


def _is_picture(shape) -> bool:
    return shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}


def _picture_asset_errors(shape, asset: dict | None, project: Path) -> list[str]:
    errors = []
    if not isinstance(asset, dict):
        return ['V63_PPTX_ASSET_LEDGER_MISSING']
    path = (project / str(asset.get('asset_path', ''))).resolve()
    expected = asset.get('asset_sha256')
    if not path.is_file() or _sha256(path) != expected or hashlib.sha256(shape.image.blob).hexdigest() != expected:
        errors.append('V63_PPTX_ASSET_HASH_MISMATCH')
    from pptx.oxml.ns import qn
    line = shape._element.spPr.find(qn('a:ln'))
    if line is not None and line.find(qn('a:noFill')) is None:
        errors.append('V63_PPTX_PICTURE_OUTLINE')
    for effects in shape._element.spPr.findall(qn('a:effectLst')):
        if len(effects):
            errors.append('V63_PPTX_PICTURE_EFFECT')
    return errors


def _expected_text(element: dict[str, Any]) -> str:
    runs = element.get("runs")
    if isinstance(runs, list):
        return "".join(str(run.get("text", "")) for run in runs if isinstance(run, dict))
    return str(element.get("text", ""))


def _skeleton_values(slide: dict[str, Any], page_number: int) -> dict[str, str]:
    points = slide.get("core_points", [])
    if isinstance(points, str):
        points = [points]
    normalized = [
        re.sub(r"^(?:[■▪▫•●□]\s*)+", "", str(point).strip()).strip()
        for point in points
        if str(point).strip()
    ]
    return {
        "chapter": str(slide.get("chapter", "")),
        "page_title": str(slide.get("title", slide.get("page_title", ""))),
        "core_judgment": "\n".join(normalized),
        "source": str(slide.get("source", "")),
        "page_number": str(slide.get("page_number", page_number)),
    }


def _normalized_text(value: Any) -> str:
    return str(value).replace("\r\n", "\n").replace("\r", "\n").strip()


def visual_bindings(project_dir: str | Path, pptx_path: str | Path) -> dict[str, str | None]:
    project = Path(project_dir).resolve()
    fixed = ['project_brief.json', '.build/slides.json', '.build/formal_blueprint_manifest.json',
             '.build/v63_source_body_rois.json', '.build/v63_visual_review_tiles.json',
             '.build/v63_visual_census.json', '.build/v63_scene_graph.json', '.build/v63_asset_ledger.json']
    fixed.extend(p.relative_to(project).as_posix() for p in (project / 'blueprints').glob('S*.png'))
    fixed.extend(p.relative_to(project).as_posix() for p in (project / '.build/rendered/current').glob('S*.png'))
    fixed.extend(p.relative_to(project).as_posix() for p in (project / '.build/assets').rglob('*.png'))
    result = {relative: _sha256(project / relative) if (project / relative).is_file() else None for relative in fixed}
    result['pptx'] = _sha256(Path(pptx_path))
    result['template'] = _sha256(Path(__file__).parents[1] / 'assets/company_template.pptx')
    return result


def read_bound_visual_review(project_dir: str | Path, pptx_path: str | Path) -> dict | None:
    project = Path(project_dir).resolve()
    path = project / '.build/v63_visual_review.json'
    if not path.is_file():
        return None
    value = _read_json(path)
    if value.get('bindings') != visual_bindings(project, pptx_path):
        # A new rendered revision naturally invalidates the previous review.
        # Treat it as missing evidence, not a COM/build failure or permission
        # to reuse the old verdict. Pure finalization can resume without build 3.
        return None
    graph = _read_json(project / '.build/v63_scene_graph.json')
    pages = value.get('pages', [])
    if value.get('reviewed') is not True or {p.get('slide_id') for p in pages} != set(graph['pages']):
        raise ValueError('V63_VISUAL_REVIEW_INCOMPLETE')
    for page in pages:
        checks = page.get('object_checks', [])
        actual = {c.get('candidate_id') for c in checks}
        census = _read_json(project / '.build/v63_visual_census.json')
        expected = {c['candidate_id'] for c in census['pages'][page['slide_id']]['candidates']}
        if actual != expected or any(c.get('status') not in {'present', 'difference', 'missing'} or not c.get('evidence') for c in checks):
            raise ValueError('V63_VISUAL_OBJECT_REVIEW_INCOMPLETE')
        if not isinstance(page.get('differences'), list):
            raise ValueError('V63_VISUAL_REVIEW_INCOMPLETE')
        if any(c['status'] != 'present' for c in checks) and not page['differences']:
            raise ValueError('V63_VISUAL_REVIEW_CONTRADICTORY')
    return value


def audit_v63_pptx(
    pptx_path: str | Path,
    project_dir: str | Path,
    *,
    template_path: str | Path,
) -> dict[str, Any]:
    """Prove that the compiled deck preserves the skeleton and every scene atom."""

    pptx = Path(pptx_path).resolve()
    project = Path(project_dir).resolve()
    blockers: list[dict[str, str]] = []
    if not pptx.is_file():
        blockers.append(_issue("V63_PPTX_MISSING", "", str(pptx)))
        return {
            "schema_version": SCHEMA_VERSION,
            "deconstruction_runtime_revision": RUNTIME_REVISION,
            "ok": False,
            "status": "blocked",
            "editable_body_count": 0,
            "image_count": 0,
            "blockers": blockers,
        }

    scene_graph = _read_json(project / ".build" / "v63_scene_graph.json")
    asset_path = project / '.build/v63_asset_ledger.json'
    ledger = _read_json(asset_path) if asset_path.is_file() else {}
    assets_by_id = {item['asset_id']: item for item in ledger.get('assets', [])}
    skeleton = _load("v63_skeleton_contract")
    scene = _load("v63_scene_graph")
    try:
        skeleton_audit = skeleton.audit_pptx_skeleton(template_path, pptx)
    except Exception as exc:
        skeleton_audit = {"ok": False, "errors": [str(exc)]}
    if not skeleton_audit.get("ok"):
        blockers.append(
            _issue(
                "V63_PPTX_SKELETON_CHANGED",
                "",
                "; ".join(str(item) for item in skeleton_audit.get("errors", []))
                or "master-owned placeholder contract changed",
            )
        )

    presentation = Presentation(str(pptx))
    page_ids = sorted(scene_graph.get("pages", {}))
    slide_specs = _read_json(project / ".build" / "slides.json")
    specs_by_id = {
        str(item.get("slide_id")): item
        for item in slide_specs
        if isinstance(slide_specs, list) and isinstance(item, dict)
    }
    if len(presentation.slides) != len(page_ids):
        blockers.append(
            _issue(
                "V63_PPTX_PAGE_COUNT_MISMATCH",
                "",
                f"expected {len(page_ids)}, got {len(presentation.slides)}",
            )
        )

    body_roi_in = skeleton.read_template_contract(template_path)["body_roi_in"]
    expected_image_names: set[str] = set()
    editable_body_count = 0
    image_count = 0
    for page_index, slide_id in enumerate(page_ids):
        if page_index >= len(presentation.slides):
            break
        slide = presentation.slides[page_index]
        try:
            placeholder_shapes = skeleton.resolve_python_pptx_shapes(slide)
            expected_skeleton = _skeleton_values(
                specs_by_id.get(slide_id, {}), page_index + 1
            )
            for role, expected_value in expected_skeleton.items():
                actual_value = placeholder_shapes[role].text
                if _normalized_text(actual_value) != _normalized_text(expected_value):
                    blockers.append(
                        _issue(
                            "V63_PPTX_SKELETON_TEXT_MISMATCH",
                            slide_id,
                            f"{role}: placeholder text differs from slides.json",
                        )
                    )
        except Exception as exc:
            blockers.append(
                _issue("V63_PPTX_SKELETON_TEXT_MISMATCH", slide_id, str(exc))
            )
        shapes_by_name: dict[str, list[Any]] = {}
        for shape in slide.shapes:
            shapes_by_name.setdefault(str(shape.name), []).append(shape)
        page = scene_graph["pages"][slide_id]
        for element in page.get("elements", []):
            if not isinstance(element, dict) or element.get("type") == "group":
                continue
            element_id = str(element.get("element_id", ""))
            shape_name = f"V63_{element_id}"
            matches = shapes_by_name.get(shape_name, [])
            if len(matches) != 1:
                code = (
                    "V63_PPTX_ELEMENT_MISSING"
                    if not matches
                    else "V63_PPTX_ELEMENT_DUPLICATED"
                )
                blockers.append(_issue(code, slide_id, f"{element_id}: found {len(matches)}"))
                continue
            shape = matches[0]
            kind = str(element.get("type"))
            normalized = scene.normalize_element_geometry(element)
            expected_box = scene.pixel_box_to_slide_box(
                normalized["bbox_px"], page["body_roi_px"], body_roi_in,
                coordinate_mode=page.get('coordinate_mode', 'legacy_stretch'),
                allow_line=kind in {'line', 'connector', 'arrow'} or not element.get('closed', True),
            )
            if not _geometry_matches(_bbox_in(shape), expected_box):
                blockers.append(
                    _issue(
                        "V63_PPTX_GEOMETRY_MISMATCH",
                        slide_id,
                        f"{element_id}: compiled box differs from scene graph",
                    )
                )
            if kind == "image_crop":
                expected_image_names.add(shape_name)
                image_count += 1
                if not _is_picture(shape):
                    blockers.append(
                        _issue(
                            "V63_PPTX_CROP_NOT_IMAGE", slide_id, f"{element_id}: expected picture"
                        )
                    )
                else:
                    for code in _picture_asset_errors(shape, assets_by_id.get(element.get('asset_id')), project):
                        blockers.append(_issue(code, slide_id, element_id))
            else:
                editable_body_count += 1
                if _is_picture(shape):
                    blockers.append(
                        _issue(
                            "V63_PPTX_EDITABLE_ATOM_RASTERIZED",
                            slide_id,
                            f"{element_id}: editable atom compiled as picture",
                        )
                    )
                if kind == "text" and _expected_text(element).replace("\r", "\n") != str(shape.text).replace("\r", "\n"):
                    blockers.append(
                        _issue(
                            "V63_PPTX_TEXT_MISMATCH",
                            slide_id,
                            f"{element_id}: compiled text differs from scene graph",
                        )
                    )
        for shape in slide.shapes:
            if _is_picture(shape) and str(shape.name).startswith("V63_") and str(shape.name) not in expected_image_names:
                blockers.append(
                    _issue(
                        "V63_PPTX_UNDECLARED_RASTER",
                        slide_id,
                        f"{shape.name}: picture is not an approved local crop",
                    )
                )

    return {
        "schema_version": SCHEMA_VERSION,
        "deconstruction_runtime_revision": RUNTIME_REVISION,
        "ok": not blockers,
        "status": "pass" if not blockers else "blocked",
        "editable_body_count": editable_body_count,
        "image_count": image_count,
        "blockers": blockers,
    }


def evaluate_v63_acceptance(
    project_dir: str | Path,
    pptx_path: str | Path,
    audit: dict[str, Any],
    render_result: dict[str, Any],
    visual_delta: dict[str, Any] | None = None,
) -> dict[str, Any]:
    project = Path(project_dir).resolve()
    pptx = Path(pptx_path).resolve()
    blockers = list(audit.get("blockers", []))
    render_valid = (
        render_result.get("ok") is True
        and render_result.get("visual_verification") is True
        and render_result.get("status") in {"pass", "pass_with_warnings"}
    )
    if not render_valid:
        blockers.append(
            _issue(
                "V63_RENDER_VERIFICATION_FAILED",
                "",
                "native render did not complete with visual verification",
            )
        )
    warnings = []
    if isinstance(visual_delta, dict):
        warnings = list(visual_delta.get("warnings", []))
        if visual_delta.get('action') not in {'deliver', 'deliver_with_warnings'}:
            blockers.append(_issue('V63_VISUAL_REVIEW_NOT_ACCEPTED', '', str(visual_delta.get('action'))))
    result = {
        "schema_version": SCHEMA_VERSION,
        "deconstruction_runtime_revision": RUNTIME_REVISION,
        "accepted": not blockers,
        "status": "pass" if not blockers and not warnings else (
            "pass_with_warnings" if not blockers else "blocked"
        ),
        "pptx": str(pptx),
        "pptx_sha256": _sha256(pptx) if pptx.is_file() else None,
        "audit": audit,
        "render": render_result,
        "warnings": warnings,
        "blockers": blockers,
    }
    if isinstance(visual_delta, dict) and visual_delta.get('review_required'):
        result['visual_acceptance_passed'] = visual_delta.get('visual_acceptance_passed', False)
        result['bindings'] = visual_bindings(project, pptx)
        result['visual_review_sha256'] = _sha256(project / '.build/v63_visual_review.json') if (project / '.build/v63_visual_review.json').is_file() else None
    _write_json_atomic(project / ".build" / "deconstruction_acceptance.json", result)
    return result


def locked_v63_acceptance(
    project_dir: str | Path, pptx_path: str | Path
) -> dict[str, Any] | None:
    """Return an accepted V6.3 result only while the PPTX bytes remain unchanged."""

    project = Path(project_dir).resolve()
    pptx = Path(pptx_path).resolve()
    path = project / ".build" / "deconstruction_acceptance.json"
    if not path.is_file() or not pptx.is_file():
        return None
    try:
        value = _read_json(path)
    except (OSError, ValueError, json.JSONDecodeError):
        return None
    if not isinstance(value, dict):
        return None
    if (
        value.get("schema_version") != SCHEMA_VERSION
        or value.get("deconstruction_runtime_revision") != RUNTIME_REVISION
        or value.get("accepted") is not True
        or value.get("pptx_sha256") != _sha256(pptx)
    ):
        return None
    if 'bindings' in value:
        if value['bindings'] != visual_bindings(project, pptx):
            return None
        review_path = project / '.build/v63_visual_review.json'
        if not review_path.is_file() or value.get('visual_review_sha256') != _sha256(review_path):
            return None
    return value


def finalize_visual_review(project_dir: str | Path, pptx_path: str | Path, *, build_attempt: int,
                           audit: dict, render_result: dict, fidelity_report: dict | None = None) -> dict:
    """Finalize an already-rendered deck. This function never launches a renderer."""
    project = Path(project_dir).resolve()
    review = read_bound_visual_review(project, pptx_path)
    delta_module = _load('v63_visual_delta')
    delta = delta_module.evaluate_visual_delta(build_attempt=build_attempt, structural_ok=bool(audit.get('ok')),
        fidelity_report=fidelity_report, visual_review=review, require_visual_review=True)
    delta['review_required'] = True
    delta_module.write_visual_delta(project, delta)
    acceptance = evaluate_v63_acceptance(project, pptx_path, audit, render_result, delta)
    return {'delta': delta, 'acceptance': acceptance}
