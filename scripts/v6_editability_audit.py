from __future__ import annotations

import hashlib
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from PIL import Image


SCHEMA_VERSION = "6.0"
DECONSTRUCTION_EDITABILITY_FAILED = "DECONSTRUCTION_EDITABILITY_FAILED"
BITMAP_CONTRACT_INVALID = "BITMAP_CONTRACT_INVALID"
_EMU_PER_INCH = 914400.0
_TOLERANCE = 0.02
_PICTURE, _LINE, _TEXTBOX = 13, 9, 17
_CHART_TYPES = {"hbar_chart", "column_chart", "line_chart", "combo_chart", "donut_chart", "grouped_hbar_chart"}
_SKELETON = {"SKEL_CHAPTER", "SKEL_TITLE", "SKEL_CORE", "SKEL_SOURCE", "SKEL_PAGE_NUMBER"}
_BASIC_TYPES = {"rect", "oval", "line", "arrow", "flow"}


def _pages(value: Any) -> dict[str, Any]:
    return value.get("pages", {}) if isinstance(value, dict) and isinstance(value.get("pages"), dict) else (value if isinstance(value, dict) else {})


def _issue(code: str, slide_id: str, message: str) -> dict[str, Any]:
    return {"code": code, "severity": "blocker", "stage": "pptx_editability", "slide_id": slide_id, "message": message}


def _box(shape) -> tuple[float, float, float, float]:
    return tuple(float(value) / _EMU_PER_INCH for value in (shape.left, shape.top, shape.width, shape.height))


def _overlap(a, b) -> float:
    left, top = max(a[0], b[0]), max(a[1], b[1])
    right, bottom = min(a[0] + a[2], b[0] + b[2]), min(a[1] + a[3], b[1] + b[3])
    return max(0.0, right - left) * max(0.0, bottom - top)


def _prefix(element_id: str) -> str:
    return f"EL_{element_id}_"


def _matching(slide, element_id: str):
    return [shape for shape in slide.shapes if shape.name.startswith(_prefix(element_id))]


def _body(slide) -> tuple[float, float, float, float]:
    core = next((shape for shape in slide.shapes if shape.name == "SKEL_CORE"), None)
    source = next((shape for shape in slide.shapes if shape.name == "SKEL_SOURCE"), None)
    if core is None or source is None:
        return (0.0, 0.0, 0.0, 0.0)
    left, _, width, _ = _box(core)
    top = (float(core.top) + float(core.height)) / _EMU_PER_INCH + .12
    bottom = float(source.top) / _EMU_PER_INCH - .195
    return (left, top, width, max(0.0, bottom - top))


def _xml_text(path: str | Path) -> dict[str, str]:
    result: dict[str, str] = {}
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist():
            if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                root = ElementTree.fromstring(archive.read(name))
                index = int(Path(name).stem.removeprefix("slide"))
                result[f"S{index:02d}"] = "".join(node.text or "" for node in root.iter() if node.tag.endswith("}t"))
    return result


def _selected(page: dict[str, Any]) -> list[str]:
    decisions = page.get("text_decisions", []) if isinstance(page.get("text_decisions"), list) else []
    return [item["selected"] for item in decisions if isinstance(item, dict) and isinstance(item.get("selected"), str) and item["selected"].strip()]


def _skeleton_errors(slide, code: str, slide_id: str) -> list[dict[str, Any]]:
    errors = []
    for name in _SKELETON:
        shape = next((item for item in slide.shapes if item.name == name), None)
        if shape is None or shape.shape_type == _PICTURE or not shape.has_text_frame:
            errors.append(_issue(code, slide_id, f"missing native text skeleton shape {name}"))
    return errors


def audit_deconstruction_pptx(pptx_path: str | Path, page_specs: dict[str, Any], alignment: dict[str, Any], allowed_large_visual_assets_by_page: dict[str, list[str]] | None = None) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(pptx_path)); specs = _pages(page_specs); aligned = _pages(alignment)
    xml = _xml_text(pptx_path); blockers = []
    allowed_by_page = allowed_large_visual_assets_by_page if isinstance(allowed_large_visual_assets_by_page, dict) else {}
    for number, slide in enumerate(presentation.slides, 1):
        slide_id = f"S{number:02d}"; spec = specs.get(slide_id, {}) if isinstance(specs.get(slide_id, {}), dict) else {}; page = aligned.get(slide_id, {}) if isinstance(aligned.get(slide_id, {}), dict) else {}
        elements = [item for item in spec.get("elements", []) if isinstance(item, dict)]
        for element in elements:
            element_id = element.get("element_id")
            if not isinstance(element_id, str):
                continue
            matches = _matching(slide, element_id)
            if not matches:
                blockers.append(_issue(DECONSTRUCTION_EDITABILITY_FAILED, slide_id, f"missing element prefix {_prefix(element_id)}")); continue
            kind = element.get("type")
            if kind == "matrix" and not any(shape.has_table for shape in matches): blockers.append(_issue(DECONSTRUCTION_EDITABILITY_FAILED, slide_id, f"{element_id} requires native table"))
            if kind in _CHART_TYPES and not any(shape.has_chart for shape in matches): blockers.append(_issue(DECONSTRUCTION_EDITABILITY_FAILED, slide_id, f"{element_id} requires native chart"))
            if kind == "flow":
                node = any(shape.shape_type not in {_PICTURE, _LINE, _TEXTBOX} for shape in matches)
                connector = any(shape.shape_type == _LINE for shape in matches)
                if not node or not connector: blockers.append(_issue(DECONSTRUCTION_EDITABILITY_FAILED, slide_id, f"{element_id} requires native node and connector"))
        for selected in _selected(page):
            if selected not in xml.get(slide_id, ""): blockers.append(_issue(DECONSTRUCTION_EDITABILITY_FAILED, slide_id, f"selected text absent from OOXML: {selected!r}"))
        body = _body(slide); body_area = body[2] * body[3]
        native = False
        for element in elements:
            element_id = element.get("element_id")
            if not isinstance(element_id, str): continue
            for shape in _matching(slide, element_id):
                if _overlap(_box(shape), body) <= 0: continue
                if shape.has_text_frame and shape.text.strip() or shape.has_table or shape.has_chart or element.get("type") in _BASIC_TYPES: native = True
        for shape in slide.shapes:
            if shape.shape_type != _PICTURE or body_area <= 0 or _overlap(_box(shape), body) / body_area < .60: continue
            asset_id = next((element.get("asset_id") for element in elements if isinstance(element.get("element_id"), str) and shape.name.startswith(_prefix(element["element_id"]))), None)
            page_allowed = set(allowed_by_page.get(slide_id, [])) if isinstance(allowed_by_page.get(slide_id, []), list) else set()
            if not native and asset_id not in page_allowed: blockers.append(_issue(DECONSTRUCTION_EDITABILITY_FAILED, slide_id, f"large body picture {shape.name} has no editable native body content"))
    ppt_slide_ids = {f"S{index:02d}" for index in range(1, len(presentation.slides) + 1)}
    if not (ppt_slide_ids == set(specs) == set(aligned)):
        blockers.append(_issue(DECONSTRUCTION_EDITABILITY_FAILED, "", "PPTX slide ids must exactly equal page spec and alignment page ids"))
    return {"schema_version": SCHEMA_VERSION, "status": "blocked" if blockers else "pass", "ok": not blockers, "warnings": [], "blockers": blockers, "warning_count": 0, "blocker_count": len(blockers), "page_count": len(presentation.slides), "allowed_large_visual_assets_by_page": allowed_by_page}


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _project_path(project: Path, value: Any) -> Path | None:
    if not isinstance(value, str) or not value or Path(value).is_absolute():
        return None
    candidate = (project / value).resolve()
    try:
        candidate.relative_to(project)
    except ValueError:
        return None
    return candidate


def _matches_declared_crop(blueprint: Path | None, asset: Path | None, source_px: Any) -> bool:
    if (
        blueprint is None
        or asset is None
        or not blueprint.is_file()
        or not asset.is_file()
        or not isinstance(source_px, list)
        or len(source_px) != 4
        or any(not isinstance(value, int) or isinstance(value, bool) for value in source_px)
    ):
        return False
    try:
        with Image.open(blueprint) as source:
            left, top, right, bottom = source_px
            if not (0 <= left < right <= source.width and 0 <= top < bottom <= source.height):
                return False
            expected = source.crop((left, top, right, bottom)).convert("RGBA")
            with Image.open(asset) as actual_source:
                actual = actual_source.convert("RGBA")
                return actual.size == expected.size and actual.tobytes() == expected.tobytes()
    except OSError:
        return False


def audit_bitmap_pptx(pptx_path: str | Path, bitmap_contract: dict[str, Any], project_dir: str | Path) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(pptx_path)); pages = _pages(bitmap_contract); project = Path(project_dir).resolve(); blockers = []
    if not isinstance(bitmap_contract, dict) or bitmap_contract.get("schema_version") != "6.0" or bitmap_contract.get("pipeline_revision") != "6.0.0" or bitmap_contract.get("construction_mode") != "bitmap":
        blockers.append(_issue(BITMAP_CONTRACT_INVALID, "", "invalid V6 bitmap contract header"))
    ppt_slide_ids = {f"S{index:02d}" for index in range(1, len(presentation.slides) + 1)}
    if set(pages) != ppt_slide_ids:
        blockers.append(_issue(BITMAP_CONTRACT_INVALID, "", "PPTX slide ids must exactly equal bitmap contract page ids"))
    for number, slide in enumerate(presentation.slides, 1):
        slide_id = f"S{number:02d}"; record = pages.get(slide_id, {}) if isinstance(pages.get(slide_id, {}), dict) else {}
        blockers.extend(_skeleton_errors(slide, BITMAP_CONTRACT_INVALID, slide_id))
        asset_id = record.get("asset_id"); expected = _prefix(asset_id) if isinstance(asset_id, str) else ""
        blueprint = _project_path(project, record.get("source_blueprint")); asset = _project_path(project, record.get("asset_path"))
        if blueprint is None or not blueprint.is_file() or _sha256(blueprint) != record.get("source_blueprint_sha256"):
            blockers.append(_issue(BITMAP_CONTRACT_INVALID, slide_id, "source blueprint hash chain is invalid"))
        if asset is None or not asset.is_file() or _sha256(asset) != record.get("asset_sha256"):
            blockers.append(_issue(BITMAP_CONTRACT_INVALID, slide_id, "bitmap asset hash chain is invalid"))
        if not _matches_declared_crop(blueprint, asset, record.get("source_px")):
            blockers.append(_issue(BITMAP_CONTRACT_INVALID, slide_id, "bitmap asset does not equal the declared source crop"))
        body = _body(slide); pictures = [shape for shape in slide.shapes if shape.shape_type == _PICTURE and _overlap(_box(shape), body) > .001]
        matches = [shape for shape in pictures if shape.name.startswith(expected)]
        if len(pictures) != 1 or len(matches) != 1:
            blockers.append(_issue(BITMAP_CONTRACT_INVALID, slide_id, "runtime body requires exactly one expected named picture")); continue
        shape = matches[0]
        if hashlib.sha256(shape.image.blob).hexdigest() != record.get("asset_sha256"):
            blockers.append(_issue(BITMAP_CONTRACT_INVALID, slide_id, "body picture SHA-256 does not match contract"))
        if any(abs(getattr(shape, name, 0)) > 1e-6 for name in ("crop_left", "crop_right", "crop_top", "crop_bottom")):
            blockers.append(_issue(BITMAP_CONTRACT_INVALID, slide_id, "body picture must not be cropped"))
        left, top, width, height = _box(shape); aspect = shape.image.size[0] / shape.image.size[1]; body_aspect = body[2] / body[3] if body[3] else 0
        expected_width, expected_height = (body[2], body[2] / aspect) if aspect >= body_aspect else (body[3] * aspect, body[3])
        expected_left = body[0] + (body[2] - expected_width) / 2; expected_top = body[1] + (body[3] - expected_height) / 2
        if any(abs(actual - target) > _TOLERANCE for actual, target in ((left, expected_left), (top, expected_top), (width, expected_width), (height, expected_height))):
            blockers.append(_issue(BITMAP_CONTRACT_INVALID, slide_id, "body picture is not maximal centered contain in runtime body"))
    return {"schema_version": SCHEMA_VERSION, "status": "blocked" if blockers else "pass", "ok": not blockers, "warnings": [], "blockers": blockers, "warning_count": 0, "blocker_count": len(blockers), "page_count": len(presentation.slides)}
