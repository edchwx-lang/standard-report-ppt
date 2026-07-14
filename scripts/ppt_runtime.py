from __future__ import annotations

import json
import hashlib
import re
import shutil
import tempfile
from math import ceil
from pathlib import Path
from typing import Any, Iterable
from zipfile import ZIP_DEFLATED, ZipFile

from PIL import Image, ImageChops, UnidentifiedImageError

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Inches, Pt


MOJIBAKE_RE = re.compile(r"\?{2,}|\ufffd")
ID_RE = re.compile(r"^S\d{2}(?:_M\d{2})?$")
APPROVED_FILL_ROLES = {
    "none",
    "white",
    "deep_blue",
    "mid_blue",
    "light_blue",
    "gray",
    "light_gray",
}
APPROVED_BORDER_ROLES = {"none", "outline", "bottom_line"}
APPROVED_HEADER_ROLES = {"text_only", "line_title", "light_header"}

BODY_RECT = (0.45, 2.10, 12.43, 4.83)
FONT = "Microsoft YaHei"
LEGACY_BLUEPRINT_PROMPT_VERSION = "standard-report-ppt-v4.2"
BLUEPRINT_PROMPT_VERSION = "standard-report-ppt-v4.3"
LEGACY_BLUEPRINT_ACCEPTANCE_FLAGS = (
    "typography_hierarchy",
    "shared_left_edge",
    "core_black_dashed_box",
    "core_square_bullets",
    "source_rule_absent",
    "red_fill_restrained",
)
BLUEPRINT_ACCEPTANCE_FLAGS = (
    "body_composition_usable",
    "module_regions_identifiable",
    "asset_candidates_identified",
    "red_fill_restrained",
    "reconstructable",
)
CATASTROPHIC_RETRY_REASONS = {
    "generation_failed",
    "unreadable_image",
    "multi_page_output",
    "missing_canonical_module",
    "body_composition_unusable",
}
MAX_RED_FILL_RATIO = 0.03

COLORS = {
    "deep_blue": RGBColor(30, 56, 107),
    "light_blue": RGBColor(227, 235, 243),
    "mid_blue": RGBColor(115, 153, 197),
    "light_gray": RGBColor(217, 217, 217),
    "gray": RGBColor(166, 166, 166),
    "dark_gray": RGBColor(102, 102, 102),
    "red": RGBColor(192, 0, 0),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
}


def _blend_with_white(color: RGBColor, alpha: float) -> RGBColor:
    return RGBColor(
        *[round(255 * (1 - alpha) + channel * alpha) for channel in color]
    )


ALLOWED_RGB = {
    tuple(COLORS[name])
    for name in (
        "deep_blue",
        "mid_blue",
        "gray",
        "light_gray",
        "red",
        "white",
        "black",
        "dark_gray",
    )
}
for _base in (COLORS["deep_blue"], COLORS["mid_blue"], COLORS["gray"]):
    for _alpha in (0.1, 0.2, 0.3, 0.4, 0.5):
        ALLOWED_RGB.add(tuple(_blend_with_white(_base, _alpha)))


def _is_allowed_rgb(color: RGBColor | tuple[int, int, int]) -> bool:
    return tuple(color) in ALLOWED_RGB


def _validate_blueprint_attempt_policy(item: dict[str, Any]) -> list[str]:
    errors: list[str] = []
    attempt_count = item.get("attempt_count")
    if attempt_count not in {1, 2}:
        return ["attempt_count must be 1 or 2"]
    retry_reason = item.get("retry_reason")
    if attempt_count == 1:
        if retry_reason not in (None, ""):
            errors.append("retry_reason is only valid when attempt_count is 2")
        return errors
    if retry_reason not in CATASTROPHIC_RETRY_REASONS:
        errors.append(
            "catastrophic_retry requires one of: "
            + ", ".join(sorted(CATASTROPHIC_RETRY_REASONS))
        )
    return errors


def _rectangle_union_area(rectangles: list[list[float]]) -> float:
    if not rectangles:
        return 0.0
    xs = sorted({float(rect[0]) for rect in rectangles} | {float(rect[2]) for rect in rectangles})
    area = 0.0
    for left, right in zip(xs, xs[1:]):
        if right <= left:
            continue
        intervals = sorted(
            (float(rect[1]), float(rect[3]))
            for rect in rectangles
            if float(rect[0]) < right and float(rect[2]) > left
        )
        covered = 0.0
        current_start: float | None = None
        current_end: float | None = None
        for start, end in intervals:
            if current_start is None:
                current_start, current_end = start, end
            elif start > float(current_end):
                covered += float(current_end) - current_start
                current_start, current_end = start, end
            else:
                current_end = max(float(current_end), end)
        if current_start is not None:
            covered += float(current_end) - current_start
        area += (right - left) * covered
    return area


def _evidence_counts(modules: list[dict[str, Any]]) -> tuple[int, int]:
    evidence_items = 0
    quantified_facts = 0

    def add_value(value: Any) -> None:
        nonlocal evidence_items, quantified_facts
        if value in (None, ""):
            return
        evidence_items += 1
        if re.search(r"\d", str(value)):
            quantified_facts += 1

    for module in modules:
        for bullet in module.get("bullets", []):
            add_value(bullet)
        for step in module.get("steps", []):
            add_value(step)
        for metric in module.get("metrics", []):
            add_value(metric.get("value"))
        matrix = module.get("matrix", {})
        for row in matrix.get("rows", []):
            for cell in row:
                add_value(cell)
        chart = module.get("chart", {})
        for series in chart.get("series", []):
            for value in series.get("values", []):
                add_value(value)
        for point in chart.get("points", []):
            add_value(point.get("label"))
            for field in ("x", "y", "size"):
                if field in point:
                    add_value(point[field])
        for label in module.get("map", {}).get("labels", []):
            add_value(label.get("label"))
    return evidence_items, quantified_facts


def _validate_density_slide(
    slide: dict[str, Any], geometry_slide: dict[str, Any] | None
) -> list[str]:
    slide_id = str(slide.get("slide_id", "slide"))
    errors: list[str] = []
    if slide.get("density_profile") != "medium_high":
        return [f"{slide_id}: density_profile must be medium_high"]
    core_chars = len(re.sub(r"\s+", "", "".join(slide.get("core_points", []))))
    if not 60 <= core_chars <= 180:
        errors.append(f"{slide_id}: medium-high core text must contain 60 to 180 characters")
    modules = slide.get("modules", [])
    if len(modules) < 5:
        errors.append(f"{slide_id}: medium-high density requires at least 5 modules")
    groups = slide.get("evidence_groups")
    if not isinstance(groups, list) or len(groups) < 3:
        errors.append(f"{slide_id}: medium-high density requires at least 3 evidence groups")
    else:
        expected_ids = {module.get("module_id") for module in modules}
        grouped_ids: list[str] = []
        for group in groups:
            module_ids = group.get("module_ids", []) if isinstance(group, dict) else []
            if not module_ids:
                errors.append(f"{slide_id}: every evidence group must reference a module")
            grouped_ids.extend(str(module_id) for module_id in module_ids)
        if set(grouped_ids) != expected_ids or len(grouped_ids) != len(set(grouped_ids)):
            errors.append(f"{slide_id}: evidence groups must cover every module exactly once")
    evidence_items, quantified_facts = _evidence_counts(modules)
    if evidence_items < 8:
        errors.append(f"{slide_id}: medium-high density requires at least 8 evidence items")
    if slide.get("quantitative_evidence_expected") is not False and quantified_facts < 3:
        errors.append(f"{slide_id}: medium-high density requires at least 3 quantified facts")
    zones = []
    if isinstance(geometry_slide, dict):
        zones = [
            module.get("zone_norm")
            for module in geometry_slide.get("modules", [])
            if _valid_zone(module.get("zone_norm"))
        ]
    if _rectangle_union_area(zones) < 0.68:
        errors.append(f"{slide_id}: body coverage must be at least 68%")
    return errors


def load_json(path: str | Path) -> dict[str, Any]:
    path = Path(path)
    text = path.read_text(encoding="utf-8")
    if MOJIBAKE_RE.search(text):
        raise ValueError(f"{path}: mojibake detected")
    value = json.loads(text)
    if not isinstance(value, dict):
        raise ValueError(f"{path}: root must be an object")
    return value


def load_project(
    spec_path: str | Path, geometry_path: str | Path
) -> tuple[dict[str, Any], dict[str, Any]]:
    return load_json(spec_path), load_json(geometry_path)


def _iter_strings(value: Any) -> Iterable[str]:
    if isinstance(value, str):
        yield value
    elif isinstance(value, dict):
        for item in value.values():
            yield from _iter_strings(item)
    elif isinstance(value, list):
        for item in value:
            yield from _iter_strings(item)


def _valid_zone(zone: Any) -> bool:
    if not isinstance(zone, list) or len(zone) != 4:
        return False
    if not all(isinstance(item, (int, float)) for item in zone):
        return False
    x1, y1, x2, y2 = zone
    return 0 <= x1 < x2 <= 1 and 0 <= y1 < y2 <= 1


def _valid_blueprint_pixels(size: Any, body: Any) -> bool:
    if not isinstance(size, list) or len(size) != 2:
        return False
    if not all(isinstance(value, int) and not isinstance(value, bool) and value > 0 for value in size):
        return False
    if not isinstance(body, list) or len(body) != 4:
        return False
    if not all(isinstance(value, (int, float)) for value in body):
        return False
    x0, y0, x1, y1 = body
    return 0 <= x0 < x1 <= size[0] and 0 <= y0 < y1 <= size[1]


def _valid_pixel_rect(size: Any, rect: Any) -> bool:
    if not isinstance(size, list) or len(size) != 2:
        return False
    if not isinstance(rect, list) or len(rect) != 4:
        return False
    if not all(isinstance(value, (int, float)) for value in rect):
        return False
    x0, y0, x1, y1 = rect
    return 0 <= x0 < x1 <= size[0] and 0 <= y0 < y1 <= size[1]


def validate_contracts(
    spec: dict[str, Any],
    geometry: dict[str, Any],
    brief: dict[str, Any] | None = None,
    blueprint_manifest: dict[str, Any] | None = None,
    *,
    project_dir: str | Path | None = None,
) -> list[str]:
    errors: list[str] = []

    if not isinstance(brief, dict):
        errors.append("project_brief: project_brief.json is required")
        brief = {}
    if not isinstance(blueprint_manifest, dict):
        errors.append("blueprint manifest: blueprint_generation_manifest.json is required")
        blueprint_manifest = {}
    project_dir = Path(project_dir) if project_dir is not None else None
    spec_version = spec.get("schema_version")
    geometry_version = geometry.get("schema_version")
    manifest_version = blueprint_manifest.get("schema_version")
    is_v43 = spec_version == "4.1"

    if brief.get("schema_version") != "1.0":
        errors.append("project_brief: schema_version must be 1.0")
    production_mode = brief.get("production_mode")
    if production_mode not in {"blueprint", "fast"}:
        errors.append("project_brief: production_mode must be blueprint or fast")
    confirmation_source = brief.get("confirmation_source")
    if confirmation_source not in {"user_explicit", "user_selected"}:
        errors.append("project_brief: confirmation_source must prove explicit user choice")
    requested_page_count = brief.get("requested_page_count")
    if not isinstance(requested_page_count, int) or isinstance(requested_page_count, bool) or requested_page_count <= 0:
        errors.append("project_brief: requested_page_count must be a positive integer")
    page_mapping = brief.get("page_mapping")
    if not isinstance(page_mapping, list):
        errors.append("project_brief: page_mapping must be a list")
    else:
        mapped_pages: set[int] = set()
        for item in page_mapping:
            if not isinstance(item, dict):
                errors.append("project_brief: each page_mapping item must be an object")
                continue
            page_number = item.get("page_number")
            source_scope = item.get("source_scope")
            if not isinstance(page_number, int) or isinstance(page_number, bool):
                errors.append("project_brief: page_mapping page_number must be an integer")
            elif isinstance(requested_page_count, int) and not 1 <= page_number <= requested_page_count:
                errors.append(f"project_brief: page_mapping page_number {page_number} is out of range")
            elif page_number in mapped_pages:
                errors.append(f"project_brief: duplicate page_mapping page_number {page_number}")
            else:
                mapped_pages.add(page_number)
            if not isinstance(source_scope, str) or not source_scope.strip():
                errors.append("project_brief: page_mapping source_scope is required")

    if manifest_version not in {"4.1", "4.2"}:
        errors.append("blueprint manifest: schema_version must be 4.1 or 4.2")
    if blueprint_manifest.get("production_mode") != production_mode:
        errors.append("blueprint manifest: production_mode does not match project_brief")
    if production_mode == "blueprint" and blueprint_manifest.get("imagegen_used") is not True:
        errors.append("blueprint mode: imagegen_used must be true")
    if production_mode == "fast" and blueprint_manifest.get("imagegen_used") is not False:
        errors.append("fast mode: imagegen_used must be false")

    if spec_version not in {"4.0", "4.1"}:
        errors.append("spec: schema_version must be 4.0 or 4.1")
    if geometry_version not in {"4.0", "4.1"}:
        errors.append("geometry: schema_version must be 4.0 or 4.1")
    if is_v43 and geometry_version != "4.1":
        errors.append("V4.3 spec requires geometry schema_version 4.1")
    if is_v43 and manifest_version != "4.2":
        errors.append("V4.3 spec requires blueprint manifest schema_version 4.2")

    for text in _iter_strings(spec):
        if MOJIBAKE_RE.search(text):
            errors.append("spec: mojibake detected")
            break

    slides = spec.get("slides")
    geometry_slides = geometry.get("slides")
    if not isinstance(slides, list) or not slides:
        return errors + ["spec: slides must be a non-empty list"]
    if not isinstance(geometry_slides, list):
        return errors + ["geometry: slides must be a list"]
    if isinstance(requested_page_count, int) and requested_page_count != len(slides):
        errors.append(
            f"project_brief: requested_page_count {requested_page_count} does not match {len(slides)} spec slides"
        )

    slide_ids: set[str] = set()
    module_ids: set[str] = set()
    spec_modules_by_slide: dict[str, set[str]] = {}
    map_modules_by_slide: dict[str, set[str]] = {}

    for slide in slides:
        if not isinstance(slide, dict):
            errors.append("spec: each slide must be an object")
            continue
        slide_id = slide.get("slide_id")
        if not isinstance(slide_id, str) or not re.match(r"^S\d{2}$", slide_id):
            errors.append(f"spec: invalid slide_id {slide_id!r}")
            continue
        if slide_id in slide_ids:
            errors.append(f"spec: duplicate slide_id {slide_id}")
        slide_ids.add(slide_id)
        for field in (
            "chapter",
            "title",
            "core",
            "source",
            "page_type",
            "layout_archetype",
        ):
            if not isinstance(slide.get(field), str) or not slide[field].strip():
                errors.append(f"{slide_id}: missing {field}")
        core_points = slide.get("core_points")
        if not isinstance(core_points, list) or not core_points:
            errors.append(f"{slide_id}: core_points must contain 1 to 3 non-empty strings")
        elif len(core_points) > 3:
            errors.append(f"{slide_id}: core_points may contain at most 3 items")
        elif not all(isinstance(point, str) and point.strip() for point in core_points):
            errors.append(f"{slide_id}: core_points must contain 1 to 3 non-empty strings")

        local_ids: set[str] = set()
        local_map_ids: set[str] = set()
        modules = slide.get("modules")
        if not isinstance(modules, list) or not modules:
            errors.append(f"{slide_id}: modules must be a non-empty list")
            modules = []
        for module in modules:
            if not isinstance(module, dict):
                errors.append(f"{slide_id}: module must be an object")
                continue
            module_id = module.get("module_id")
            if not isinstance(module_id, str) or not re.match(
                rf"^{re.escape(slide_id)}_M\d{{2}}$", module_id
            ):
                errors.append(f"{slide_id}: invalid module_id {module_id!r}")
                continue
            if module_id in module_ids:
                errors.append(f"spec: duplicate module_id {module_id}")
            module_ids.add(module_id)
            local_ids.add(module_id)
            if module.get("role") == "map" or module.get("primitive") == "map":
                local_map_ids.add(module_id)
            if not isinstance(module.get("role"), str) or not module["role"].strip():
                errors.append(f"{module_id}: missing role")
            if not isinstance(module.get("title"), str):
                errors.append(f"{module_id}: missing title")
            for field in ("primitive", "fill_role"):
                if not isinstance(module.get(field), str) or not module[field].strip():
                    errors.append(f"{module_id}: missing {field}")
            if module.get("fill_role") not in APPROVED_FILL_ROLES:
                errors.append(f"{module_id}: unapproved fill_role {module.get('fill_role')}")
        spec_modules_by_slide[slide_id] = local_ids
        map_modules_by_slide[slide_id] = local_map_ids

    geometry_by_slide: dict[str, dict[str, Any]] = {}
    for slide in geometry_slides:
        if not isinstance(slide, dict):
            errors.append("geometry: each slide must be an object")
            continue
        slide_id = slide.get("slide_id")
        if slide_id not in slide_ids:
            errors.append(f"geometry: unknown slide_id {slide_id}")
            continue
        if slide_id in geometry_by_slide:
            errors.append(f"geometry: duplicate slide_id {slide_id}")
        geometry_by_slide[slide_id] = slide
        if not _valid_blueprint_pixels(
            slide.get("blueprint_size_px"), slide.get("body_source_px")
        ):
            errors.append(f"{slide_id}: invalid blueprint_size_px or body_source_px")

        geometry_ids: set[str] = set()
        visual_strategy_by_module: dict[str, str] = {}
        geometry_modules = slide.get("modules")
        if not isinstance(geometry_modules, list):
            errors.append(f"{slide_id}: geometry modules must be a list")
            geometry_modules = []
        for module in geometry_modules:
            if not isinstance(module, dict):
                errors.append(f"{slide_id}: geometry module must be an object")
                continue
            module_id = module.get("module_id")
            if module_id not in spec_modules_by_slide.get(slide_id, set()):
                if module_id in module_ids:
                    errors.append(
                        f"geometry: module_id {module_id} belongs to another slide"
                    )
                else:
                    errors.append(f"geometry: unknown module_id {module_id}")
                continue
            if module_id in geometry_ids:
                errors.append(f"{slide_id}: duplicate geometry module_id {module_id}")
                continue
            geometry_ids.add(module_id)
            if not _valid_zone(module.get("zone_norm")):
                errors.append(f"{module_id}: invalid zone_norm")
            content_zone = module.get("content_zone_norm")
            if content_zone is not None:
                if not _valid_zone(content_zone):
                    errors.append(f"{module_id}: invalid content_zone_norm")
                elif _valid_zone(module.get("zone_norm")):
                    zone = module["zone_norm"]
                    if not (
                        zone[0] <= content_zone[0]
                        and zone[1] <= content_zone[1]
                        and content_zone[2] <= zone[2]
                        and content_zone[3] <= zone[3]
                    ):
                        errors.append(f"{module_id}: content_zone_norm must stay inside zone_norm")
            for field in ("primitive", "border_role", "fill_role", "header_role"):
                if not isinstance(module.get(field), str) or not module[field].strip():
                    errors.append(f"{module_id}: missing geometry {field}")
            if module.get("fill_role") not in APPROVED_FILL_ROLES:
                errors.append(f"{module_id}: unapproved geometry fill_role {module.get('fill_role')}")
            if module.get("border_role") not in APPROVED_BORDER_ROLES:
                errors.append(f"{module_id}: unapproved border_role {module.get('border_role')}")
            if module.get("header_role") not in APPROVED_HEADER_ROLES:
                errors.append(f"{module_id}: unapproved header_role {module.get('header_role')}")
            if is_v43 and production_mode == "blueprint":
                visual_strategy = module.get("visual_strategy")
                if visual_strategy not in {"native_primitives", "blueprint_crop", "hybrid"}:
                    errors.append(
                        f"{module_id}: visual_strategy must be native_primitives, blueprint_crop, or hybrid"
                    )
                else:
                    visual_strategy_by_module[module_id] = visual_strategy

        missing = spec_modules_by_slide.get(slide_id, set()) - geometry_ids
        for module_id in sorted(missing):
            errors.append(f"{slide_id}: missing geometry for {module_id}")

        if production_mode == "blueprint":
            if slide.get("geometry_source") != "imagegen_deconstruction":
                errors.append(
                    f"{slide_id}: blueprint mode geometry_source must be imagegen_deconstruction"
                )
            acceptance = slide.get("blueprint_acceptance")
            if not isinstance(acceptance, dict):
                errors.append(f"{slide_id}: blueprint_acceptance must be an object")
            else:
                acceptance_flags = (
                    BLUEPRINT_ACCEPTANCE_FLAGS if is_v43 else LEGACY_BLUEPRINT_ACCEPTANCE_FLAGS
                )
                for flag in acceptance_flags:
                    if acceptance.get(flag) is not True:
                        errors.append(
                            f"{slide_id}: blueprint_acceptance flag {flag} must be true"
                        )

            regions = slide.get("blueprint_regions")
            if not isinstance(regions, list):
                errors.append(f"{slide_id}: blueprint_regions must be a list")
            else:
                region_ids: set[str] = set()
                for region in regions:
                    if not isinstance(region, dict):
                        errors.append(f"{slide_id}: blueprint_regions item must be an object")
                        continue
                    module_id = region.get("module_id")
                    if module_id not in geometry_ids:
                        errors.append(
                            f"{slide_id}: blueprint_regions references unknown module_id {module_id}"
                        )
                        continue
                    if module_id in region_ids:
                        errors.append(
                            f"{slide_id}: duplicate blueprint_regions module_id {module_id}"
                        )
                    region_ids.add(module_id)
                    if not _valid_pixel_rect(
                        slide.get("blueprint_size_px"), region.get("source_px")
                    ):
                        errors.append(
                            f"{slide_id}/{module_id}: invalid blueprint_regions source_px"
                        )
                for module_id in sorted(geometry_ids - region_ids):
                    errors.append(
                        f"{slide_id}: blueprint_regions missing module {module_id}"
                    )
                for module_id in sorted(region_ids - geometry_ids):
                    errors.append(
                        f"{slide_id}: blueprint_regions unknown module {module_id}"
                    )

        assets = slide.get("assets", [])
        if not isinstance(assets, list):
            errors.append(f"{slide_id}: assets must be a list")
        else:
            asset_ids: set[str] = set()
            asset_module_ids: set[str] = set()
            background_asset_modules: set[str] = set()
            for asset in assets:
                if not isinstance(asset, dict):
                    errors.append(f"{slide_id}: asset must be an object")
                    continue
                asset_id = asset.get("asset_id")
                if not isinstance(asset_id, str) or not asset_id:
                    errors.append(f"{slide_id}: missing asset_id")
                elif asset_id in asset_ids:
                    errors.append(f"{slide_id}: duplicate asset_id {asset_id}")
                else:
                    asset_ids.add(asset_id)
                if asset.get("module_id") not in spec_modules_by_slide.get(slide_id, set()):
                    errors.append(f"{slide_id}: asset references unknown module_id")
                elif isinstance(asset.get("module_id"), str):
                    asset_module_ids.add(asset["module_id"])
                if not _valid_zone(asset.get("zone_norm")):
                    errors.append(f"{slide_id}/{asset_id}: invalid zone_norm")
                if not isinstance(asset.get("path"), str) or not asset["path"].strip():
                    errors.append(f"{slide_id}/{asset_id}: missing asset path")
                elif project_dir is not None and asset.get("required", True):
                    asset_path = Path(asset["path"])
                    if not asset_path.is_absolute():
                        asset_path = project_dir / asset_path
                    if not asset_path.is_file() or asset_path.stat().st_size == 0:
                        errors.append(f"{slide_id}/{asset_id}: required asset file missing or empty")
                if "required" in asset and not isinstance(asset["required"], bool):
                    errors.append(f"{slide_id}/{asset_id}: required must be boolean")
                layer = asset.get("layer", "foreground")
                if layer not in {"background", "foreground"}:
                    errors.append(f"{slide_id}/{asset_id}: layer must be background or foreground")
                if layer == "background" and isinstance(asset.get("module_id"), str):
                    background_asset_modules.add(asset["module_id"])
                if is_v43 and production_mode == "blueprint":
                    source_type = asset.get("source_type")
                    if source_type not in {"blueprint_crop", "external_reference"}:
                        errors.append(
                            f"{slide_id}/{asset_id}: source_type must be blueprint_crop or external_reference"
                        )
                    if source_type == "blueprint_crop":
                        if not _valid_pixel_rect(
                            slide.get("blueprint_size_px"), asset.get("source_px")
                        ):
                            errors.append(f"{slide_id}/{asset_id}: invalid asset source_px")
                        else:
                            width, height = slide["blueprint_size_px"]
                            left, top, right, bottom = asset["source_px"]
                            if ((right - left) * (bottom - top)) / (width * height) > 0.35:
                                errors.append(f"{slide_id}/{asset_id}: blueprint crop exceeds 35% of page")
                        asset_path_value = asset.get("path")
                        if isinstance(asset_path_value, str) and asset_path_value:
                            asset_path = Path(asset_path_value)
                            if asset_path.is_absolute() or asset_path.parts[:1] != ("blueprint_assets",):
                                errors.append(
                                    f"{slide_id}/{asset_id}: blueprint crop path must be under blueprint_assets"
                                )
                        digest = asset.get("sha256")
                        pixel_size = asset.get("pixel_size")
                        if not isinstance(digest, str) or not re.fullmatch(r"[0-9a-f]{64}", digest):
                            errors.append(f"{slide_id}/{asset_id}: missing valid asset sha256")
                        if (
                            not isinstance(pixel_size, list)
                            or len(pixel_size) != 2
                            or not all(isinstance(value, int) and value > 0 for value in pixel_size)
                        ):
                            errors.append(f"{slide_id}/{asset_id}: missing valid asset pixel_size")
                        if project_dir is not None and isinstance(asset_path_value, str):
                            source = project_dir / asset_path_value
                            if source.is_file():
                                actual_digest = hashlib.sha256(source.read_bytes()).hexdigest()
                                if isinstance(digest, str) and actual_digest != digest:
                                    errors.append(f"{slide_id}/{asset_id}: asset sha256 mismatch")
                                try:
                                    with Image.open(source) as image:
                                        if list(image.size) != pixel_size:
                                            errors.append(f"{slide_id}/{asset_id}: asset pixel_size mismatch")
                                        blueprint_value = slide.get("blueprint_file")
                                        if isinstance(blueprint_value, str):
                                            blueprint_source = project_dir / blueprint_value
                                            if blueprint_source.is_file() and _valid_pixel_rect(
                                                slide.get("blueprint_size_px"), asset.get("source_px")
                                            ):
                                                with Image.open(blueprint_source) as blueprint_image:
                                                    expected_crop = blueprint_image.crop(
                                                        tuple(asset["source_px"])
                                                    ).convert("RGB")
                                                    actual_crop = image.convert("RGB")
                                                    if (
                                                        expected_crop.size != actual_crop.size
                                                        or ImageChops.difference(
                                                            expected_crop, actual_crop
                                                        ).getbbox()
                                                        is not None
                                                    ):
                                                        errors.append(
                                                            f"{slide_id}/{asset_id}: asset does not match blueprint source pixels"
                                                        )
                                except UnidentifiedImageError:
                                    errors.append(f"{slide_id}/{asset_id}: asset is not a readable image")
            for module_id in sorted(map_modules_by_slide.get(slide_id, set()) - background_asset_modules):
                errors.append(f"{slide_id}: map module {module_id} requires a background asset")
            if is_v43 and production_mode == "blueprint":
                for module_id, strategy in visual_strategy_by_module.items():
                    if strategy in {"blueprint_crop", "hybrid"} and module_id not in asset_module_ids:
                        errors.append(
                            f"{slide_id}/{module_id}: visual_strategy {strategy} requires a declared asset"
                        )

    for slide_id in sorted(slide_ids - set(geometry_by_slide)):
        errors.append(f"geometry: missing slide {slide_id}")

    if is_v43:
        for slide in slides:
            if isinstance(slide, dict):
                errors.extend(
                    _validate_density_slide(slide, geometry_by_slide.get(slide.get("slide_id")))
                )

    manifest_slides = blueprint_manifest.get("slides")
    if not isinstance(manifest_slides, list):
        errors.append("blueprint manifest: slides must be a list")
        manifest_slides = []
    manifest_by_slide: dict[str, dict[str, Any]] = {}
    for item in manifest_slides:
        if not isinstance(item, dict):
            errors.append("blueprint manifest: each slide must be an object")
            continue
        slide_id = item.get("slide_id")
        if slide_id in manifest_by_slide:
            errors.append(f"blueprint manifest: duplicate slide_id {slide_id}")
        elif slide_id not in slide_ids:
            errors.append(f"blueprint manifest: unknown slide_id {slide_id}")
        else:
            manifest_by_slide[slide_id] = item

    for slide_id in sorted(slide_ids):
        geometry_slide = geometry_by_slide.get(slide_id, {})
        manifest_slide = manifest_by_slide.get(slide_id)
        if manifest_slide is None:
            errors.append(f"blueprint manifest: missing slide {slide_id}")
            continue
        blueprint_file = geometry_slide.get("blueprint_file")
        if not isinstance(blueprint_file, str) or not blueprint_file:
            errors.append(f"{slide_id}: missing blueprint_file")
            continue
        if production_mode == "fast":
            if not blueprint_file.startswith("runtime_archetype/"):
                errors.append(f"{slide_id}: fast mode requires runtime_archetype geometry")
            if manifest_slide.get("status") != "runtime_archetype":
                errors.append(f"{slide_id}: fast mode manifest status must be runtime_archetype")
            expected_archetype = blueprint_file.partition("/")[2]
            if manifest_slide.get("archetype") != expected_archetype:
                errors.append(f"{slide_id}: fast mode archetype does not match geometry")
        elif production_mode == "blueprint":
            if manifest_slide.get("generator") != "imagegen":
                errors.append(f"{slide_id}: blueprint mode generator must be imagegen")
            expected_prompt_version = (
                BLUEPRINT_PROMPT_VERSION if is_v43 else LEGACY_BLUEPRINT_PROMPT_VERSION
            )
            if manifest_slide.get("prompt_version") != expected_prompt_version:
                errors.append(
                    f"{slide_id}: blueprint mode prompt_version must be {expected_prompt_version}"
                )
            if is_v43:
                errors.extend(
                    f"{slide_id}: {error}"
                    for error in _validate_blueprint_attempt_policy(manifest_slide)
                )
            if blueprint_file.startswith("runtime_archetype/"):
                errors.append(f"{slide_id}: blueprint mode forbids runtime_archetype geometry")
                continue
            if manifest_slide.get("status") != "generated":
                errors.append(f"{slide_id}: blueprint mode manifest status must be generated")
            if manifest_slide.get("file") != blueprint_file:
                errors.append(f"{slide_id}: blueprint mode manifest file does not match geometry")
            if project_dir is None:
                errors.append(f"{slide_id}: blueprint mode requires project_dir for file validation")
                continue
            image_path = project_dir / blueprint_file
            if not image_path.is_file() or image_path.stat().st_size == 0:
                errors.append(f"{slide_id}: blueprint mode file missing or empty: {blueprint_file}")
                continue
            digest = hashlib.sha256(image_path.read_bytes()).hexdigest()
            if geometry_slide.get("blueprint_sha256") != digest:
                errors.append(f"{slide_id}: blueprint geometry hash mismatch")
            if manifest_slide.get("sha256") != digest:
                errors.append(f"{slide_id}: blueprint manifest hash mismatch")
            try:
                with Image.open(image_path) as image:
                    if list(image.size) != geometry_slide.get("blueprint_size_px"):
                        errors.append(f"{slide_id}: blueprint_size_px does not match image")
                    image.verify()
            except (UnidentifiedImageError, OSError):
                errors.append(f"{slide_id}: blueprint file is not a valid image")

    return errors


def _set_run_font(run: Any, size: float, color: RGBColor, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = rpr.find(qn(tag))
        if node is None:
            node = rpr.makeelement(qn(tag), {"typeface": FONT})
            rpr.append(node)
        else:
            node.set("typeface", FONT)


def _add_text(
    slide: Any,
    text: str,
    rect: tuple[float, float, float, float],
    *,
    size: float = 9,
    color: RGBColor | None = None,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    name: str,
    vertical: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.MIDDLE,
    margin_left_cm: float = 0.10,
    body_paragraph: bool = False,
    bullet: bool = False,
) -> Any:
    x, y, w, h = rect
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Cm(margin_left_cm)
    tf.margin_right = Cm(0.08)
    tf.margin_top = Cm(0.03)
    tf.margin_bottom = Cm(0.03)
    tf.vertical_anchor = vertical
    p = tf.paragraphs[0]
    value = str(text)
    if bullet and value and not value.lstrip().startswith("■"):
        value = f"■ {value}"
    p.text = value
    p.alignment = align
    if body_paragraph:
        p.alignment = PP_ALIGN.JUSTIFY
        p.space_before = Pt(0)
        p.space_after = Pt(6)
        p.line_spacing = 1.2
        ppr = p._p.get_or_add_pPr()
        ppr.set("marL", str(Cm(0.64)))
        ppr.set("indent", str(-Cm(0.64)))
    for run in p.runs:
        _set_run_font(run, size, color or COLORS["black"], bold)
    return shape


def _add_core_text(
    slide: Any,
    points: list[str],
    rect: tuple[float, float, float, float],
    *,
    name: str,
) -> Any:
    x, y, w, h = rect
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Cm(0.18)
    tf.margin_right = Cm(0.08)
    tf.margin_top = Cm(0.03)
    tf.margin_bottom = Cm(0.03)
    tf.vertical_anchor = (
        MSO_VERTICAL_ANCHOR.MIDDLE
        if len(points) == 1
        else MSO_VERTICAL_ANCHOR.TOP
    )
    for index, point in enumerate(points):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        value = str(point).strip()
        if not value.startswith("■"):
            value = f"■ {value}"
        paragraph.text = value
        paragraph.alignment = PP_ALIGN.JUSTIFY
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(6)
        paragraph.line_spacing = 1.2
        ppr = paragraph._p.get_or_add_pPr()
        ppr.set("marL", str(Cm(0.64)))
        ppr.set("indent", str(-Cm(0.64)))
        for run in paragraph.runs:
            _set_run_font(run, 12, COLORS["black"])
    return shape


def _add_rect(
    slide: Any,
    rect: tuple[float, float, float, float],
    *,
    fill: RGBColor | None,
    line: RGBColor | None,
    name: str,
    line_width: float = 0.75,
    dashed: bool = False,
) -> Any:
    x, y, w, h = rect
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.name = name
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
        if dashed:
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return shape


def _add_line(
    slide: Any,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: RGBColor,
    width: float = 1.2,
    name: str,
) -> Any:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.name = name
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def _remove_existing_slides(prs: Presentation) -> None:
    while len(prs.slides):
        slide_id = prs.slides._sldIdLst[0]
        relationship_id = slide_id.rId
        prs.part.drop_rel(relationship_id)
        del prs.slides._sldIdLst[0]


def _blank_layout(prs: Presentation) -> Any:
    named = [layout for layout in prs.slide_layouts if "blank" in layout.name.lower()]
    if named:
        return named[0]
    return min(prs.slide_layouts, key=lambda layout: len(layout.placeholders))


def _fill_color(role: str | None) -> RGBColor | None:
    if role in (None, "none"):
        return None
    return COLORS.get(role, COLORS["white"])


def _map_zone(zone: list[float]) -> tuple[float, float, float, float]:
    x1, y1, x2, y2 = zone
    bx, by, bw, bh = BODY_RECT
    return bx + x1 * bw, by + y1 * bh, (x2 - x1) * bw, (y2 - y1) * bh


def _add_skeleton(slide: Any, slide_spec: dict[str, Any], page_number: int) -> None:
    sid = slide_spec["slide_id"]
    _add_text(
        slide,
        slide_spec["chapter"],
        (0.45, Cm(0.35).inches, 12.43, 0.35),
        size=20,
        color=COLORS["deep_blue"],
        bold=True,
        name=f"{sid}_SECTION_TITLE",
        margin_left_cm=0.18,
    )
    _add_rect(
        slide,
        (0.45, Cm(1.5).inches, 12.43, Cm(1).inches),
        fill=COLORS["deep_blue"],
        line=COLORS["deep_blue"],
        name=f"{sid}_PAGE_TITLE_BAR",
    )
    _add_text(
        slide,
        slide_spec["title"],
        (0.45, Cm(1.5).inches, 12.43, Cm(1).inches),
        size=16,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.LEFT,
        name=f"{sid}_PAGE_TITLE_TEXT",
        margin_left_cm=0.18,
    )
    _add_rect(
        slide,
        (0.45, Cm(2.85).inches, 12.43, 0.92),
        fill=COLORS["white"],
        line=COLORS["black"],
        name=f"{sid}_CORE_BOX",
        line_width=1.0,
        dashed=True,
    )
    _add_core_text(
        slide,
        slide_spec["core_points"],
        (0.45, Cm(2.85).inches + 0.03, 12.43, 0.86),
        name=f"{sid}_CORE_TEXT",
    )
    _add_text(
        slide,
        slide_spec["source"],
        (0.45, 7.12, 10.6, 0.18),
        size=7,
        color=COLORS["dark_gray"],
        name=f"{sid}_SOURCE",
    )
    _add_text(
        slide,
        str(page_number),
        (12.15, 7.07, 0.55, 0.22),
        size=8,
        color=COLORS["black"],
        align=PP_ALIGN.RIGHT,
        name=f"{sid}_PAGE_NUMBER",
    )


def _module_root(slide: Any, module_id: str, rect: tuple[float, float, float, float]) -> Any:
    return _add_rect(slide, rect, fill=None, line=None, name=module_id)


def _render_card(
    slide: Any,
    module: dict[str, Any],
    rect: tuple[float, float, float, float],
    content_rect: tuple[float, float, float, float] | None = None,
) -> None:
    module_id = module["module_id"]
    x, y, w, h = rect
    fill = _fill_color(module.get("fill_role")) or COLORS["white"]
    _add_rect(
        slide,
        rect,
        fill=fill,
        line=COLORS["mid_blue"],
        name=f"{module_id}_CARD",
        line_width=1.0,
    )
    _add_text(
        slide,
        module.get("title", ""),
        (x + 0.12, y + 0.06, w - 0.24, min(0.38, h * 0.22)),
        size=11,
        color=COLORS["deep_blue"],
        bold=True,
        name=f"{module_id}_TITLE",
    )
    cx, cy, cw, ch = content_rect or rect
    content_bottom = min(y + h, cy + ch)
    cursor = max(y + min(0.48, h * 0.27), cy)
    metrics = module.get("metrics", [])
    if metrics:
        cols = max(1, len(metrics))
        metric_w = (cw - 0.18) / cols
        if module.get("metric_layout") == "vertical" or metric_w < 0.65:
            row_h = min(0.42, max(0.30, (content_bottom - cursor) * 0.55 / cols))
            for index, metric in enumerate(metrics):
                my = cursor + index * row_h
                _add_text(
                    slide,
                    metric.get("value", ""),
                    (cx + 0.04, my, cw * 0.48, row_h),
                    size=11,
                    color=COLORS["red"],
                    bold=True,
                    align=PP_ALIGN.LEFT,
                    name=f"{module_id}_METRIC_{index + 1}",
                )
                _add_text(
                    slide,
                    metric.get("label", ""),
                    (cx + cw * 0.52, my + 0.02, cw * 0.44, row_h - 0.02),
                    size=8,
                    color=COLORS["dark_gray"],
                    align=PP_ALIGN.LEFT,
                    name=f"{module_id}_METRIC_LABEL_{index + 1}",
                )
            cursor += row_h * cols + 0.05
        else:
            for index, metric in enumerate(metrics):
                mx = cx + 0.09 + index * metric_w
                _add_text(
                    slide,
                    metric.get("value", ""),
                    (mx, cursor, metric_w - 0.05, 0.32),
                    size=13,
                    color=COLORS["red"],
                    bold=True,
                    align=PP_ALIGN.CENTER,
                    name=f"{module_id}_METRIC_{index + 1}",
                )
                _add_text(
                    slide,
                    metric.get("label", ""),
                    (mx, cursor + 0.28, metric_w - 0.05, 0.24),
                    size=8,
                    color=COLORS["dark_gray"],
                    align=PP_ALIGN.CENTER,
                    name=f"{module_id}_METRIC_LABEL_{index + 1}",
                )
            cursor += 0.60
    bullets = module.get("bullets", [])
    available = max(0.25, content_bottom - cursor - 0.08)
    row_h = min(0.38, available / max(1, len(bullets)))
    for index, bullet in enumerate(bullets):
        _add_text(
            slide,
            str(bullet),
            (cx + 0.14, cursor + index * row_h, cw - 0.28, row_h),
            size=8.5,
            color=COLORS["black"],
            name=f"{module_id}_BULLET_{index + 1}",
            body_paragraph=True,
            bullet=True,
        )


def _render_line_group(
    slide: Any,
    module: dict[str, Any],
    rect: tuple[float, float, float, float],
    content_rect: tuple[float, float, float, float] | None = None,
) -> None:
    module_id = module["module_id"]
    x, y, w, h = rect
    _add_text(
        slide,
        module.get("title", ""),
        (x, y, w, 0.34),
        size=11,
        color=COLORS["deep_blue"],
        bold=True,
        name=f"{module_id}_TITLE",
    )
    _add_line(
        slide,
        x,
        y + 0.34,
        x + w,
        y + 0.34,
        color=COLORS["deep_blue"],
        name=f"{module_id}_LINE",
    )
    bullets = module.get("bullets", [])
    cx, cy, cw, ch = content_rect or rect
    content_top = max(y + 0.43, cy)
    content_bottom = min(y + h, cy + ch)
    row_h = min(
        0.50,
        max(0.30, (content_bottom - content_top) / max(1, len(bullets))),
    )
    for index, bullet in enumerate(bullets):
        _add_text(
            slide,
            str(bullet),
            (cx + 0.08, content_top + index * row_h, cw - 0.16, row_h),
            size=9,
            color=COLORS["black"],
            name=f"{module_id}_BULLET_{index + 1}",
            body_paragraph=True,
            bullet=True,
        )


def _render_process(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> None:
    module_id = module["module_id"]
    x, y, w, h = rect
    _add_text(
        slide,
        module.get("title", ""),
        (x, y, w, min(0.30, h * 0.30)),
        size=10.5,
        color=COLORS["deep_blue"],
        bold=True,
        name=f"{module_id}_TITLE",
    )
    steps = module.get("steps", [])
    gap = 0.06
    cell_w = (w - gap * max(0, len(steps) - 1)) / max(1, len(steps))
    sy = y + min(0.38, h * 0.34)
    sh = max(0.34, h - (sy - y) - 0.04)
    for index, step in enumerate(steps):
        sx = x + index * (cell_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.CHEVRON, Inches(sx), Inches(sy), Inches(cell_w), Inches(sh)
        )
        shape.name = f"{module_id}_STEP_BG_{index + 1}"
        shape.fill.solid()
        shape.fill.fore_color.rgb = (
            COLORS["deep_blue"] if index == 0 else COLORS["mid_blue"]
        )
        shape.line.fill.background()
        _add_text(
            slide,
            step,
            (sx + 0.05, sy + 0.02, cell_w - 0.12, sh - 0.04),
            size=9,
            color=COLORS["white"],
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"{module_id}_STEP_{index + 1}",
        )


def _render_matrix(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> None:
    module_id = module["module_id"]
    x, y, w, h = rect
    _add_text(
        slide,
        module.get("title", ""),
        (x, y, w, 0.32),
        size=10.5,
        color=COLORS["deep_blue"],
        bold=True,
        name=f"{module_id}_TITLE",
    )
    matrix = module.get("matrix", {})
    headers = matrix.get("headers", [])
    rows = matrix.get("rows", [])
    all_rows = [headers] + rows if headers else rows
    if not all_rows:
        return
    columns = max(len(row) for row in all_rows)
    cell_w = w / max(1, columns)
    cell_h = (h - 0.38) / max(1, len(all_rows))
    for row_index, row in enumerate(all_rows):
        for col_index in range(columns):
            value = str(row[col_index]) if col_index < len(row) else ""
            cell = (
                x + col_index * cell_w,
                y + 0.38 + row_index * cell_h,
                cell_w,
                cell_h,
            )
            header = row_index == 0 and bool(headers)
            _add_rect(
                slide,
                cell,
                fill=COLORS["deep_blue"] if header else COLORS["white"],
                line=COLORS["gray"],
                name=f"{module_id}_CELL_BG_{row_index + 1}_{col_index + 1}",
                line_width=0.5,
            )
            _add_text(
                slide,
                value,
                cell,
                size=8.5,
                color=COLORS["white"] if header else COLORS["black"],
                bold=header or col_index == 0,
                align=PP_ALIGN.CENTER,
                name=f"{module_id}_CELL_{row_index + 1}_{col_index + 1}",
            )


def _render_bar_chart(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> None:
    module_id = module["module_id"]
    x, y, w, h = rect
    _add_text(
        slide,
        module.get("title", ""),
        (x, y, w, 0.34),
        size=10.5,
        color=COLORS["deep_blue"],
        bold=True,
        name=f"{module_id}_TITLE",
    )
    chart = module.get("chart", {})
    categories = chart.get("categories", [])
    series = chart.get("series", [])
    values = series[0].get("values", []) if series else []
    if not categories or len(categories) != len(values):
        return
    left = x + 0.42
    bottom = y + h - 0.45
    top = y + 0.55
    right = x + w - 0.14
    _add_line(slide, left, top, left, bottom, color=COLORS["gray"], name=f"{module_id}_Y_AXIS")
    _add_line(slide, left, bottom, right, bottom, color=COLORS["gray"], name=f"{module_id}_X_AXIS")
    maximum = max(float(value) for value in values) or 1.0
    slot = (right - left) / len(values)
    bar_w = slot * 0.54
    for index, (category, raw_value) in enumerate(zip(categories, values)):
        value = float(raw_value)
        bar_h = (bottom - top - 0.12) * value / maximum
        bx = left + index * slot + (slot - bar_w) / 2
        by = bottom - bar_h
        _add_rect(
            slide,
            (bx, by, bar_w, bar_h),
            fill=COLORS["deep_blue"] if index < len(values) - 1 else COLORS["mid_blue"],
            line=None,
            name=f"{module_id}_BAR_{index + 1}",
        )
        _add_text(
            slide,
            f"{raw_value}",
            (bx - 0.08, by - 0.25, bar_w + 0.16, 0.22),
            size=8,
            color=COLORS["red"] if index == len(values) - 1 else COLORS["black"],
            bold=True,
            align=PP_ALIGN.CENTER,
            name=f"{module_id}_VALUE_{index + 1}",
        )
        _add_text(
            slide,
            str(category),
            (left + index * slot, bottom + 0.04, slot, 0.24),
            size=8,
            color=COLORS["dark_gray"],
            align=PP_ALIGN.CENTER,
            name=f"{module_id}_CATEGORY_{index + 1}",
        )
    if chart.get("unit"):
        _add_text(
            slide,
            str(chart["unit"]),
            (x + w - 0.7, y + 0.35, 0.64, 0.20),
            size=7,
            color=COLORS["dark_gray"],
            align=PP_ALIGN.RIGHT,
            name=f"{module_id}_UNIT",
        )


def _chart_frame(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> tuple[float, float, float, float]:
    module_id = module["module_id"]
    x, y, w, h = rect
    _add_text(
        slide,
        module.get("title", ""),
        (x, y, w, 0.34),
        size=10.5,
        color=COLORS["deep_blue"],
        bold=True,
        name=f"{module_id}_TITLE",
    )
    left, top, right, bottom = x + 0.42, y + 0.55, x + w - 0.14, y + h - 0.45
    _add_line(slide, left, top, left, bottom, color=COLORS["gray"], name=f"{module_id}_Y_AXIS")
    _add_line(slide, left, bottom, right, bottom, color=COLORS["gray"], name=f"{module_id}_X_AXIS")
    unit = module.get("chart", {}).get("unit")
    if unit:
        _add_text(slide, str(unit), (right - 0.65, y + 0.34, 0.65, 0.20), size=7, color=COLORS["dark_gray"], align=PP_ALIGN.RIGHT, name=f"{module_id}_UNIT")
    return left, top, right, bottom


def _add_point(
    slide: Any,
    x: float,
    y: float,
    diameter: float,
    *,
    fill: RGBColor,
    name: str,
) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x - diameter / 2),
        Inches(y - diameter / 2),
        Inches(diameter),
        Inches(diameter),
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = COLORS["white"]
    shape.line.width = Pt(0.5)
    return shape


def _render_line_chart(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> None:
    module_id = module["module_id"]
    chart = module.get("chart", {})
    categories = chart.get("categories", [])
    series = chart.get("series", [])
    values = series[0].get("values", []) if series else []
    if not categories or len(categories) != len(values):
        return
    left, top, right, bottom = _chart_frame(slide, module, rect)
    _add_text(slide, str(series[0].get("name", "")), (left, top - 0.28, 1.4, 0.20), size=8, color=COLORS["deep_blue"], bold=True, name=f"{module_id}_LINE_LEGEND")
    maximum, minimum = max(map(float, values)), min(map(float, values))
    span = maximum - minimum or 1.0
    points = []
    for index, (category, raw_value) in enumerate(zip(categories, values)):
        px = left + (right - left) * index / max(1, len(values) - 1)
        py = bottom - (bottom - top - 0.12) * (float(raw_value) - minimum) / span
        points.append((px, py))
        _add_point(slide, px, py, 0.12, fill=COLORS["deep_blue"], name=f"{module_id}_LINE_POINT_{index + 1}")
        _add_text(slide, str(raw_value), (px - 0.25, py - 0.25, 0.5, 0.20), size=8, color=COLORS["red"] if index == len(values) - 1 else COLORS["black"], bold=True, align=PP_ALIGN.CENTER, name=f"{module_id}_LINE_VALUE_{index + 1}")
        _add_text(slide, str(category), (px - 0.35, bottom + 0.04, 0.7, 0.22), size=8, color=COLORS["dark_gray"], align=PP_ALIGN.CENTER, name=f"{module_id}_LINE_CATEGORY_{index + 1}")
    for index, ((x1, y1), (x2, y2)) in enumerate(zip(points, points[1:]), start=1):
        _add_line(slide, x1, y1, x2, y2, color=COLORS["deep_blue"], width=1.8, name=f"{module_id}_LINE_SEGMENT_{index}")


def _render_stacked_chart(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> None:
    module_id = module["module_id"]
    chart = module.get("chart", {})
    categories = chart.get("categories", [])
    series = chart.get("series", [])
    if not categories or not series:
        return
    left, top, right, bottom = _chart_frame(slide, module, rect)
    totals = [sum(float(item.get("values", [0] * len(categories))[i]) for item in series) for i in range(len(categories))]
    maximum = max(totals) or 1.0
    slot = (right - left) / len(categories)
    bar_w = slot * 0.50
    palette = [COLORS["deep_blue"], COLORS["mid_blue"], COLORS["gray"]]
    for category_index, category in enumerate(categories):
        cursor = bottom
        for series_index, item in enumerate(series):
            value = float(item.get("values", [0] * len(categories))[category_index])
            height = (bottom - top - 0.12) * value / maximum
            cursor -= height
            bar_x = left + category_index * slot + (slot - bar_w) / 2
            _add_rect(slide, (bar_x, cursor, bar_w, height), fill=palette[series_index % len(palette)], line=None, name=f"{module_id}_STACK_{category_index + 1}_{series_index + 1}")
            _add_text(slide, str(item.get("values", [0] * len(categories))[category_index]), (bar_x, cursor + max(0.0, height / 2 - 0.10), bar_w, 0.20), size=8, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER, name=f"{module_id}_STACK_VALUE_{category_index + 1}_{series_index + 1}")
        _add_text(slide, str(category), (left + category_index * slot, bottom + 0.04, slot, 0.22), size=8, color=COLORS["dark_gray"], align=PP_ALIGN.CENTER, name=f"{module_id}_STACK_CATEGORY_{category_index + 1}")
    for index, item in enumerate(series):
        _add_text(slide, str(item.get("name", "")), (left + index * 1.1, top - 0.28, 1.05, 0.20), size=8, color=palette[index % len(palette)], bold=True, name=f"{module_id}_STACK_LEGEND_{index + 1}")


def _render_combo_chart(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> None:
    module_id = module["module_id"]
    chart = module.get("chart", {})
    categories = chart.get("categories", [])
    series = chart.get("series", [])
    if not categories or len(series) < 2:
        return
    bars, line_values = series[0].get("values", []), series[1].get("values", [])
    if len(bars) != len(categories) or len(line_values) != len(categories):
        return
    left, top, right, bottom = _chart_frame(slide, module, rect)
    for index, item in enumerate(series[:2]):
        _add_text(slide, str(item.get("name", "")), (left + index * 1.45, top - 0.28, 1.4, 0.20), size=8, color=COLORS["mid_blue" if index == 0 else "red"], bold=True, name=f"{module_id}_COMBO_LEGEND_{index + 1}")
    slot = (right - left) / len(categories)
    maximum = max(map(float, bars)) or 1.0
    line_max, line_min = max(map(float, line_values)), min(map(float, line_values))
    line_span = line_max - line_min or 1.0
    points = []
    for index, category in enumerate(categories):
        bar_h = (bottom - top - 0.12) * float(bars[index]) / maximum
        bx = left + index * slot + slot * 0.25
        _add_rect(slide, (bx, bottom - bar_h, slot * 0.5, bar_h), fill=COLORS["mid_blue"], line=None, name=f"{module_id}_COMBO_BAR_{index + 1}")
        px = left + index * slot + slot / 2
        py = bottom - (bottom - top - 0.12) * (float(line_values[index]) - line_min) / line_span
        points.append((px, py))
        _add_point(slide, px, py, 0.11, fill=COLORS["red"], name=f"{module_id}_COMBO_POINT_{index + 1}")
        _add_text(slide, str(bars[index]), (bx - 0.05, bottom - bar_h - 0.22, slot * 0.6, 0.20), size=8, color=COLORS["black"], bold=True, align=PP_ALIGN.CENTER, name=f"{module_id}_COMBO_BAR_VALUE_{index + 1}")
        _add_text(slide, str(line_values[index]), (px + 0.04, py - 0.20, 0.48, 0.20), size=8, color=COLORS["red"], bold=True, name=f"{module_id}_COMBO_LINE_VALUE_{index + 1}")
        _add_text(slide, str(category), (left + index * slot, bottom + 0.04, slot, 0.22), size=8, color=COLORS["dark_gray"], align=PP_ALIGN.CENTER, name=f"{module_id}_COMBO_CATEGORY_{index + 1}")
    for index, ((x1, y1), (x2, y2)) in enumerate(zip(points, points[1:]), start=1):
        _add_line(slide, x1, y1, x2, y2, color=COLORS["red"], width=1.6, name=f"{module_id}_COMBO_LINE_{index}")


def _render_scatter_or_bubble(
    slide: Any,
    module: dict[str, Any],
    rect: tuple[float, float, float, float],
    *,
    bubble: bool,
) -> None:
    module_id = module["module_id"]
    points = module.get("chart", {}).get("points", [])
    if not points:
        return
    left, top, right, bottom = _chart_frame(slide, module, rect)
    xs, ys = [float(item["x"]) for item in points], [float(item["y"]) for item in points]
    min_x, max_x, min_y, max_y = min(xs), max(xs), min(ys), max(ys)
    span_x, span_y = max_x - min_x or 1.0, max_y - min_y or 1.0
    max_size = max(float(item.get("size", 1)) for item in points) or 1.0
    prefix = "BUBBLE" if bubble else "SCATTER"
    plot_left, plot_right = left + 0.12, right - 0.95
    plot_top, plot_bottom = top + 0.16, bottom - 0.12
    for index, item in enumerate(points, start=1):
        px = plot_left + (plot_right - plot_left) * (float(item["x"]) - min_x) / span_x
        py = plot_bottom - (plot_bottom - plot_top) * (float(item["y"]) - min_y) / span_y
        diameter = 0.12 if not bubble else 0.12 + 0.24 * float(item.get("size", 1)) / max_size
        _add_point(slide, px, py, diameter, fill=COLORS["mid_blue" if bubble else "deep_blue"], name=f"{module_id}_{prefix}_POINT_{index}")
        details = f"{item.get('label', '')} ({item.get('x')},{item.get('y')}"
        if bubble:
            details += f";{item.get('size', 1)}"
        details += ")"
        _add_text(slide, details, (px + diameter / 2, py - 0.10, 0.85, 0.20), size=8, color=COLORS["black"], name=f"{module_id}_{prefix}_LABEL_{index}")


def _render_map(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> None:
    module_id = module["module_id"]
    x, y, w, h = rect
    _add_text(slide, module.get("title", ""), (x, y, w, 0.34), size=10.5, color=COLORS["deep_blue"], bold=True, name=f"{module_id}_TITLE")
    _add_rect(slide, (x, y + 0.40, w, h - 0.40), fill=None, line=COLORS["gray"], name=f"{module_id}_MAP_BASE")
    for index, item in enumerate(module.get("map", {}).get("labels", []), start=1):
        px = x + float(item.get("x", 0.5)) * w
        py = y + 0.40 + float(item.get("y", 0.5)) * (h - 0.40)
        _add_point(slide, px, py, 0.10, fill=COLORS["red"], name=f"{module_id}_MAP_POINT_{index}")
        _add_text(slide, str(item.get("label", "")), (px + 0.06, py - 0.10, 0.75, 0.20), size=8.5, color=COLORS["deep_blue"], bold=True, name=f"{module_id}_MAP_LABEL_{index}")


def _render_chart(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> None:
    role = str(module.get("role", ""))
    chart_type = str(module.get("chart", {}).get("type") or role)
    if chart_type in {"line", "line_chart"}:
        _render_line_chart(slide, module, rect)
    elif chart_type in {"stacked", "stacked_chart"}:
        _render_stacked_chart(slide, module, rect)
    elif chart_type in {"combo", "combo_chart"}:
        _render_combo_chart(slide, module, rect)
    elif chart_type in {"scatter", "scatter_chart"}:
        _render_scatter_or_bubble(slide, module, rect, bubble=False)
    elif chart_type in {"bubble", "bubble_chart"}:
        _render_scatter_or_bubble(slide, module, rect, bubble=True)
    else:
        _render_bar_chart(slide, module, rect)


def _render_generic(
    slide: Any, module: dict[str, Any], rect: tuple[float, float, float, float]
) -> None:
    module_id = module["module_id"]
    x, y, w, h = rect
    _add_rect(
        slide,
        rect,
        fill=_fill_color(module.get("fill_role")) or COLORS["white"],
        line=COLORS["mid_blue"],
        name=f"{module_id}_PANEL",
    )
    _add_text(
        slide,
        module.get("title", ""),
        (x + 0.10, y + 0.06, w - 0.20, 0.34),
        size=10.5,
        color=COLORS["deep_blue"],
        bold=True,
        name=f"{module_id}_TITLE",
    )
    content: list[str] = []
    content.extend(module.get("bullets", []))
    content.extend(module.get("steps", []))
    row_h = min(0.46, max(0.28, (h - 0.48) / max(1, len(content))))
    for index, value in enumerate(content):
        _add_text(
            slide,
            str(value),
            (x + 0.12, y + 0.46 + index * row_h, w - 0.24, row_h),
            size=8.5,
            color=COLORS["black"],
            name=f"{module_id}_ITEM_{index + 1}",
            body_paragraph=True,
            bullet=True,
        )


def _render_module(
    slide: Any,
    module: dict[str, Any],
    geometry_module: dict[str, Any],
) -> None:
    rect = _map_zone(geometry_module["zone_norm"])
    content_rect = None
    if _valid_zone(geometry_module.get("content_zone_norm")):
        content_rect = _map_zone(geometry_module["content_zone_norm"])
    _module_root(slide, module["module_id"], rect)
    role = module.get("role", "")
    primitive = geometry_module.get("primitive") or module.get("primitive")
    if role in ("metric_card", "card") or primitive == "card":
        _render_card(slide, module, rect, content_rect)
    elif role in ("bullet_group", "line_group") or primitive == "line_group":
        _render_line_group(slide, module, rect, content_rect)
    elif role == "process" or primitive == "flow_arrow":
        _render_process(slide, module, rect)
    elif role == "matrix" or primitive == "table_grid":
        _render_matrix(slide, module, rect)
    elif role in {
        "bar_chart",
        "line_chart",
        "stacked_chart",
        "combo_chart",
        "scatter_chart",
        "bubble_chart",
    } or primitive == "chart":
        _render_chart(slide, module, rect)
    elif role == "map" or primitive == "map":
        _render_map(slide, module, rect)
    else:
        _render_generic(slide, module, rect)


def _render_assets(
    slide: Any,
    geometry_slide: dict[str, Any],
    geometry_path: Path,
    *,
    layer: str,
) -> None:
    for asset in geometry_slide.get("assets", []):
        if asset.get("layer", "foreground") != layer:
            continue
        source = Path(asset.get("path", ""))
        if not source.is_absolute():
            source = geometry_path.parent / source
        if not source.exists():
            continue
        x, y, w, h = _map_zone(asset["zone_norm"])
        picture = slide.shapes.add_picture(str(source), Inches(x), Inches(y), Inches(w), Inches(h))
        picture.name = f"A_ASSET_{asset['asset_id']}"


def _effect_patterns() -> tuple[re.Pattern[bytes], ...]:
    tags = (
        b"effectRef",
        b"outerShdw",
        b"innerShdw",
        b"refl",
        b"glow",
        b"softEdge",
        b"scene3d",
        b"sp3d",
    )
    return tuple(
        re.compile(rb"<[^>]*" + tag + rb"\b[^>]*(?:/>|>.*?</[^>]+>)", re.DOTALL)
        for tag in tags
    )


def scrub_visual_effects(pptx_path: str | Path) -> None:
    pptx_path = Path(pptx_path)
    with tempfile.TemporaryDirectory(prefix="pptx_scrub_") as tmp:
        tmp_path = Path(tmp) / pptx_path.name
        with ZipFile(pptx_path, "r") as source, ZipFile(
            tmp_path, "w", ZIP_DEFLATED
        ) as target:
            for item in source.infolist():
                data = source.read(item.filename)
                if item.filename.startswith("ppt/slides/") and item.filename.endswith(".xml"):
                    for pattern in _effect_patterns():
                        data = pattern.sub(b"", data)
                target.writestr(item, data)
        shutil.move(str(tmp_path), str(pptx_path))


def build_deck(
    spec_path: str | Path,
    geometry_path: str | Path,
    output_path: str | Path,
    template_path: str | Path,
    *,
    brief_path: str | Path,
    blueprint_manifest_path: str | Path,
) -> Path:
    spec_path = Path(spec_path)
    geometry_path = Path(geometry_path)
    output_path = Path(output_path)
    template_path = Path(template_path)
    spec, geometry = load_project(spec_path, geometry_path)
    brief = load_json(brief_path)
    blueprint_manifest = load_json(blueprint_manifest_path)
    contract_errors = validate_contracts(
        spec,
        geometry,
        brief,
        blueprint_manifest,
        project_dir=spec_path.parent,
    )
    if contract_errors:
        raise ValueError("Contract validation failed:\n" + "\n".join(contract_errors))
    if not template_path.exists():
        raise FileNotFoundError(template_path)

    prs = Presentation(str(template_path))
    _remove_existing_slides(prs)
    layout = _blank_layout(prs)
    geometry_by_slide = {item["slide_id"]: item for item in geometry["slides"]}

    for page_number, slide_spec in enumerate(spec["slides"], start=1):
        slide = prs.slides.add_slide(layout)
        for placeholder in list(slide.placeholders):
            element = placeholder._element
            element.getparent().remove(element)
        _add_skeleton(slide, slide_spec, page_number)
        module_by_id = {item["module_id"]: item for item in slide_spec["modules"]}
        geometry_slide = geometry_by_slide[slide_spec["slide_id"]]
        _render_assets(slide, geometry_slide, geometry_path, layer="background")
        for geometry_module in geometry_slide["modules"]:
            _render_module(
                slide,
                module_by_id[geometry_module["module_id"]],
                geometry_module,
            )
        _render_assets(slide, geometry_slide, geometry_path, layer="foreground")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    scrub_visual_effects(output_path)
    return output_path


def _normalize_text(value: str) -> str:
    return re.sub(r"\s+", "", value or "")


def _format_rgb(format_object: Any) -> tuple[int, int, int] | None:
    try:
        rgb = format_object.fore_color.rgb
    except (AttributeError, TypeError, ValueError):
        return None
    return tuple(rgb) if rgb is not None else None


def _canonical_texts(slide_spec: dict[str, Any]) -> list[str]:
    texts = [
        slide_spec["chapter"],
        slide_spec["title"],
        slide_spec["source"],
    ]
    texts.extend(slide_spec.get("core_points", []))
    for module in slide_spec["modules"]:
        texts.append(module.get("title", ""))
        texts.extend(module.get("bullets", []))
        texts.extend(module.get("steps", []))
        for metric in module.get("metrics", []):
            texts.extend([metric.get("label", ""), metric.get("value", "")])
        matrix = module.get("matrix", {})
        texts.extend(matrix.get("headers", []))
        for row in matrix.get("rows", []):
            texts.extend(str(item) for item in row)
        chart = module.get("chart", {})
        texts.extend(str(item) for item in chart.get("categories", []))
        for series in chart.get("series", []):
            texts.append(series.get("name", ""))
            texts.extend(str(item) for item in series.get("values", []))
        for point in chart.get("points", []):
            texts.extend(
                str(point.get(field, ""))
                for field in ("label", "x", "y", "size")
                if field in point
            )
        if chart.get("unit"):
            texts.append(str(chart["unit"]))
        for label in module.get("map", {}).get("labels", []):
            texts.append(str(label.get("label", "")))
    return [text for text in texts if str(text).strip()]


def validate_pptx(
    pptx_path: str | Path,
    spec_path: str | Path,
    geometry_path: str | Path,
    brief_path: str | Path,
    blueprint_manifest_path: str | Path,
) -> list[str]:
    pptx_path = Path(pptx_path)
    spec, geometry = load_project(spec_path, geometry_path)
    brief = load_json(brief_path)
    blueprint_manifest = load_json(blueprint_manifest_path)
    errors = validate_contracts(
        spec,
        geometry,
        brief,
        blueprint_manifest,
        project_dir=Path(spec_path).parent,
    )
    if not pptx_path.exists():
        return errors + [f"pptx missing: {pptx_path}"]

    prs = Presentation(str(pptx_path))
    template_ref = spec.get("deck", {}).get("template")
    template_path = Path(spec_path).parent / template_ref if isinstance(template_ref, str) else None
    if template_path is None or not template_path.exists():
        template_path = Path(__file__).resolve().parents[1] / "assets" / "company_template.pptx"
    if template_path.exists():
        template_prs = Presentation(str(template_path))
        if (prs.slide_width, prs.slide_height) != (
            template_prs.slide_width,
            template_prs.slide_height,
        ):
            errors.append("pptx: slide size differs from company template")
        if [master.part.blob for master in prs.slide_masters] != [
            master.part.blob for master in template_prs.slide_masters
        ]:
            errors.append("pptx: slide masters differ from company template")
    if len(prs.slides) != len(spec["slides"]):
        errors.append(f"pptx: expected {len(spec['slides'])} slides, got {len(prs.slides)}")
        return errors

    geometry_by_slide = {item["slide_id"]: item for item in geometry["slides"]}
    for index, (slide, slide_spec) in enumerate(zip(prs.slides, spec["slides"]), start=1):
        sid = slide_spec["slide_id"]
        by_name = {shape.name: shape for shape in slide.shapes}
        required = {
            f"{sid}_SECTION_TITLE",
            f"{sid}_PAGE_TITLE_BAR",
            f"{sid}_PAGE_TITLE_TEXT",
            f"{sid}_CORE_BOX",
            f"{sid}_CORE_TEXT",
            f"{sid}_SOURCE",
            f"{sid}_PAGE_NUMBER",
        }
        required.update(item["module_id"] for item in slide_spec["modules"])
        for name in sorted(required - set(by_name)):
            errors.append(f"{sid}: missing shape {name}")

        expected_tops = {
            f"{sid}_SECTION_TITLE": Cm(0.35),
            f"{sid}_PAGE_TITLE_BAR": Cm(1.5),
            f"{sid}_CORE_BOX": Cm(2.85),
        }
        tolerance = Cm(0.01)
        for name, expected in expected_tops.items():
            if name in by_name and abs(by_name[name].top - expected) > tolerance:
                errors.append(f"{sid}: incorrect top for {name}")

        shared_left = Inches(0.45)
        for name in (
            f"{sid}_SECTION_TITLE",
            f"{sid}_PAGE_TITLE_TEXT",
            f"{sid}_CORE_TEXT",
        ):
            if name in by_name and abs(by_name[name].left - shared_left) > tolerance:
                errors.append(f"{sid}: {name} does not use the shared left edge")

        core_box = by_name.get(f"{sid}_CORE_BOX")
        if core_box is not None:
            core_line = _format_rgb(getattr(getattr(core_box, "line", None), "fill", None))
            if core_line != tuple(COLORS["black"]):
                errors.append(f"{sid}: core box border must be black")
            if core_box.line.width != Pt(1):
                errors.append(f"{sid}: core box border must be 1 pt")
            if core_box.line.dash_style != MSO_LINE_DASH_STYLE.DASH:
                errors.append(f"{sid}: core box border must use short dashes")

        core_text = by_name.get(f"{sid}_CORE_TEXT")
        if core_text is not None and getattr(core_text, "has_text_frame", False):
            core_paragraphs = [
                paragraph
                for paragraph in core_text.text_frame.paragraphs
                if paragraph.text.strip()
            ]
            if len(core_paragraphs) != len(slide_spec.get("core_points", [])):
                errors.append(f"{sid}: core paragraph count does not match core_points")
            if not 1 <= len(core_paragraphs) <= 3:
                errors.append(f"{sid}: core must contain 1 to 3 square-bullet points")
            for paragraph in core_paragraphs:
                if not paragraph.text.lstrip().startswith("■"):
                    errors.append(f"{sid}: core paragraph must start with a square bullet")
        if f"{sid}_CORE_LABEL" in by_name:
            errors.append(f"{sid}: core judgment label is forbidden")

        slide_text = "".join(
            shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
        )
        normalized_slide_text = _normalize_text(slide_text)
        for text in _canonical_texts(slide_spec):
            if _normalize_text(str(text)) not in normalized_slide_text:
                errors.append(f"{sid}: missing canonical text {text!r}")

        red_fill_area = 0
        source_shape = by_name.get(f"{sid}_SOURCE")
        source_top = source_shape.top if source_shape is not None else Inches(7.12)
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0:
                errors.append(f"{sid}: negative position {shape.name}")
            if shape.left + shape.width > prs.slide_width + 2 or shape.top + shape.height > prs.slide_height + 2:
                errors.append(f"{sid}: out of bounds {shape.name}")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not shape.name.startswith("A_ASSET_"):
                errors.append(f"{sid}: unexpected picture {shape.name}")
            for label, color in (
                ("fill", _format_rgb(getattr(shape, "fill", None))),
                ("line", _format_rgb(getattr(getattr(shape, "line", None), "fill", None))),
            ):
                if color is not None and not _is_allowed_rgb(color):
                    errors.append(f"{sid}: color outside company palette in {shape.name} {label}: {color}")
            fill_color = _format_rgb(getattr(shape, "fill", None))
            if fill_color == tuple(COLORS["red"]):
                red_fill_area += shape.width * shape.height
                upper_name = shape.name.upper()
                if any(
                    token in upper_name
                    for token in ("CARD", "HEADER", "TITLE", "PANEL", "CELL_BG")
                ):
                    errors.append(
                        f"{sid}: red fill is forbidden for card/header/background shape {shape.name}"
                    )
            is_source_rule = (
                shape.name == f"{sid}_SOURCE_RULE"
                or (
                    shape.width >= Inches(5)
                    and Inches(6.85) <= shape.top <= source_top
                    and (
                        shape.shape_type == MSO_SHAPE_TYPE.LINE
                        or shape.height <= Pt(3)
                    )
                )
            )
            if is_source_rule:
                errors.append(f"{sid}: source separator line is forbidden")
            if getattr(shape, "has_text_frame", False):
                allow_small = shape.name.endswith(("_SOURCE", "_PAGE_NUMBER", "_UNIT"))
                for paragraph in shape.text_frame.paragraphs:
                    if shape.name.endswith(("_CORE_TEXT",)) or "_BULLET_" in shape.name or "_ITEM_" in shape.name:
                        ppr = paragraph._p.get_or_add_pPr()
                        if ppr.get("marL") != str(Cm(0.64)) or ppr.get("indent") != str(-Cm(0.64)):
                            errors.append(f"{sid}: incorrect hanging indent in {shape.name}")
                        if paragraph.space_after != Pt(6) or paragraph.line_spacing != 1.2:
                            errors.append(f"{sid}: incorrect paragraph spacing in {shape.name}")
                        if paragraph.alignment != PP_ALIGN.JUSTIFY:
                            errors.append(f"{sid}: body paragraph is not justified in {shape.name}")
                    for run in paragraph.runs:
                        if run.text and run.font.size is None:
                            errors.append(f"{sid}: unset font size in {shape.name}")
                        elif run.text and run.font.size is not None:
                            minimum = 7 if allow_small else 8
                            if run.font.size.pt < minimum - 0.01:
                                errors.append(f"{sid}: font below {minimum} pt in {shape.name}")
                        if run.text and run.font.name not in (FONT, "微软雅黑"):
                            errors.append(f"{sid}: wrong font in {shape.name}")
                        run_color = None
                        try:
                            if run.font.color.rgb is not None:
                                run_color = tuple(run.font.color.rgb)
                        except (AttributeError, TypeError, ValueError):
                            pass
                        if run_color is not None and not _is_allowed_rgb(run_color):
                            errors.append(f"{sid}: font color outside company palette in {shape.name}: {run_color}")

        slide_area = prs.slide_width * prs.slide_height
        if slide_area and red_fill_area / slide_area > MAX_RED_FILL_RATIO:
            errors.append(
                f"{sid}: large red fill area exceeds {MAX_RED_FILL_RATIO:.0%} of the slide"
            )

        exact_sizes = {
            f"{sid}_SECTION_TITLE": 20,
            f"{sid}_PAGE_TITLE_TEXT": 16,
            f"{sid}_CORE_TEXT": 12,
            f"{sid}_SOURCE": 7,
            f"{sid}_PAGE_NUMBER": 8,
        }
        for name, expected_size in exact_sizes.items():
            shape = by_name.get(name)
            if shape is None or not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text and (run.font.size is None or abs(run.font.size.pt - expected_size) > 0.01):
                        errors.append(f"{sid}: {name} must be {expected_size} pt")

        expected_assets = {
            f"A_ASSET_{asset['asset_id']}"
            for asset in geometry_by_slide[sid].get("assets", [])
            if asset.get("required", True)
        }
        for name in expected_assets - set(by_name):
            errors.append(f"{sid}: missing required asset {name}")

    forbidden = (
        b"effectRef",
        b"outerShdw",
        b"innerShdw",
        b"refl",
        b"glow",
        b"softEdge",
        b"scene3d",
        b"sp3d",
    )
    with ZipFile(pptx_path) as archive:
        for name in archive.namelist():
            if not name.startswith("ppt/slides/") or not name.endswith(".xml"):
                continue
            data = archive.read(name)
            for token in forbidden:
                if token in data:
                    errors.append(f"pptx: forbidden visual effect {token.decode()} in {name}")

    return errors
