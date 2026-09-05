from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


SCHEMA_VERSION = "6.3"
ROLE_ORDER = (
    "chapter",
    "page_title",
    "core_judgment",
    "source",
    "page_number",
)
_PYTHON_PLACEHOLDERS = {
    PP_PLACEHOLDER.TITLE: "chapter",
    PP_PLACEHOLDER.BODY: "page_title",
    PP_PLACEHOLDER.OBJECT: "core_judgment",
    PP_PLACEHOLDER.FOOTER: "source",
    PP_PLACEHOLDER.SLIDE_NUMBER: "page_number",
}
_COM_PLACEHOLDERS = {int(key): value for key, value in _PYTHON_PLACEHOLDERS.items()}


def _sha256_text(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8")).hexdigest()


def _style_xml(shape) -> str:
    fragments = [shape.element.spPr.xml]
    tx_body = getattr(shape.element, "txBody", None)
    if tx_body is not None:
        body_pr = getattr(tx_body, "bodyPr", None)
        if body_pr is not None:
            fragments.append(body_pr.xml)
        for paragraph in getattr(tx_body, "p_lst", []):
            p_pr = getattr(paragraph, "pPr", None)
            if p_pr is not None:
                fragments.append(p_pr.xml)
            end_pr = getattr(paragraph, "endParaRPr", None)
            if end_pr is not None:
                fragments.append(end_pr.xml)
    return "".join(fragments)


def _shape_record(shape) -> dict[str, Any]:
    text_frame = shape.text_frame
    return {
        "shape_id": int(shape.shape_id),
        "name": str(shape.name),
        "placeholder_type": int(shape.placeholder_format.type),
        "geometry_emu": [
            int(shape.left),
            int(shape.top),
            int(shape.width),
            int(shape.height),
        ],
        "rotation": float(shape.rotation or 0),
        "text_frame": {
            "margin_left": int(text_frame.margin_left or 0),
            "margin_right": int(text_frame.margin_right or 0),
            "margin_top": int(text_frame.margin_top or 0),
            "margin_bottom": int(text_frame.margin_bottom or 0),
            "vertical_anchor": (
                int(text_frame.vertical_anchor)
                if text_frame.vertical_anchor is not None
                else None
            ),
            "word_wrap": text_frame.word_wrap,
            "auto_size": (
                int(text_frame.auto_size) if text_frame.auto_size is not None else None
            ),
        },
        "style_sha256": _sha256_text(_style_xml(shape)),
    }


def resolve_python_pptx_shapes(slide) -> dict[str, Any]:
    resolved: dict[str, Any] = {}
    for shape in slide.shapes:
        if not shape.is_placeholder:
            continue
        role = _PYTHON_PLACEHOLDERS.get(shape.placeholder_format.type)
        if role is None:
            continue
        if role in resolved:
            raise ValueError(f"V63_SKELETON_DUPLICATE_ROLE: {role}")
        resolved[role] = shape
    missing = [role for role in ROLE_ORDER if role not in resolved]
    if missing:
        raise ValueError("V63_SKELETON_MISSING_ROLE: " + ", ".join(missing))
    return {role: resolved[role] for role in ROLE_ORDER}


def _body_roi(shapes: dict[str, Any]) -> list[float]:
    page_title = shapes["page_title"]
    core = shapes["core_judgment"]
    source = shapes["source"]
    emu_per_inch = 914400.0
    left = float(page_title.left) / emu_per_inch
    top = float(core.top + core.height) / emu_per_inch + 0.08
    right = float(page_title.left + page_title.width) / emu_per_inch
    bottom = float(source.top) / emu_per_inch - 0.08
    if right <= left or bottom <= top:
        raise ValueError("V63_SKELETON_BODY_ROI_INVALID")
    return [left, top, right - left, bottom - top]


def read_template_contract(template_path: str | Path) -> dict[str, Any]:
    presentation = Presentation(str(Path(template_path).resolve()))
    if len(presentation.slides) < 1:
        raise ValueError("V63_SKELETON_TEMPLATE_EMPTY")
    shapes = resolve_python_pptx_shapes(presentation.slides[0])
    return {
        "schema_version": SCHEMA_VERSION,
        "roles": {role: _shape_record(shape) for role, shape in shapes.items()},
        "body_roi_in": _body_roi(shapes),
        "slide_size_in": [
            float(presentation.slide_width) / 914400.0,
            float(presentation.slide_height) / 914400.0,
        ],
    }


def _set_shape_text(shape, value: str) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    lines = str(value).splitlines() or [""]
    text_frame.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def update_pptx_skeleton(
    template_path: str | Path,
    output_path: str | Path,
    values: dict[str, str],
) -> Path:
    missing = [role for role in ROLE_ORDER if role not in values]
    if missing:
        raise ValueError("V63_SKELETON_TEXT_MISSING: " + ", ".join(missing))
    presentation = Presentation(str(Path(template_path).resolve()))
    shapes = resolve_python_pptx_shapes(presentation.slides[0])
    for role in ROLE_ORDER:
        _set_shape_text(shapes[role], str(values[role]))
    destination = Path(output_path).resolve()
    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(destination)
    return destination


def audit_pptx_skeleton(
    template_path: str | Path, candidate_path: str | Path
) -> dict[str, Any]:
    expected = read_template_contract(template_path)
    actual = read_template_contract(candidate_path)
    errors: list[str] = []
    for role in ROLE_ORDER:
        if actual["roles"][role] != expected["roles"][role]:
            errors.append(f"{role}: master-owned skeleton fingerprint changed")
    return {
        "schema_version": SCHEMA_VERSION,
        "ok": not errors,
        "errors": errors,
        "body_roi_in": actual["body_roi_in"],
    }


def resolve_com_shapes(slide) -> dict[str, Any]:
    resolved: dict[str, Any] = {}
    for index in range(1, int(slide.Shapes.Count) + 1):
        shape = slide.Shapes(index)
        try:
            placeholder_type = int(shape.PlaceholderFormat.Type)
        except Exception:
            continue
        role = _COM_PLACEHOLDERS.get(placeholder_type)
        if role is None:
            continue
        if role in resolved:
            raise ValueError(f"V63_SKELETON_DUPLICATE_ROLE: {role}")
        resolved[role] = shape
    missing = [role for role in ROLE_ORDER if role not in resolved]
    if missing:
        raise ValueError("V63_SKELETON_MISSING_ROLE: " + ", ".join(missing))
    return {role: resolved[role] for role in ROLE_ORDER}


def update_com_skeleton(slide, values: dict[str, str]) -> dict[str, Any]:
    missing = [role for role in ROLE_ORDER if role not in values]
    if missing:
        raise ValueError("V63_SKELETON_TEXT_MISSING: " + ", ".join(missing))
    shapes = resolve_com_shapes(slide)
    for role, shape in shapes.items():
        shape.TextFrame.TextRange.Text = str(values[role]).replace("\n", "\r")
    return shapes
