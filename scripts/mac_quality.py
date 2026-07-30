from __future__ import annotations

import hashlib
import json
from pathlib import Path
from zipfile import ZipFile

from lxml import etree
from pptx import Presentation


FORBIDDEN = {
    "effectLst", "effectDag", "scene3d", "sp3d",
    "glow", "reflection", "softEdge",
}


def _forbidden_effects(path: Path) -> list[str]:
    findings: list[str] = []
    with ZipFile(path) as archive:
        for name in archive.namelist():
            if not name.endswith(".xml") or not name.startswith(
                ("ppt/slides/", "ppt/slideMasters/", "ppt/slideLayouts/")
            ):
                continue
            root = etree.fromstring(archive.read(name))
            for element in root.iter():
                local_name = etree.QName(element).localname
                if local_name in FORBIDDEN:
                    findings.append(f"{name}:{local_name}")
    return findings


def _geometry_errors(presentation) -> list[str]:
    errors: list[str] = []
    width = int(presentation.slide_width)
    height = int(presentation.slide_height)
    for page_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            left = int(shape.left)
            top = int(shape.top)
            right = left + int(shape.width)
            bottom = top + int(shape.height)
            if left < 0 or top < 0 or right > width or bottom > height:
                errors.append(f"page {page_number}: {shape.name} is off slide")
    return errors


def _skeleton_errors(presentation) -> list[str]:
    required = {
        "SKEL_CHAPTER", "SKEL_TITLE", "SKEL_CORE",
        "SKEL_SOURCE", "SKEL_PAGE_NUMBER",
    }
    errors: list[str] = []
    for page_number, slide in enumerate(presentation.slides, start=1):
        names = {shape.name for shape in slide.shapes}
        for missing in sorted(required - names):
            errors.append(f"page {page_number}: missing {missing}")
    return errors


def _text_capacity_errors(project: Path) -> list[str]:
    path = project / ".build" / "mac_text_metrics_report.json"
    if not path.is_file():
        return ["Mac text metrics report is missing"]
    payload = json.loads(path.read_text(encoding="utf-8"))
    return [
        f"text capacity exceeded for {item.get('shape_name', '?')}"
        for item in payload.get("items", [])
        if isinstance(item, dict) and item.get("fits") is not True
    ]


def _asset_errors(presentation, project: Path) -> list[str]:
    path = project / ".build" / "mac_asset_report.json"
    if not path.is_file():
        return ["Mac asset report is missing"]
    payload = json.loads(path.read_text(encoding="utf-8"))
    picture_hashes: dict[str, list[str]] = {}
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                picture_hashes.setdefault(shape.name, []).append(
                    hashlib.sha256(shape.image.blob).hexdigest()
                )
    errors: list[str] = []
    for record in payload.get("assets", []):
        raw_names = record.get("shape_names")
        names = (
            [str(name) for name in raw_names if isinstance(name, str) and name]
            if isinstance(raw_names, list)
            else []
        )
        if not names:
            legacy_name = record.get("shape_name")
            names = [legacy_name] if isinstance(legacy_name, str) and legacy_name else []
        name = ", ".join(names)
        expected = str(record.get("source_sha256", ""))
        found = [
            digest
            for shape_name in names
            for digest in picture_hashes.get(shape_name, [])
        ]
        if len(found) != 1:
            errors.append(f"{name}: expected exactly one embedded picture")
        elif found[0] != expected:
            errors.append(
                f"{name}: embedded picture hash differs from extracted asset"
            )
    return errors


def audit_mac_pptx(
    pptx_path: str | Path,
    *,
    expected_page_count: int,
    render_result: dict,
    font_fallbacks: list[dict],
    project_dir: str | Path | None = None,
) -> dict:
    path = Path(pptx_path)
    errors: list[str] = []
    warnings: list[str] = []
    try:
        presentation = Presentation(path)
    except Exception as exc:
        return {
            "schema_version": "5.9",
            "builder_backend": "mac_python_pptx_v1",
            "status": "blocked",
            "ok": False,
            "errors": [f"PPTX reopen failed: {exc}"],
            "warnings": [],
        }
    if len(presentation.slides) != expected_page_count:
        errors.append(
            f"expected {expected_page_count} slides, got {len(presentation.slides)}"
        )
    errors.extend(_geometry_errors(presentation))
    if expected_page_count:
        errors.extend(_skeleton_errors(presentation))
    if project_dir is not None:
        project = Path(project_dir).resolve()
        errors.extend(_text_capacity_errors(project))
        errors.extend(_asset_errors(presentation, project))
    effects = _forbidden_effects(path)
    if effects:
        errors.append("forbidden effects remain: " + ", ".join(effects))
    if font_fallbacks:
        warnings.append("font fallback used")
    if render_result.get("renderer") == "libreoffice":
        warnings.append("visual verification used LibreOffice")
    if render_result.get("warning"):
        warnings.append(str(render_result["warning"]))
    if not render_result.get("visual_verification"):
        status = "structurally_valid_unrendered" if not errors else "blocked"
    elif errors:
        status = "blocked"
    elif warnings:
        status = "pass_with_warnings"
    else:
        status = "pass"
    return {
        "schema_version": "5.9",
        "builder_backend": "mac_python_pptx_v1",
        "render_backend": render_result.get("renderer"),
        "measurement_backend": "pillow_font_metrics",
        "audit_sources": [
            "python-pptx reopen",
            "DrawingML geometry and effects",
            "font metric allocation report",
            "embedded media hashes",
        ],
        "font_fallbacks": font_fallbacks,
        "visual_verification": bool(render_result.get("visual_verification")),
        "status": status,
        "ok": not errors,
        "errors": errors,
        "warnings": warnings,
    }
