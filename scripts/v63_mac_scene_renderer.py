from __future__ import annotations

import importlib.util
import json
import re
from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


SCHEMA_VERSION = "6.3"
RUNTIME_REVISION = "6.3.1"
_SHAPE_TYPES = {
    "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    "round_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
}
_ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}
_VERTICAL = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "center": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def _load(name: str):
    path = Path(__file__).with_name(f"{name}.py")
    spec = importlib.util.spec_from_file_location(f"v63_mac_{name}", path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def _read_json(path: Path) -> dict[str, Any] | list[Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def _write_json_atomic(path: Path, value: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_suffix(path.suffix + ".tmp")
    temporary.write_text(
        json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    temporary.replace(path)


def _rgb(value: str) -> RGBColor:
    token = str(value).lstrip("#")
    return RGBColor.from_string(token.upper())


def _remove_last_slide(presentation: Presentation) -> None:
    slide_id = presentation.slides._sldIdLst[-1]
    relationship_id = slide_id.rId
    presentation.part.drop_rel(relationship_id)
    del presentation.slides._sldIdLst[-1]


def _remove_body_shapes(slide) -> None:
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        shape._element.getparent().remove(shape._element)


def _skeleton_values(slide: dict[str, Any], page_number: int) -> dict[str, str]:
    raw_points = slide.get("core_points", [])
    if isinstance(raw_points, str):
        raw_points = [raw_points]
    points = []
    for point in raw_points if isinstance(raw_points, list) else []:
        value = re.sub(r"^(?:[■▪▫•●□]\s*)+", "", str(point).strip()).strip()
        points.append(value)
    return {
        "chapter": str(slide.get("chapter", "")),
        "page_title": str(slide.get("title", slide.get("page_title", ""))),
        "core_judgment": "\n".join(points),
        "source": str(slide.get("source", "")),
        "page_number": str(slide.get("page_number", page_number)),
    }


def _map_point(
    point: list[float], body_roi_px: list[float], body_roi_in: list[float]
) -> list[float]:
    x, y = [float(value) for value in point]
    body_x, body_y, body_width, body_height = [float(value) for value in body_roi_px]
    slide_x, slide_y, slide_width, slide_height = [float(value) for value in body_roi_in]
    return [
        slide_x + (x - body_x) / body_width * slide_width,
        slide_y + (y - body_y) / body_height * slide_height,
    ]


def _apply_fill_and_line(shape, style: dict[str, Any]) -> None:
    fill = style.get("fill")
    if fill in {None, "none"}:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(str(fill))
        if style.get("transparency") is not None:
            shape.fill.transparency = float(style["transparency"])
    line = style.get("line")
    if line in {None, "none"}:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(str(line))
        shape.line.width = Pt(float(style.get("line_width", 0.8)))


def _set_arrowhead(shape, command: dict[str, Any]) -> None:
    if command.get("type") != "arrow" and not command.get("style", {}).get("arrow_end"):
        return
    line = shape._element.spPr.get_or_add_ln()
    for child in list(line):
        if child.tag == qn("a:tailEnd"):
            line.remove(child)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    line.append(tail)


def _render_text(slide, command: dict[str, Any]):
    x, y, width, height = command["bbox_in"]
    style = command.get("style", {})
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    _apply_fill_and_line(shape, style)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = bool(style.get('word_wrap', True))
    margin = float(style.get("margin", 1.5))
    frame.margin_left = Pt(float(style.get("margin_left", margin)))
    frame.margin_right = Pt(float(style.get("margin_right", margin)))
    frame.margin_top = Pt(float(style.get("margin_top", 0.8)))
    frame.margin_bottom = Pt(float(style.get("margin_bottom", 0.8)))
    frame.vertical_anchor = _VERTICAL.get(str(style.get("valign", "top")).lower(), MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = _ALIGNMENTS.get(str(style.get("align", "left")).lower(), PP_ALIGN.LEFT)
    runs = command.get("runs") if isinstance(command.get("runs"), list) else None
    run_payloads = runs if runs is not None else [{"text": str(command.get("text", ""))}]
    for run_index, payload in enumerate(run_payloads):
        if not isinstance(payload, dict):
            continue
        run = paragraph.add_run()
        run.text = str(payload.get("text", ""))
        run_style = {**style, **(payload.get("style", {}) if isinstance(payload.get("style"), dict) else {})}
        run.font.name = str(run_style.get("font_name", "Microsoft YaHei"))
        run.font.size = Pt(float(run_style.get("font_size", 10.0)))
        run.font.bold = bool(run_style.get("bold"))
        run.font.color.rgb = _rgb(str(run_style.get("color", "#111111")))
    if style.get('fit') == 'shrink_to_box':
        # Font metrics are unavailable here. Ask Office to fit on opening;
        # this is NOT evidence of Mac native text/visual verification.
        from pptx.enum.text import MSO_AUTO_SIZE
        frame.word_wrap = False
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return shape


def _render_shape(slide, command: dict[str, Any]):
    x, y, width, height = command["bbox_in"]
    shape = slide.shapes.add_shape(
        _SHAPE_TYPES[command["type"]], Inches(x), Inches(y), Inches(width), Inches(height)
    )
    _apply_fill_and_line(shape, command.get("style", {}))
    return shape


def _render_line(slide, command: dict[str, Any]):
    points = command.get("points_in")
    if isinstance(points, list) and len(points) > 2:
        shape = _render_freeform(slide, {**command, 'closed': False})
        _set_arrowhead(shape, command)
        return shape
    if not isinstance(points, list) or len(points) < 2:
        x, y, width, height = command["bbox_in"]
        points = [[x, y], [x + width, y + height]]
    start, end = points[0], points[-1]
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start[0]),
        Inches(start[1]),
        Inches(end[0]),
        Inches(end[1]),
    )
    style = command.get("style", {})
    shape.line.color.rgb = _rgb(str(style.get("line", "#111111")))
    shape.line.width = Pt(float(style.get("line_width", 1.0)))
    _set_arrowhead(shape, command)
    return shape


def _render_freeform(slide, command: dict[str, Any]):
    points = command.get("points_in")
    if not isinstance(points, list) or len(points) < (3 if command.get('closed', True) else 2):
        raise ValueError(f"MAC_RECONSTRUCTION_UNSUPPORTED: {command.get('element_id')} freeform")
    emu_points = [(Inches(point[0]), Inches(point[1])) for point in points]
    builder = slide.shapes.build_freeform(emu_points[0][0], emu_points[0][1])
    remaining = emu_points[1:-1] if emu_points[-1] == emu_points[0] else emu_points[1:]
    builder.add_line_segments(remaining, close=command.get('closed', True))
    shape = builder.convert_to_shape()
    _apply_fill_and_line(shape, command.get("style", {}))
    _set_arrowhead(shape, command)
    return shape


def _render_picture(slide, command: dict[str, Any], asset_path: Path):
    x, y, width, height = command["bbox_in"]
    shape = slide.shapes.add_picture(
        str(asset_path), Inches(x), Inches(y), width=Inches(width), height=Inches(height)
    )
    shape.line.fill.background()
    return shape


def _render_command(slide, command: dict[str, Any], assets: dict[str, Path]):
    kind = command["type"]
    if kind == "text":
        shape = _render_text(slide, command)
    elif kind in _SHAPE_TYPES:
        shape = _render_shape(slide, command)
    elif kind in {"line", "connector", "arrow"}:
        shape = _render_line(slide, command)
    elif kind == "freeform":
        shape = _render_freeform(slide, command)
    elif kind == "image_crop":
        asset_id = str(command.get("asset_id", ""))
        if asset_id not in assets:
            raise ValueError(f"MAC_ASSET_CONTRACT_MISMATCH: {asset_id}")
        shape = _render_picture(slide, command, assets[asset_id])
    else:
        raise ValueError(f"MAC_RECONSTRUCTION_UNSUPPORTED: {kind}")
    shape.name = f"V63_{command['element_id']}"
    if command.get("style", {}).get("rotation") is not None:
        shape.rotation = float(command["style"]["rotation"])
    return shape


def render_plan(page: dict[str, Any], body_roi_in: list[float]) -> list[dict[str, Any]]:
    mapper = _load("v63_scene_graph")
    plan: list[dict[str, Any]] = []
    for element in sorted(page.get("elements", []), key=lambda item: item.get("z_order", 0)):
        if element.get("type") == "group":
            continue
        command = mapper.normalize_element_geometry(element)
        command["bbox_in"] = mapper.pixel_box_to_slide_box(
            command["bbox_px"], page["body_roi_px"], body_roi_in,
            coordinate_mode=page.get('coordinate_mode', 'legacy_stretch'),
            allow_line=command['type'] in {'line', 'connector', 'arrow'} or not command.get('closed', True),
        )
        if isinstance(command.get("points_px"), list):
            command["points_in"] = [
                (mapper.map_source_point(point, mapper.body_contain_transform(page['body_roi_px'], body_roi_in))
                 if page.get('coordinate_mode') == 'source_pixels_contain'
                 else _map_point(point, page["body_roi_px"], body_roi_in))
                for point in command["points_px"]
            ]
        plan.append(command)
    return plan


def build_deck(
    project_dir: str | Path,
    output_path: str | Path,
    *,
    template_path: str | Path,
) -> Path:
    project = Path(project_dir).resolve()
    output = Path(output_path).resolve()
    template = Path(template_path).resolve()
    scene_graph = _read_json(project / ".build" / "v63_scene_graph.json")
    slides = _read_json(project / ".build" / "slides.json")
    ledger = _read_json(project / ".build" / "v63_asset_ledger.json")
    if not isinstance(slides, list):
        raise ValueError("V63_SLIDES_INVALID")
    slide_by_id = {str(item.get("slide_id")): item for item in slides if isinstance(item, dict)}
    page_ids = sorted(scene_graph.get("pages", {}))
    if set(page_ids) != set(slide_by_id):
        raise ValueError("V63_SCENE_SLIDE_SET_MISMATCH")
    assets = {
        str(item["asset_id"]): (project / str(item["asset_path"])).resolve()
        for item in ledger.get("assets", [])
        if isinstance(item, dict) and item.get("asset_id") and item.get("asset_path")
    }
    skeleton = _load("v63_skeleton_contract")
    contract = skeleton.read_template_contract(template)
    presentation = Presentation(str(template))
    first_layout = presentation.slides[0].slide_layout
    while len(presentation.slides) < len(page_ids):
        added = presentation.slides.add_slide(first_layout)
        # python-pptx intentionally skips footer/date/page-number placeholders
        # when adding from a layout. Clone the seed's original text placeholders
        # so every page has exactly the same five master-owned shape contracts.
        for shape in list(added.shapes):
            shape._element.getparent().remove(shape._element)
        for shape in presentation.slides[0].shapes:
            if shape.is_placeholder:
                added.shapes._spTree.insert_element_before(deepcopy(shape._element), 'p:extLst')
    while len(presentation.slides) > len(page_ids):
        _remove_last_slide(presentation)
    objects: list[dict[str, Any]] = []
    for page_index, slide_id in enumerate(page_ids, start=1):
        slide = presentation.slides[page_index - 1]
        _remove_body_shapes(slide)
        shapes = skeleton.resolve_python_pptx_shapes(slide)
        for role, value in _skeleton_values(slide_by_id[slide_id], page_index).items():
            skeleton._set_shape_text(shapes[role], value)
        for command in render_plan(scene_graph["pages"][slide_id], contract["body_roi_in"]):
            shape = _render_command(slide, command, assets)
            objects.append(
                {
                    "slide_id": slide_id,
                    "element_id": command["element_id"],
                    "shape_name": shape.name,
                    "type": command["type"],
                    "editable": command["type"] != "image_crop",
                    "asset_id": command.get("asset_id"),
                }
            )
    output.parent.mkdir(parents=True, exist_ok=True)
    temporary = output.with_name(f".{output.stem}.building.pptx")
    if temporary.is_file():
        temporary.unlink()
    presentation.save(temporary)
    temporary.replace(output)
    report = {
        "schema_version": SCHEMA_VERSION,
        "deconstruction_runtime_revision": RUNTIME_REVISION,
        "builder_backend": "mac_python_pptx_v2",
        "ok": True,
        "status": "structurally_valid_unrendered",
        "mac_native_render_unverified": True,
        "pptx": str(output),
        "object_count": len(objects),
        "editable_body_count": sum(item["editable"] for item in objects),
        "image_count": sum(not item["editable"] for item in objects),
        "objects": objects,
    }
    _write_json_atomic(project / ".build" / "v63_mac_build_report.json", report)
    _write_json_atomic(project / ".build" / "v63_object_ledger.json", report)
    return output
