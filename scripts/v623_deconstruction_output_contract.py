from __future__ import annotations

import hashlib
import json
import re
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt


SCHEMA_VERSION = "6.2.3"
TEMPLATE_MISMATCH = "D623_TEMPLATE_MISMATCH"
SKELETON_MISMATCH = "D623_SKELETON_MISMATCH"
FORBIDDEN_EFFECT = "D623_FORBIDDEN_EFFECT"
_EMU_PER_INCH = 914400.0
_POSITION_TOLERANCE = 0.025
_SKELETON = {
    "SKEL_CHAPTER", "SKEL_TITLE", "SKEL_CORE", "SKEL_SOURCE", "SKEL_PAGE_NUMBER"
}
_CORE_POINT = re.compile(r"^■\s+(?![■▪▫•●□])\S")
_ACTIVE_EFFECTS = {
    "outerShdw", "innerShdw", "glow", "reflection", "softEdge", "scene3d", "sp3d"
}


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _issue(code: str, message: str, *, slide_id: str = "", scope: str = "") -> dict[str, str]:
    return {
        "code": code,
        "severity": "blocker",
        "stage": "deconstruction_output_contract",
        "slide_id": slide_id,
        "scope": scope,
        "message": message,
    }


def _layout_inventory(presentation: Presentation) -> list[list[str]]:
    return [
        [str(layout.name) for layout in master.slide_layouts]
        for master in presentation.slide_masters
    ]


def _theme_signatures(presentation: Presentation) -> list[list[tuple[str, tuple, str]]]:
    signatures = []
    semantic_names = {
        "clrScheme", "fontScheme", "dk1", "lt1", "dk2", "lt2", "accent1",
        "accent2", "accent3", "accent4", "accent5", "accent6", "hlink",
        "folHlink", "latin", "ea", "cs", "srgbClr", "sysClr", "schemeClr",
    }
    for part in presentation.part.package.iter_parts():
        if "/theme/" not in str(part.partname):
            continue
        root = etree.fromstring(part.blob)
        signatures.append(
            [
                (
                    etree.QName(node).localname,
                    tuple(sorted(node.attrib.items())),
                    node.text or "",
                )
                for node in root.iter()
                if etree.QName(node).localname in semantic_names
            ]
        )
    return signatures


def _shape_by_name(slide, name: str):
    return next((shape for shape in slide.shapes if shape.name == name), None)


def _inches(value: int) -> float:
    return float(value) / _EMU_PER_INCH


def _near(actual: int, expected: float) -> bool:
    return abs(_inches(actual) - expected) <= _POSITION_TOLERANCE


def _font_sizes(shape) -> list[float]:
    sizes: list[float] = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.size is not None:
                sizes.append(float(run.font.size.pt))
    return sizes


def _rgb(shape) -> RGBColor | None:
    try:
        return shape.line.color.rgb
    except (AttributeError, TypeError, ValueError):
        return None


def _dash(shape) -> str | None:
    line = shape._element.spPr.find(qn("a:ln"))
    if line is None:
        return None
    dash = line.find(qn("a:prstDash"))
    return dash.get("val") if dash is not None else None


def _skeleton_blockers(slide, slide_id: str) -> list[dict[str, str]]:
    blockers: list[dict[str, str]] = []
    shapes = {name: _shape_by_name(slide, name) for name in _SKELETON}
    for name, shape in shapes.items():
        if shape is None or not getattr(shape, "has_text_frame", False):
            blockers.append(_issue(SKELETON_MISMATCH, f"missing native skeleton shape {name}", slide_id=slide_id))
    if blockers:
        return blockers

    expected_tops = {
        "SKEL_CHAPTER": 0.1575,
        "SKEL_TITLE": 0.5906,
        "SKEL_CORE": 1.09 if _shape_by_name(slide, "SKEL_CORE_BORDER") else 1.063,
    }
    for name, expected in expected_tops.items():
        if not _near(shapes[name].top, expected):
            blockers.append(_issue(SKELETON_MISMATCH, f"{name} top must be {expected:.4f}in", slide_id=slide_id))
    lefts = [_inches(shapes[name].left) for name in ("SKEL_CHAPTER", "SKEL_TITLE", "SKEL_CORE")]
    if max(lefts) - min(lefts) > 0.085:
        blockers.append(_issue(SKELETON_MISMATCH, "chapter, title, and core must use the template anchor", slide_id=slide_id))

    for name, expected in (("SKEL_CHAPTER", 20), ("SKEL_TITLE", 16), ("SKEL_CORE", 12)):
        sizes = _font_sizes(shapes[name])
        if not sizes or any(abs(value - expected) > 0.1 for value in sizes):
            blockers.append(_issue(SKELETON_MISMATCH, f"{name} font size must be {expected}pt", slide_id=slide_id))

    core = shapes["SKEL_CORE"]
    paragraphs = [p.text.strip() for p in core.text_frame.paragraphs if p.text.strip()]
    if not paragraphs or any(not _CORE_POINT.match(text) for text in paragraphs):
        blockers.append(_issue(SKELETON_MISMATCH, "each core paragraph must contain exactly one square bullet", slide_id=slide_id))
    if any(p.alignment not in (None, PP_ALIGN.LEFT) for p in core.text_frame.paragraphs if p.text.strip()):
        blockers.append(_issue(SKELETON_MISMATCH, "core judgment must be left aligned", slide_id=slide_id))

    border = _shape_by_name(slide, "SKEL_CORE_BORDER") or core
    width_pt = float(border.line.width.pt) if border.line.width is not None else 0.0
    if _rgb(border) != RGBColor(0, 0, 0) or abs(width_pt - 1.0) > 0.05 or _dash(border) != "dash":
        blockers.append(_issue(SKELETON_MISMATCH, "core border must be black, 1pt, dashed, and effect-free", slide_id=slide_id))
    return blockers


def _effect_blockers(root, scope: str, slide_id: str = "") -> list[dict[str, str]]:
    blockers: list[dict[str, str]] = []
    for element in root.iter():
        local = etree.QName(element).localname
        active = local in _ACTIVE_EFFECTS
        if local == "effectRef":
            try:
                active = int(element.get("idx", "0")) != 0
            except ValueError:
                active = True
        if active:
            blockers.append(_issue(FORBIDDEN_EFFECT, f"forbidden {local} in {scope}", slide_id=slide_id, scope=scope))
    return blockers


def audit_deconstruction_output(pptx_path: str | Path, template_path: str | Path) -> dict[str, Any]:
    pptx = Path(pptx_path).resolve()
    template = Path(template_path).resolve()
    blockers: list[dict[str, str]] = []
    if not pptx.is_file() or not template.is_file():
        missing = pptx if not pptx.is_file() else template
        blockers.append(_issue(TEMPLATE_MISMATCH, f"missing file: {missing}"))
        return _result(pptx, template, blockers, 0)

    actual = Presentation(pptx)
    expected = Presentation(template)
    if (actual.slide_width, actual.slide_height) != (expected.slide_width, expected.slide_height):
        blockers.append(_issue(TEMPLATE_MISMATCH, "slide size does not exactly match company template"))
    if (
        _layout_inventory(actual) != _layout_inventory(expected)
        or _theme_signatures(actual) != _theme_signatures(expected)
    ):
        blockers.append(_issue(TEMPLATE_MISMATCH, "master/layout/theme does not match company template"))
    if not expected.slides:
        blockers.append(_issue(TEMPLATE_MISMATCH, "company template has no seed slide"))
        expected_layout = None
    else:
        expected_layout = expected.slides[0].slide_layout.name

    for index, slide in enumerate(actual.slides, 1):
        slide_id = f"S{index:02d}"
        if expected_layout is not None and slide.slide_layout.name != expected_layout:
            blockers.append(_issue(TEMPLATE_MISMATCH, f"slide must inherit seed layout {expected_layout!r}", slide_id=slide_id))
        blockers.extend(_skeleton_blockers(slide, slide_id))
        blockers.extend(_effect_blockers(slide.element, f"slide:{slide_id}", slide_id))
    for master_index, master in enumerate(actual.slide_masters, 1):
        blockers.extend(_effect_blockers(master.element, f"master:{master_index}"))
        for layout_index, layout in enumerate(master.slide_layouts, 1):
            blockers.extend(_effect_blockers(layout.element, f"layout:{master_index}:{layout_index}"))
    return _result(pptx, template, blockers, len(actual.slides))


def _result(pptx: Path, template: Path, blockers: list[dict[str, str]], page_count: int) -> dict[str, Any]:
    return {
        "schema_version": SCHEMA_VERSION,
        "construction_mode": "deconstruct",
        "status": "blocked" if blockers else "pass",
        "ok": not blockers,
        "pptx": str(pptx),
        "pptx_sha256": _sha256(pptx) if pptx.is_file() else None,
        "template": str(template),
        "template_sha256": _sha256(template) if template.is_file() else None,
        "page_count": page_count,
        "blockers": blockers,
        "blocker_count": len(blockers),
        "warnings": [],
        "warning_count": 0,
    }


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("pptx")
    parser.add_argument("template")
    parser.add_argument("--output")
    arguments = parser.parse_args()
    report = audit_deconstruction_output(arguments.pptx, arguments.template)
    payload = json.dumps(report, ensure_ascii=False, indent=2) + "\n"
    if arguments.output:
        Path(arguments.output).write_text(payload, encoding="utf-8")
    else:
        print(payload, end="")
    raise SystemExit(0 if report["ok"] else 1)
