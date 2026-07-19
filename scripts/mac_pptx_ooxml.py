from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile

from lxml import etree
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn


def remove_shape(shape) -> None:
    parent = shape._element.getparent()
    if parent is not None:
        parent.remove(shape._element)


def remove_last_slide(presentation) -> None:
    if not presentation.slides:
        raise ValueError("cannot remove a slide from an empty presentation")
    slide_id = presentation.slides._sldIdLst[-1]
    presentation.part.drop_rel(slide_id.rId)
    del presentation.slides._sldIdLst[-1]


def set_shape_name(shape, name: str) -> None:
    for element in shape._element.iter():
        if etree.QName(element).localname == "cNvPr":
            element.set("name", name)
            return
    raise ValueError("shape has no cNvPr name node")


def set_east_asian_font(run, name: str) -> None:
    rpr = run._r.get_or_add_rPr()
    east_asian = rpr.find(qn("a:ea"))
    if east_asian is None:
        east_asian = etree.SubElement(rpr, qn("a:ea"))
    east_asian.set("typeface", name)


def set_line_dash(shape, dash: str) -> None:
    allowed = {"solid", "dash", "dashDot", "lgDash", "sysDot"}
    if dash not in allowed:
        raise ValueError(f"unsupported dash style: {dash}")
    line = shape._element.spPr.get_or_add_ln()
    for child in list(line):
        if child.tag == qn("a:prstDash"):
            line.remove(child)
    line.append(parse_xml(f'<a:prstDash {nsdecls("a")} val="{dash}"/>'))


def set_arrow_end(shape, value: str = "triangle") -> None:
    line = shape._element.spPr.get_or_add_ln()
    for child in list(line):
        if child.tag == qn("a:tailEnd"):
            line.remove(child)
    line.append(parse_xml(f'<a:tailEnd {nsdecls("a")} type="{value}" w="med" len="med"/>'))


FORBIDDEN_LOCAL_NAMES = {
    "effectLst", "effectDag", "scene3d", "sp3d", "glow", "reflection", "softEdge"
}


def clear_forbidden_effects(root) -> int:
    removed = 0
    for element in list(root.iter()):
        if etree.QName(element).localname in FORBIDDEN_LOCAL_NAMES:
            parent = element.getparent()
            if parent is not None:
                parent.remove(element)
                removed += 1
    return removed


def validate_pptx_package(path: str | Path) -> dict:
    path = Path(path)
    with ZipFile(path) as archive:
        bad_member = archive.testzip()
        if bad_member:
            raise ValueError(f"PPTX ZIP integrity failed at {bad_member}")
        etree.fromstring(archive.read("ppt/presentation.xml"))
    presentation = Presentation(path)
    return {
        "ok": True, "slides": len(presentation.slides),
        "width": presentation.slide_width, "height": presentation.slide_height,
    }
