"""Audit V5.6 complete one-object bitmap placement in a generated PowerPoint deck."""

from __future__ import annotations

import argparse
import hashlib
import importlib.util
import json
from pathlib import Path


SCHEMA_VERSION = "5.6"
TOLERANCE_IN = 0.02
ASPECT_TOLERANCE = 0.02
EMU_PER_POINT = 12700.0


def sha256_file(path: str | Path) -> str:
    digest = hashlib.sha256()
    with Path(path).open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _load_extractor():
    path = Path(__file__).with_name("extract_direct_assets.py")
    spec = importlib.util.spec_from_file_location("standard_report_direct_assets", path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def audit_manifest(
    manifest: dict,
    asset_crops: dict,
    asset_report: dict,
    *,
    census_crop_ids: set[str] | None = None,
) -> dict:
    errors: list[str] = []
    report_by_id = {record.get("asset_id"): record for record in asset_report.get("assets", []) if isinstance(record, dict)}
    pages = {page.get("slide_id"): page for page in manifest.get("pages", []) if isinstance(page, dict)}
    for asset_id, crop in asset_crops.items():
        slide_id = crop.get("slide_id")
        page = pages.get(slide_id)
        if page is None:
            errors.append(f"{asset_id}: slide {slide_id} is missing from PPT manifest")
            continue
        shapes = page.get("assets", {}).get(asset_id, [])
        if len(shapes) != 1:
            errors.append(f"{asset_id}: expected exactly one named picture, found {len(shapes)}")
            continue
        record = report_by_id.get(asset_id)
        if not record:
            errors.append(f"{asset_id}: extraction report entry is missing")
            continue
        shape = shapes[0]
        left, top, width, height = [float(shape[key]) / 72.0 for key in ("left", "top", "width", "height")]
        if height <= 0 or width <= 0:
            errors.append(f"{asset_id}: picture dimensions must be positive")
            continue
        expected_aspect = float(record.get("aspect_ratio", 0))
        actual_aspect = width / height
        if expected_aspect <= 0 or abs(actual_aspect / expected_aspect - 1.0) > ASPECT_TOLERANCE:
            errors.append(f"{asset_id}: picture aspect ratio does not match extracted PNG")
        box_x, box_y, box_width, box_height = [float(value) for value in crop.get("target_box_in", [])]
        if crop.get("target_coord_space") == "body":
            box_x += 0.56
            box_y += float(page.get("core_bottom", 0)) / 72.0 + 0.12
        if left < box_x - TOLERANCE_IN or top < box_y - TOLERANCE_IN or left + width > box_x + box_width + TOLERANCE_IN or top + height > box_y + box_height + TOLERANCE_IN:
            errors.append(f"{asset_id}: picture is outside its target box")
        core_bottom = float(page.get("core_bottom", 0)) / 72.0
        footer_top = float(page.get("footer_top", 99999)) / 72.0
        if top < core_bottom - TOLERANCE_IN or top + height > footer_top + TOLERANCE_IN:
            errors.append(f"{asset_id}: picture must stay inside the body between core and footer")
    declared = set(asset_crops)
    expected_from_census = set(census_crop_ids or set())
    for missing in sorted(expected_from_census - declared):
        errors.append(f"{missing}: crop census subject is missing from ASSET_CROPS")
    for extra in sorted(declared - expected_from_census):
        if census_crop_ids is not None:
            errors.append(f"{extra}: ASSET_CROPS entry is missing from the crop census")
    found = {
        asset_id
        for page in pages.values()
        for asset_id in page.get("assets", {})
    }
    for extra in sorted(found - declared):
        errors.append(f"{extra}: undeclared ASSET_ picture found")
    inserted_declared = found & declared
    complete_inventory = not errors and inserted_declared == declared
    return {
        "schema_version": SCHEMA_VERSION,
        "ok": complete_inventory,
        "errors": errors,
        "declared_assets": len(declared),
        "inserted_assets": len(inserted_declared),
        "complete_inventory": complete_inventory,
        "assets": len(declared),
        "census_crop_assets": len(expected_from_census),
    }


def extract_manifest(pptx_path: str | Path) -> dict:
    import win32com.client

    app = win32com.client.DispatchEx("PowerPoint.Application")
    presentation = app.Presentations.Open(str(Path(pptx_path).resolve()), False, True, False)
    pages = []
    try:
        for page_no in range(1, presentation.Slides.Count + 1):
            slide = presentation.Slides(page_no)
            assets: dict[str, list[dict]] = {}
            core_bottom = 0.0
            footer_top = float(presentation.PageSetup.SlideHeight)
            for index in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(index)
                name = str(shape.Name)
                if name == "SKEL_CORE":
                    core_bottom = float(shape.Top) + float(shape.Height)
                elif name == "SKEL_SOURCE":
                    footer_top = float(shape.Top)
                elif name.startswith("ASSET_"):
                    asset_id = name[len("ASSET_"):]
                    assets.setdefault(asset_id, []).append({
                        "left": float(shape.Left),
                        "top": float(shape.Top),
                        "width": float(shape.Width),
                        "height": float(shape.Height),
                    })
            pages.append({"slide_id": f"S{page_no:02d}", "assets": assets, "core_bottom": core_bottom, "footer_top": footer_top})
    finally:
        presentation.Close()
        app.Quit()
    return {"pages": pages}


def extract_manifest_python_pptx(pptx_path: str | Path) -> dict:
    from pptx import Presentation

    presentation = Presentation(str(Path(pptx_path).resolve()))
    pages = []
    for page_no, slide in enumerate(presentation.slides, start=1):
        assets: dict[str, list[dict]] = {}
        core_bottom = 0.0
        footer_top = float(presentation.slide_height) / EMU_PER_POINT
        for shape in slide.shapes:
            name = str(shape.name)
            left = float(shape.left) / EMU_PER_POINT
            top = float(shape.top) / EMU_PER_POINT
            width = float(shape.width) / EMU_PER_POINT
            height = float(shape.height) / EMU_PER_POINT
            if name == "SKEL_CORE":
                core_bottom = top + height
            elif name == "SKEL_SOURCE":
                footer_top = top
            elif name.startswith("ASSET_"):
                asset_id = name[len("ASSET_"):]
                assets.setdefault(asset_id, []).append(
                    {
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                    }
                )
        pages.append(
            {
                "slide_id": f"S{page_no:02d}",
                "assets": assets,
                "core_bottom": core_bottom,
                "footer_top": footer_top,
            }
        )
    return {"pages": pages}


def audit_pptx(pptx_path: str | Path, generator_path: str | Path, project_dir: str | Path, output_path: str | Path | None = None) -> dict:
    extractor = _load_extractor()
    _, _, asset_crops = extractor.load_generator_contract(generator_path)
    project = Path(project_dir)
    report = json.loads((project / ".build" / "direct_asset_report.json").read_text(encoding="utf-8"))
    census_crop_ids: set[str] | None = None
    brief_path = project / "project_brief.json"
    brief: dict = {}
    if brief_path.is_file():
        brief = json.loads(brief_path.read_text(encoding="utf-8"))
        if brief.get("pipeline_revision") in {"5.9.1", "5.9.2", "5.9.4", "5.9.5", "5.9.6"}:
            visual_manifest = json.loads(
                (project / ".build" / "visual_manifest.json").read_text(encoding="utf-8")
            )
            census_crop_ids = {
                str(visual["asset_id"])
                for page in visual_manifest.get("pages", {}).values()
                if isinstance(page, dict)
                for visual in page.get("visuals", [])
                if isinstance(visual, dict)
                and visual.get("treatment", visual.get("disposition")) == "crop"
                and isinstance(visual.get("asset_id"), str)
                and visual.get("asset_id")
            }
    ppt_manifest = (
        extract_manifest_python_pptx(pptx_path)
        if brief.get("pipeline_revision") in {"5.9.4", "5.9.5", "5.9.6"}
        else extract_manifest(pptx_path)
    )
    result = audit_manifest(
        ppt_manifest,
        asset_crops,
        report,
        census_crop_ids=census_crop_ids,
    )
    result["pptx_sha256"] = sha256_file(pptx_path)
    if output_path:
        destination = Path(output_path)
        destination.parent.mkdir(parents=True, exist_ok=True)
        destination.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    return result


def main() -> None:
    parser = argparse.ArgumentParser(description="Audit V5.5 complete one-object bitmap placement in a PowerPoint deck.")
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--generator", required=True, type=Path)
    parser.add_argument("--project", required=True, type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()
    result = audit_pptx(args.pptx, args.generator, args.project, args.output)
    print(json.dumps(result, ensure_ascii=False, indent=2))
    raise SystemExit(0 if result["ok"] else 1)


if __name__ == "__main__":
    main()
