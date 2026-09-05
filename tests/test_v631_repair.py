from __future__ import annotations

import hashlib
import importlib.util
import json
import os
import tempfile
import unittest
from pathlib import Path

from PIL import Image

ROOT = Path(__file__).resolve().parents[1]


def load(name):
    spec = importlib.util.spec_from_file_location(name, ROOT / 'scripts' / f'{name}.py')
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


class CoordinateTests(unittest.TestCase):
    def test_contain_has_one_scale_and_keeps_labels_registered(self):
        module = load('v63_scene_graph')
        transform = module.body_contain_transform([100, 200, 1000, 500], [0.5, 2, 10, 4])
        self.assertAlmostEqual(0.008, transform['scale'])
        point = module.map_source_point([100, 200], transform)
        self.assertAlmostEqual(1.5, point[0])
        self.assertAlmostEqual(2, point[1])

    def test_real_source_roi_does_not_follow_master_ratios(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            (project / 'blueprints').mkdir()
            Image.new('RGB', (1200, 675), 'white').save(project / 'blueprints/S01.png')
            report = load('v63_visual_tiles').generate_review_tiles(
                project, source_body_rois={'S01': [10, 100, 1180, 500]})
            self.assertEqual([10, 100, 1180, 500], report['pages']['S01']['body_roi_px'])
            self.assertIn('PAGE', [t['tile_id'] for t in report['pages']['S01']['tiles']])

    def test_line_geometry_is_derived_and_zero_height_is_valid(self):
        module = load('v63_scene_graph')
        element = module.normalize_element_geometry({'type': 'line',
            'bbox_px': [0, 0, 100, 100], 'points_px': [[10, 20], [90, 20]]})
        self.assertEqual([10., 20., 90., 20.], element['bbox_px'])
        with self.assertRaises(ValueError):
            module.normalize_element_geometry({'type': 'line', 'points_px': [[10, 20], [10, 20]]})


class CensusTests(unittest.TestCase):
    def test_fractional_stroke_box_is_not_truncated_to_zero(self):
        self.assertTrue(load('v63_visual_census')._valid_box([10,20,10.5,40], [0,0,100,100]))
    def test_parent_cannot_absorb_missing_logo(self):
        errors = load('v63_visual_census').validate_candidate_hierarchy([
            {'candidate_id': 'P', 'kind': 'panel', 'child_candidate_ids': ['LOGO']}])
        self.assertIn('V63_CENSUS_CHILD_MISSING', {e['code'] for e in errors})

    def test_world_map_cannot_be_relabelled_basic_geometry(self):
        errors = load('v63_visual_census').validate_candidate_hierarchy([
            {'candidate_id': 'M', 'observed_subject': 'world_map', 'kind': 'basic_geometry',
             'expected_treatment': 'editable'}])
        self.assertIn('V63_CENSUS_SUBJECT_SUBSTITUTED', {e['code'] for e in errors})

    def test_amendment_allows_evidenced_addition_not_removal(self):
        module = load('v63_visual_census')
        original = {'pages': {'S01': {'candidates': [{'candidate_id': 'A', 'kind': 'text'}]}}}
        amended = {'pages': {'S01': {'candidates': [{'candidate_id': 'A', 'kind': 'text'},
            {'candidate_id': 'B', 'kind': 'logo', 'bbox_px': [1, 2, 3, 4]}]}}}
        change = {'slide_id': 'S01', 'candidate_id': 'B', 'action': 'add',
                  'reason': 'visible missing logo', 'evidence_px': [1, 2, 3, 4]}
        self.assertEqual([], module.validate_census_amendment(original, amended, [change]))
        self.assertTrue(module.validate_census_amendment(amended, original, []))


class CropTests(unittest.TestCase):
    def test_mask_removes_overlay_without_mutating_source(self):
        module = load('v63_extract_scene_assets')
        source = Image.new('RGBA', (20, 20), (10, 40, 120, 255))
        before = source.tobytes()
        result = module.apply_crop_recipe(source, [0, 0, 20, 20], {
            'mode': 'masked_crop', 'exclude_regions': [{'polygon_px': [[2,2],[6,2],[6,6],[2,6]],
                                                       'overlay_element_ids': ['LABEL']}]})
        self.assertEqual(0, result.getpixel((4,4))[3])
        self.assertEqual((10,40,120,255), result.getpixel((12,12)))
        self.assertEqual(before, source.tobytes())

    def test_unknown_recipe_and_unbound_mask_rejected(self):
        module = load('v63_extract_scene_assets')
        source = Image.new('RGBA', (20,20), 'blue')
        for recipe in ({'mode': 'regenerate'}, {'mode': 'masked_crop', 'exclude_regions': [
            {'polygon_px': [[0,0],[30,0],[30,30]], 'overlay_element_ids': []}]}):
            with self.assertRaises(ValueError):
                module.apply_crop_recipe(source, [0,0,20,20], recipe)


@unittest.skipUnless(os.name == 'nt', 'Windows COM only')
class WindowsRunsTests(unittest.TestCase):
    def test_real_com_preserves_mixed_chinese_and_number_runs(self):
        import pythoncom
        import win32com.client
        from pptx import Presentation
        pythoncom.CoInitialize()
        app = win32com.client.DispatchEx('PowerPoint.Application')
        presentation = None
        try:
            presentation = app.Presentations.Add(False)
            slide = presentation.Slides.Add(1, 12)
            shape = load('v63_windows_scene_renderer')._render_text(slide, {
                'bbox_in': [1, 1, 5, 1], 'style': {'font_size': 20, 'color': '#123456'},
                'runs': [{'text': '年均增长 '}, {'text': '5.1%', 'style': {'bold': True, 'color': '#C00000'}}]})
            shape.Name = 'runs_probe'
            renderer = load('v63_windows_scene_renderer')
            compact = renderer._render_text(slide, {'bbox_in': [1, 3, 0.65, 0.3], 'text': '27.2%',
                'style': {'font_size': 30, 'word_wrap': False, 'fit': 'shrink_to_box', 'margin': 0,
                          'margin_top': 0, 'margin_bottom': 0}})
            self.assertEqual(0, compact.TextFrame2.WordWrap)
            self.assertLessEqual(compact.TextFrame2.TextRange.BoundWidth, compact.Width + 0.5)
            self.assertLessEqual(compact.TextFrame2.TextRange.BoundHeight, compact.Height + 0.5)
            arrow = renderer._render_freeform(slide, {'bbox_in': [1,4,2,1], 'points_in': [[1,4],[1.5,4.8],[3,5]],
                'closed': False, 'style': {'line': '#112233', 'arrow_end': True}})
            self.assertNotEqual(1, arrow.Line.EndArrowheadStyle)
            with tempfile.TemporaryDirectory() as directory:
                output = Path(directory) / 'runs.pptx'
                presentation.SaveAs(str(output), 24)
                presentation.Close()
                presentation = None
                saved = next(s for s in Presentation(output).slides[0].shapes if s.name == 'runs_probe')
                self.assertEqual('年均增长 5.1%', saved.text)
                highlight = [r for p in saved.text_frame.paragraphs for r in p.runs if '5.1%' in r.text]
                self.assertTrue(highlight)
                self.assertEqual('C00000', str(highlight[0].font.color.rgb))
                self.assertTrue(highlight[0].font.bold)
        finally:
            if presentation is not None:
                presentation.Close()
            app.Quit()
            pythoncom.CoUninitialize()


class AcceptanceTests(unittest.TestCase):
    def test_picture_bytes_and_outline_are_audited(self):
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            path = project / 'crop.png'
            Image.new('RGB', (20,20), 'blue').save(path)
            deck = Presentation()
            slide = deck.slides.add_slide(deck.slide_layouts[6])
            picture = slide.shapes.add_picture(str(path), Inches(1), Inches(1))
            picture.line.fill.background()
            asset = {'asset_path':'crop.png', 'asset_sha256':hashlib.sha256(path.read_bytes()).hexdigest()}
            check = load('v63_acceptance')._picture_asset_errors
            self.assertEqual([], check(picture, asset, project))
            picture.line.color.rgb = RGBColor(255,0,0)
            self.assertIn('V63_PPTX_PICTURE_OUTLINE', check(picture, asset, project))
            asset['asset_sha256'] = '0'*64
            self.assertIn('V63_PPTX_ASSET_HASH_MISMATCH', check(picture, asset, project))

    def test_pending_review_finalizes_and_caches_without_a_second_build(self):
        from tests.test_v63_deconstruction import write_project
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_project(project)
            graph_path = project / '.build/v63_scene_graph.json'
            graph = json.loads(graph_path.read_text(encoding='utf-8'))
            graph['pages']['S01']['coordinate_mode'] = 'source_pixels_contain'
            graph_path.write_text(json.dumps(graph), encoding='utf-8')
            output = project / 'probe.pptx'
            output.write_bytes(b'state-only-checkpoint')
            module = load('v63_acceptance')
            bindings = module.visual_bindings(project, output)
            (project / '.build/v63_pending_review.json').write_text(json.dumps({
                'attempt_count': 1, 'audit': {'ok': True, 'blockers': []},
                'render': {'ok': True, 'status': 'pass', 'visual_verification': True},
                'bindings': bindings}), encoding='utf-8')
            (project / '.build/v6_build_attempt.json').write_text(json.dumps({
                'attempt_count': 1, 'builder_backend': 'windows_com_v584', 'status': 'awaiting_visual_review'}), encoding='utf-8')
            pipeline = load('project_pipeline')
            brief = json.loads((project / 'project_brief.json').read_text(encoding='utf-8'))
            first = pipeline._finish_v631_review(project, brief, output, auto_package=False)
            self.assertFalse(first['ok'])
            self.assertEqual('awaiting_visual_review', first['status'])
            review = {'reviewed': True, 'bindings': bindings, 'pages': [{'slide_id': 'S01',
                'object_checks': [{'candidate_id': 'C1', 'status': 'present', 'evidence': 'fixture panel'}], 'differences': []}]}
            (project / '.build/v63_visual_review.json').write_text(json.dumps(review), encoding='utf-8')
            final = pipeline._finish_v631_review(project, brief, output, auto_package=False)
            self.assertTrue(final['ok'])
            result = pipeline._run_v6_project(project, brief, output, catastrophic_repair=False, user_revision=False, auto_package=False)
            self.assertTrue(result['cached'])
            self.assertEqual(1, result['build_attempt'])
            output.write_bytes(b'changed')
            self.assertIsNone(module.locked_v63_acceptance(project, output))
            self.assertIsNone(module.read_bound_visual_review(project, output))

    def test_refinement_cannot_be_accepted(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            output = project / 'probe.pptx'
            output.write_bytes(b'state-test-not-a-real-pptx')
            module = load('v63_acceptance')
            result = module.evaluate_v63_acceptance(project, output,
                {'ok': True, 'blockers': []}, {'ok': True, 'status': 'pass', 'visual_verification': True},
                {'action': 'targeted_refinement', 'warnings': []})
            self.assertFalse(result['accepted'])
            self.assertIsNone(module.locked_v63_acceptance(project, output))

    def test_good_metric_does_not_overrule_missing_visual_subject(self):
        result = load('v63_visual_delta').evaluate_visual_delta(build_attempt=1, structural_ok=True,
            fidelity_report={'passed': True}, visual_review={'reviewed': True,
                'pages': [{'slide_id': 'S01', 'differences': [{'candidate_id': 'LOGO', 'severity': 'material', 'message': 'missing'}]}]})
        self.assertEqual('targeted_refinement', result['action'])

    def test_missing_review_waits_without_rebuild(self):
        result = load('v63_visual_delta').evaluate_visual_delta(build_attempt=1, structural_ok=True,
            fidelity_report={'passed': True}, visual_review=None, require_visual_review=True)
        self.assertEqual('awaiting_visual_review', result['action'])


class MacMappingTests(unittest.TestCase):
    def test_multiple_pages_keep_all_five_original_placeholders(self):
        from pptx import Presentation
        with tempfile.TemporaryDirectory() as directory:
            project=Path(directory)
            (project/'.build').mkdir()
            (project/'.build/v63_scene_graph.json').write_text(json.dumps({'pages':{
                sid:{'body_roi_px':[0,0,100,100],'elements':[]} for sid in ('S01','S02')}}),encoding='utf-8')
            (project/'.build/slides.json').write_text(json.dumps([{'slide_id':sid} for sid in ('S01','S02')]),encoding='utf-8')
            (project/'.build/v63_asset_ledger.json').write_text('{"assets":[]}',encoding='utf-8')
            output=project/'two.pptx'
            load('v63_mac_scene_renderer').build_deck(project,output,template_path=ROOT/'assets/company_template.pptx')
            skeleton=load('v63_skeleton_contract')
            self.assertEqual(2,len(Presentation(output).slides))
            self.assertTrue(skeleton.audit_pptx_skeleton(ROOT/'assets/company_template.pptx',output)['ok'])

    def test_explicit_text_margin_uses_points_not_inches(self):
        from pptx import Presentation
        slide = Presentation().slides.add_slide(Presentation().slide_layouts[6])
        shape = load('v63_mac_scene_renderer')._render_text(slide, {
            'bbox_in': [1,1,4,1], 'text': '测试', 'style': {'margin_left': 1.5}})
        self.assertAlmostEqual(1.5, shape.text_frame.margin_left.pt)

    def test_mac_and_windows_keep_source_aspect_ratio(self):
        page = {'body_roi_px': [100,200,1000,500], 'coordinate_mode': 'source_pixels_contain',
                'elements': [{'element_id': 'A', 'type': 'rect', 'bbox_px': [100,200,200,300],
                              'z_order': 1, 'style': {}, 'source_candidate_ids': ['A']}]}
        for name in ('v63_windows_scene_renderer', 'v63_mac_scene_renderer'):
            box = load(name).render_plan(page, [0.5,2,10,4])[0]['bbox_in']
            self.assertEqual([1.5,2.,0.8,0.8], box)


if __name__ == '__main__':
    unittest.main()
