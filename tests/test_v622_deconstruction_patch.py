from __future__ import annotations

import importlib.util
import os
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path
from types import SimpleNamespace

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def add_skeleton(slide) -> None:
    definitions = (
        ("SKEL_CHAPTER", .56, .10, 11.13, .20, "Chapter"),
        ("SKEL_TITLE", .56, .30, 11.13, .20, "Title"),
        ("SKEL_CORE", .56, .50, 11.13, .10, "■ Core"),
        ("SKEL_SOURCE", .56, 6.80, 11.13, .10, "Source"),
        ("SKEL_PAGE_NUMBER", .56, 6.95, 1.00, .10, "1"),
    )
    for name, left, top, width, height, text in definitions:
        shape = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.name = name
        shape.text = text


class ShrinkingFrame:
    def __init__(self, shape) -> None:
        object.__setattr__(self, "_shape", shape)
        object.__setattr__(self, "AutoSize", 1)
        object.__setattr__(
            self,
            "TextRange",
            SimpleNamespace(
                Text="",
                Font=SimpleNamespace(
                    Name="",
                    NameFarEast="",
                    Size=0,
                    Bold=0,
                    Fill=SimpleNamespace(ForeColor=SimpleNamespace(RGB=0)),
                ),
                ParagraphFormat=SimpleNamespace(Alignment=0),
            ),
        )

    def __setattr__(self, name, value) -> None:
        object.__setattr__(self, name, value)
        if name.startswith("Margin") and self.AutoSize != 0:
            self._shape.Height = 24.719


class V622DeconstructionPatchTests(unittest.TestCase):
    def test_skill_requires_prompt_guard_before_imagegen_and_preserves_bitmap(self):
        skill = (ROOT / "SKILL.md").read_text(encoding="utf-8")
        freedom = (ROOT / "prompts" / "v6_pre_imagegen_freedom_prompt.md").read_text(
            encoding="utf-8"
        )

        self.assertIn("# Standard Report PPT V6.2.2", skill)
        self.assertIn("scripts/v622_prompt_guard.py", skill)
        self.assertIn("V6.2.1 行为不变", skill)
        self.assertIn("must pass before the ImageGen call", freedom)
        self.assertIn("scripts/v622_prompt_guard.py", freedom)

    def test_prompt_guard_rejects_agent_added_visual_bans(self):
        guard = load_module(
            "v622_prompt_guard",
            ROOT / "scripts" / "v622_prompt_guard.py",
        )
        blocked = guard.validate_prompt(
            "企业汇报蓝图。为便于解构，无图标照片人物地图logo，全部用可编辑矩形。"
        )
        self.assertFalse(blocked["ok"])
        self.assertIn("V622_AGENT_VISUAL_BAN", {item["code"] for item in blocked["blockers"]})

        allowed = guard.validate_prompt(
            "企业汇报蓝图。用产业链节点和园区实景作为视觉锚点，突出结论与关键证据。"
        )
        self.assertTrue(allowed["ok"], allowed["blockers"])

    def test_windows_deconstruct_render_uses_same_runtime_environment_as_bitmap(self):
        pipeline = load_module(
            "v622_project_pipeline",
            ROOT / "scripts" / "project_pipeline.py",
        )
        with tempfile.TemporaryDirectory() as directory:
            dependencies = Path(directory) / "dependencies"
            python = dependencies / "python" / "python.exe"
            node = dependencies / "node" / "bin" / "node.exe"
            modules = dependencies / "node" / "node_modules"
            override = dependencies / "bin" / "override"
            python.parent.mkdir(parents=True)
            python.touch()
            node.parent.mkdir(parents=True)
            node.touch()
            modules.mkdir(parents=True)
            override.mkdir(parents=True)

            expected = pipeline._bitmap_render_environment(
                {"PATH": "existing"}, executable=python
            )
            actual = pipeline._windows_render_environment_for_mode(
                "deconstruct", {"PATH": "existing"}, executable=python
            )
            bitmap = pipeline._windows_render_environment_for_mode(
                "bitmap", {"PATH": "existing"}, executable=python
            )

            self.assertEqual(expected, actual)
            self.assertEqual(expected, bitmap)

    def test_windows_deconstruct_textbox_preserves_requested_height(self):
        template = (ROOT / "assets" / "direct_blueprint_generator_template.py").read_text(
            encoding="utf-8"
        )
        self.assertIn('DECK_META.get("construction_mode") == "deconstruct"', template)

        class Shape:
            def __init__(self, height) -> None:
                self.Height = height
                self.Fill = SimpleNamespace(
                    Visible=0,
                    Solid=lambda: None,
                    ForeColor=SimpleNamespace(RGB=0),
                )
                self.Line = SimpleNamespace(
                    Visible=0,
                    ForeColor=SimpleNamespace(RGB=0),
                    Weight=0,
                )
                self.TextFrame2 = ShrinkingFrame(self)

        class Shapes:
            def AddTextbox(self, _orientation, _x, _y, _w, height):
                return Shape(height)

        module_source = template.replace("__DECK_META__", "{}").replace(
            "__SLIDES__", "[]"
        ).replace("__DESIGN_DRAFTS__", "{}").replace(
            "__ASSET_CROPS__", "{}"
        ).replace("__PAGE_SPECS__", "{}").replace(
            "__PAGE_BUILDERS__", "{}"
        ).replace("__BUILD_SLIDE_FUNCTIONS__", "")
        namespace: dict[str, object] = {}
        exec(compile(module_source, "template", "exec"), namespace)
        namespace["clear_shape_effects"] = lambda _shape: None

        shape = namespace["add_textbox"](
            SimpleNamespace(Shapes=Shapes()),
            "Two lines of body text",
            0,
            0,
            2,
            1,
            font_size=9,
            margin_left=.04,
            margin_right=.04,
            margin_top=.02,
            margin_bottom=.02,
            preserve_requested_height=True,
        )

        self.assertAlmostEqual(72.0, shape.Height, places=3)

    def test_deconstruction_audit_rejects_collapsed_text_geometry_on_both_backends(self):
        audit = load_module(
            "v622_editability_audit",
            ROOT / "scripts" / "v6_editability_audit.py",
        )
        with tempfile.TemporaryDirectory() as directory:
            deck = Path(directory) / "collapsed.pptx"
            presentation = Presentation(str(ROOT / "assets" / "company_template.pptx"))
            slide = presentation.slides[0]
            for shape in list(slide.shapes):
                slide.shapes._spTree.remove(shape._element)
            add_skeleton(slide)
            text = slide.shapes.add_textbox(
                Inches(.56), Inches(.72), Inches(2), Inches(.3433)
            )
            text.name = "EL_BODY_TEXT_1"
            text.text = "Body copy"
            presentation.save(deck)

            specs = {
                "S01": {
                    "elements": [
                        {
                            "element_id": "BODY_TEXT",
                            "type": "text",
                            "box": [0, 0, 2, 1],
                            "text": "Body copy",
                        }
                    ]
                }
            }
            alignment = {"pages": {"S01": {"text_decisions": []}}}

            for backend in ("windows_com_v584", "mac_python_pptx_v2"):
                with self.subTest(backend=backend):
                    report = audit.audit_deconstruction_pptx(
                        deck, specs, alignment, builder_backend=backend
                    )
                    self.assertFalse(report["ok"])
                    self.assertTrue(
                        any(
                            "text geometry differs from resolved page spec"
                            in item["message"]
                            for item in report["blockers"]
                        )
                    )

    @unittest.skipUnless(os.name == "nt", "requires Windows PowerPoint COM")
    def test_real_windows_deconstruct_build_and_render_preserve_text_height(self):
        from tests.test_v6_windows_compile import COMPILER, project_files

        pipeline = load_module(
            "v622_windows_e2e_pipeline",
            ROOT / "scripts" / "project_pipeline.py",
        )
        runtime = load_module(
            "v622_windows_e2e_runtime",
            ROOT / "scripts" / "ensure_windows_runtime.py",
        )
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            project_files(
                project,
                "deconstruct",
                [
                    {
                        "element_id": "BODY_TEXT",
                        "type": "text",
                        "box": [0, 0, 2, 1],
                        "text": "Two lines of body text for geometry verification",
                    }
                ],
            )
            generator = COMPILER.compile_project(project)
            runtime.ensure_windows_runtime(project_dir=project, probe_com=True)
            output = project / "output" / "report.pptx"
            output.parent.mkdir()
            build = subprocess.run(
                [sys.executable, str(generator), "--output", str(output)],
                check=False,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=180,
            )
            self.assertEqual(0, build.returncode, build.stdout + build.stderr)

            presentation = Presentation(str(output))
            body = next(
                shape
                for shape in presentation.slides[0].shapes
                if shape.name.startswith("EL_BODY_TEXT_")
            )
            self.assertAlmostEqual(1.0, float(body.height) / 914400.0, places=2)

            render = subprocess.run(
                [
                    sys.executable,
                    str(ROOT / "scripts" / "render_slides.py"),
                    str(output),
                    "--project",
                    str(project),
                    "--expected",
                    "1",
                    "--timeout",
                    "45",
                ],
                check=False,
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=180,
                env=pipeline._windows_render_environment_for_mode("deconstruct"),
            )
            self.assertEqual(0, render.returncode, render.stdout + render.stderr)
            self.assertTrue((project / ".build" / "rendered" / "current" / "S01.png").is_file())


if __name__ == "__main__":
    unittest.main()
