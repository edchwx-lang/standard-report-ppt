from __future__ import annotations

import importlib.util
import tempfile
import unittest
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]


def load(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


OOXML = load("v623_ooxml_tests", ROOT / "scripts" / "mac_pptx_ooxml.py")


class V623DeconstructionOutputContractTests(unittest.TestCase):
    def setUp(self):
        self.subject = load(
            "v623_output_contract_tests",
            ROOT / "scripts" / "v623_deconstruction_output_contract.py",
        )

    def _deck(self, path: Path) -> Path:
        template = ROOT / "assets" / "company_template.pptx"
        presentation = Presentation(template)
        slide = presentation.slides[0]
        for shape in list(slide.shapes):
            shape._element.getparent().remove(shape._element)
        for root in [slide.element, *[m.element for m in presentation.slide_masters]]:
            OOXML.clear_forbidden_effects(root)
        for master in presentation.slide_masters:
            for layout in master.slide_layouts:
                OOXML.clear_forbidden_effects(layout.element)

        def text_box(name, text, box, size, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
            shape = slide.shapes.add_textbox(*[Inches(v) for v in box])
            OOXML.set_shape_name(shape, name)
            paragraph = shape.text_frame.paragraphs[0]
            paragraph.alignment = align
            run = paragraph.add_run()
            run.text = text
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.color.rgb = color
            return shape

        text_box("SKEL_CHAPTER", "第一章", (0.56, 0.1575, 12.20, 0.47), 20)
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.56), Inches(0.5906), Inches(12.20), Inches(0.394)
        )
        OOXML.set_shape_name(title_bar, "SKEL_TITLE_BAR")
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = RGBColor(0x1E, 0x38, 0x6B)
        title_bar.line.fill.background()
        text_box(
            "SKEL_TITLE", "页面标题", (0.60, 0.60, 12.12, 0.34), 16,
            RGBColor(255, 255, 255),
        )
        core_border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.56), Inches(1.063), Inches(12.20), Inches(0.64)
        )
        OOXML.set_shape_name(core_border, "SKEL_CORE_BORDER")
        core_border.fill.solid()
        core_border.fill.fore_color.rgb = RGBColor(255, 255, 255)
        core_border.line.color.rgb = RGBColor(0, 0, 0)
        core_border.line.width = Pt(1)
        OOXML.set_line_dash(core_border, "dash")
        text_box("SKEL_CORE", "■ 核心判断", (0.64, 1.09, 12.04, 0.59), 12)
        text_box("SKEL_SOURCE", "来源：测试", (0.56, 7.218, 10.56, 0.282), 7)
        text_box(
            "SKEL_PAGE_NUMBER", "1", (12.267, 7.218, 0.50, 0.282), 8,
            align=PP_ALIGN.RIGHT,
        )
        OOXML.clear_forbidden_effects(slide.element)
        presentation.save(path)
        return path

    def _audit(self, deck: Path):
        return self.subject.audit_deconstruction_output(
            deck, ROOT / "assets" / "company_template.pptx"
        )

    def test_reference_deck_passes_template_skeleton_and_effect_contract(self):
        with tempfile.TemporaryDirectory() as directory:
            result = self._audit(self._deck(Path(directory) / "good.pptx"))
            self.assertTrue(result["ok"], result["blockers"])

    def test_wrong_seed_layout_is_a_template_blocker(self):
        with tempfile.TemporaryDirectory() as directory:
            path = self._deck(Path(directory) / "wrong-layout.pptx")
            presentation = Presentation(path)
            presentation.slides.add_slide(presentation.slide_layouts[0])
            presentation.save(path)
            result = self._audit(path)
            self.assertIn("D623_TEMPLATE_MISMATCH", {x["code"] for x in result["blockers"]})

    def test_nonstandard_core_border_is_a_skeleton_blocker(self):
        with tempfile.TemporaryDirectory() as directory:
            path = self._deck(Path(directory) / "bad-core.pptx")
            presentation = Presentation(path)
            core = next(x for x in presentation.slides[0].shapes if x.name == "SKEL_CORE_BORDER")
            core.line.color.rgb = RGBColor(0x46, 0x3A, 0x30)
            core.line.width = Pt(0.8)
            presentation.save(path)
            result = self._audit(path)
            self.assertIn("D623_SKELETON_MISMATCH", {x["code"] for x in result["blockers"]})

    def test_inherited_effect_reference_is_a_blocker(self):
        with tempfile.TemporaryDirectory() as directory:
            path = self._deck(Path(directory) / "effect-ref.pptx")
            presentation = Presentation(path)
            shape = presentation.slides[0].shapes[0]
            style = etree.SubElement(shape._element, qn("p:style"))
            etree.SubElement(style, qn("a:effectRef"), idx="2")
            presentation.save(path)
            result = self._audit(path)
            self.assertIn("D623_FORBIDDEN_EFFECT", {x["code"] for x in result["blockers"]})

    def test_mac_scrubber_normalizes_effect_reference_and_removes_shadows(self):
        with tempfile.TemporaryDirectory() as directory:
            path = self._deck(Path(directory) / "scrub.pptx")
            presentation = Presentation(path)
            shape = presentation.slides[0].shapes[0]
            style = etree.SubElement(shape._element, qn("p:style"))
            effect_ref = etree.SubElement(style, qn("a:effectRef"), idx="3")
            etree.SubElement(effect_ref, qn("a:schemeClr"), val="accent1")
            effects = etree.SubElement(shape._element.spPr, qn("a:effectLst"))
            etree.SubElement(effects, qn("a:outerShdw"), blurRad="1000")
            OOXML.clear_forbidden_effects(shape._element)
            self.assertEqual("0", effect_ref.get("idx"))
            self.assertEqual([], shape._element.xpath(".//a:outerShdw"))


if __name__ == "__main__":
    unittest.main()
