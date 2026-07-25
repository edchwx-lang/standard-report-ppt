from __future__ import annotations

import importlib.util
import json
import copy
import hashlib
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
RUNTIME = Path(r"C:\Users\edchw\.cache\codex-runtimes\codex-primary-runtime\dependencies\python")


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def reviewed(decisions=None):
    return {"pages": {"S01": {"text_decisions": decisions or [
        {"role": "chapter", "selected": "Chapter"}, {"role": "title", "selected": "Title"},
        {"role": "core_point", "selected": "Core"}, {"role": "source", "selected": "Source"},
        {"role": "page_number", "selected": "1"},
    ], "reconstruction_contract": {"module_bindings": []}, "visuals": []}}}


def add_skeleton(slide):
    definitions = (("SKEL_CHAPTER", .56, .1, 11.13, .2, "Chapter"), ("SKEL_TITLE", .56, .3, 11.13, .2, "Title"), ("SKEL_CORE", .56, .5, 11.13, .1, "■ Core"), ("SKEL_SOURCE", .56, 6.8, 11.13, .1, "Source"), ("SKEL_PAGE_NUMBER", .56, 6.95, 1, .1, "1"))
    for name, left, top, width, height, text in definitions:
        shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        shape.name, shape.text = name, text


def add_picture(slide, asset, name, left=.56, top=.88, width=11.13, height=5.565):
    shape = slide.shapes.add_picture(str(asset), Inches(left), Inches(top), Inches(width), Inches(height))
    shape.name = name
    return shape


def add_theme_picture_outline(shape):
    style = OxmlElement("p:style")
    line_ref = OxmlElement("a:lnRef")
    line_ref.set("idx", "2")
    scheme = OxmlElement("a:schemeClr")
    scheme.set("val", "dk1")
    line_ref.append(scheme)
    style.append(line_ref)
    shape._element.append(style)


class V6EditabilityAuditTests(unittest.TestCase):
    def setUp(self):
        self.subject = load_module("v6_editability_audit", ROOT / "scripts" / "v6_editability_audit.py")
        self.bitmap = load_module("v6_bitmap", ROOT / "scripts" / "v6_bitmap.py")
        self.prebuild = load_module("v6_deconstruction", ROOT / "scripts" / "v6_deconstruction.py")

    def deck(self, directory):
        ppt = Presentation(); ppt.slide_width = Inches(12.25); ppt.slide_height = Inches(7.5)
        slide = ppt.slides.add_slide(ppt.slide_layouts[6]); add_skeleton(slide)
        return ppt, slide, Path(directory) / "deck.pptx"

    def bitmap_contract(self, directory):
        project = Path(directory) / "project"; (project / "blueprints").mkdir(parents=True, exist_ok=True)
        blueprint = Image.new("RGB", (400, 200), "navy")
        blueprint.paste("red", (0, 0, 200, 200))
        blueprint.save(project / "blueprints" / "S01.png")
        review = self.bitmap.prepare_bitmap_review(project)
        alignment = {"schema_version": "6.0", "pipeline_revision": "6.0.0", "construction_mode": "bitmap", "pages": {"S01": {"reviewed_full_page": True, "blueprint_sha256": review["pages"]["S01"]["blueprint_sha256"], "source_px": [10, 20, 390, 190], "excluded_skeleton_regions": list(self.bitmap.EXCLUDED_SKELETON_REGIONS)}}}
        (project / ".build" / "bitmap_alignment.json").write_text(json.dumps(alignment), encoding="utf-8")
        contract = self.bitmap.materialize_bitmap_assets(project)
        return project, contract, project / contract["pages"]["S01"]["asset_path"]

    def test_postbuild_uses_alignment_text_and_prefix_element_names(self):
        with tempfile.TemporaryDirectory() as directory:
            ppt, slide, deck = self.deck(directory)
            text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(.3)); text.name = "EL_TEXT_1"; text.text = "Body copy"
            table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(3), Inches(1)); table.name = "EL_TABLE_1"
            data = CategoryChartData(); data.categories = ["A"]; data.add_series("S", (1,))
            chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5), Inches(2), Inches(2), Inches(1), data); chart.name = "EL_CHART_1"
            node = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(1), Inches(.3)); node.name = "EL_FLOW_1"
            line = slide.shapes.add_connector(1, Inches(2), Inches(4.1), Inches(3), Inches(4.1)); line.name = "EL_FLOW_2"
            ppt.save(deck)
            specs = {"S01": {"elements": [{"element_id": "TEXT", "type": "text"}, {"element_id": "TABLE", "type": "matrix"}, {"element_id": "CHART", "type": "column_chart"}, {"element_id": "FLOW", "type": "flow"}]}}
            alignment = reviewed(); alignment["pages"]["S01"]["text_decisions"].append({"role": "body", "selected": "Body copy"})
            self.assertTrue(self.subject.audit_deconstruction_pptx(deck, specs, alignment)["ok"])
            slide.shapes._spTree.remove(line._element); ppt.save(deck)
            self.assertFalse(self.subject.audit_deconstruction_pptx(deck, specs, alignment)["ok"])
            ppt, slide, deck = self.deck(directory)
            text_only = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(2), Inches(.3)); text_only.name = "EL_FLOW_1"; text_only.text = "not a flow node"
            ppt.save(deck)
            self.assertFalse(self.subject.audit_deconstruction_pptx(deck, {"S01": {"elements": [{"element_id": "FLOW", "type": "flow"}]}}, reviewed())["ok"])

    def test_picture_cannot_masquerade_as_native_rect_line_or_text(self):
        with tempfile.TemporaryDirectory() as directory:
            asset = Path(directory) / "body.png"; Image.new("RGB", (400, 200), "navy").save(asset)
            expectations = {"rect": "native AutoShape", "line": "native connector", "text": "native text frame"}
            for kind, message in expectations.items():
                with self.subTest(kind=kind):
                    ppt, slide, deck = self.deck(directory); add_picture(slide, asset, f"EL_FAKE_1"); ppt.save(deck)
                    report = self.subject.audit_deconstruction_pptx(deck, {"S01": {"elements": [{"element_id": "FAKE", "type": kind}]}}, reviewed())
                    self.assertFalse(report["ok"])
                    self.assertTrue(any(message in item["message"] for item in report["blockers"]))
                    self.assertTrue(any("large body picture" in item["message"] for item in report["blockers"]))

    def test_chart_or_table_cannot_masquerade_as_flow_node(self):
        with tempfile.TemporaryDirectory() as directory:
            for impostor in ("chart", "table"):
                with self.subTest(impostor=impostor):
                    ppt, slide, deck = self.deck(directory)
                    if impostor == "chart":
                        data = CategoryChartData(); data.categories = ["A"]; data.add_series("S", (1,))
                        shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(2), Inches(1), data)
                    else:
                        shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(2), Inches(1))
                    shape.name = "EL_FLOW_1"
                    connector = slide.shapes.add_connector(1, Inches(3), Inches(2.5), Inches(4), Inches(2.5)); connector.name = "EL_FLOW_2"
                    ppt.save(deck)
                    report = self.subject.audit_deconstruction_pptx(deck, {"S01": {"elements": [{"element_id": "FLOW", "type": "flow"}]}}, reviewed())
                    self.assertFalse(report["ok"])
                    self.assertTrue(any("requires native AutoShape node and connector" in item["message"] for item in report["blockers"]))

    def test_windows_backend_explicitly_accepts_editable_chart_primitives(self):
        with tempfile.TemporaryDirectory() as directory:
            ppt, slide, deck = self.deck(directory)
            for index, left in enumerate((1, 2), start=1):
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(2), Inches(.5), Inches(1)); shape.name = f"EL_CHART_{index}"
            ppt.save(deck)
            specs = {"S01": {"elements": [{"element_id": "CHART", "type": "column_chart"}]}}
            self.assertFalse(self.subject.audit_deconstruction_pptx(deck, specs, reviewed())["ok"])
            self.assertTrue(self.subject.audit_deconstruction_pptx(deck, specs, reviewed(), builder_backend="windows_com_v584")["ok"])

    def test_large_body_picture_fails_for_editability_not_missing_text(self):
        with tempfile.TemporaryDirectory() as directory:
            asset = Path(directory) / "map.png"; Image.new("RGB", (400, 200), "navy").save(asset)
            ppt, slide, deck = self.deck(directory); add_picture(slide, asset, "EL_MAP_1"); ppt.save(deck)
            specs = {"S01": {"elements": [{"element_id": "MAP", "asset_id": "MAP_ASSET", "type": "asset"}]}}
            report = self.subject.audit_deconstruction_pptx(deck, specs, reviewed())
            self.assertFalse(report["ok"])
            self.assertTrue(any("large body picture" in item["message"] for item in report["blockers"]))
            alignment = reviewed()
            alignment["pages"]["S01"].update({"structure_modules": [{"module_id": "map", "module_kind": "pure_visual", "contains_editable_text": False}], "visuals": [{"asset_id": "MAP_ASSET", "kind": "map"}], "reconstruction_contract": {"module_bindings": [{"module_id": "map", "element_ids": ["MAP"]}]}})
            prebuild = self.prebuild.validate_deconstruction_prebuild({"schema_version": "6.0", "pipeline_revision": "6.0.0", "construction_mode": "deconstruct"}, specs, alignment, "windows_com_v584")
            self.assertTrue(self.subject.audit_deconstruction_pptx(deck, specs, alignment, prebuild["allowed_large_visual_assets_by_page"])["ok"])

    def test_large_asset_allowlist_is_scoped_to_its_slide_and_page_set_is_exact(self):
        with tempfile.TemporaryDirectory() as directory:
            asset = Path(directory) / "map.png"; Image.new("RGB", (400, 200), "navy").save(asset)
            ppt, slide, deck = self.deck(directory); add_picture(slide, asset, "EL_MAP_1")
            second = ppt.slides.add_slide(ppt.slide_layouts[6]); add_skeleton(second); add_picture(second, asset, "EL_MAP_1")
            ppt.save(deck)
            specs = {"S01": {"elements": [{"element_id": "MAP", "asset_id": "MAP_ASSET", "type": "asset"}]}, "S02": {"elements": [{"element_id": "MAP", "asset_id": "MAP_ASSET", "type": "asset"}]}}
            alignment = {"pages": {"S01": reviewed()["pages"]["S01"], "S02": reviewed()["pages"]["S01"]}}
            report = self.subject.audit_deconstruction_pptx(deck, specs, alignment, {"S01": ["MAP_ASSET"]})
            self.assertFalse(report["ok"])
            self.assertTrue(any(item["slide_id"] == "S02" and "large body picture" in item["message"] for item in report["blockers"]))

    def test_deconstruction_page_ids_must_match_specs_and_alignment(self):
        with tempfile.TemporaryDirectory() as directory:
            ppt, slide, deck = self.deck(directory)
            second = ppt.slides.add_slide(ppt.slide_layouts[6]); add_skeleton(second); ppt.save(deck)
            report = self.subject.audit_deconstruction_pptx(deck, {"S01": {"elements": []}}, reviewed(), {})
            self.assertFalse(report["ok"])
            self.assertTrue(any("page spec and alignment page ids" in item["message"] for item in report["blockers"]))

    def test_bitmap_audit_consumes_task1_contract_and_rejects_bad_contain(self):
        with tempfile.TemporaryDirectory() as directory:
            project, contract, asset = self.bitmap_contract(directory)
            for case, geometry, crop in (("good", (.56, 1.1725, 11.13, 4.98), 0), ("small", (.56, 1.1725, 10, 4.48), 0), ("off_center", (.70, 1.1725, 11.13, 4.98), 0), ("stretch", (.56, 1.1725, 11.13, 4.7), 0), ("crop", (.56, 1.1725, 11.13, 4.98), 10000)):
                with self.subTest(case=case):
                    ppt, slide, deck = self.deck(directory)
                    picture = add_picture(slide, asset, "EL_S01_BODY_BITMAP_1", *geometry)
                    picture.crop_left = crop
                    ppt.save(deck)
                    report = self.subject.audit_bitmap_pptx(deck, contract, project)
                    self.assertEqual(case == "good", report["ok"])
            ppt, slide, deck = self.deck(directory); add_picture(slide, asset, "EL_S01_BODY_BITMAP_1"); add_picture(slide, asset, "EL_OTHER_1", 1, 1, 2, 1); ppt.save(deck)
            self.assertFalse(self.subject.audit_bitmap_pptx(deck, contract, project)["ok"])

    def test_bitmap_audit_rejects_duplicate_core_bullet_and_theme_picture_outline(self):
        with tempfile.TemporaryDirectory() as directory:
            project, contract, asset = self.bitmap_contract(directory)
            ppt, slide, deck = self.deck(directory)
            core = next(shape for shape in slide.shapes if shape.name == "SKEL_CORE")
            core.text = "■ ■ Core"
            picture = add_picture(
                slide,
                asset,
                "EL_S01_BODY_BITMAP_1",
                .56,
                1.1725,
                11.13,
                4.98,
            )
            add_theme_picture_outline(picture)
            ppt.save(deck)

            report = self.subject.audit_bitmap_pptx(deck, contract, project)

            self.assertFalse(report["ok"])
            messages = [item["message"] for item in report["blockers"]]
            self.assertTrue(any("exactly one square bullet" in item for item in messages))
            self.assertTrue(any("picture outline must be none" in item for item in messages))

    def test_windows_deconstruct_audit_rejects_theme_picture_outline(self):
        with tempfile.TemporaryDirectory() as directory:
            asset = Path(directory) / "icon.png"
            Image.new("RGB", (80, 80), "navy").save(asset)
            ppt, slide, deck = self.deck(directory)
            picture = add_picture(
                slide,
                asset,
                "EL_ICON_1",
                1.0,
                2.0,
                0.6,
                0.6,
            )
            add_theme_picture_outline(picture)
            ppt.save(deck)
            specs = {
                "S01": {
                    "elements": [
                        {
                            "element_id": "ICON",
                            "asset_id": "ICON_ASSET",
                            "type": "asset",
                        }
                    ]
                }
            }

            report = self.subject.audit_deconstruction_pptx(
                deck,
                specs,
                reviewed(),
                builder_backend="windows_com_v584",
            )

            self.assertFalse(report["ok"])
            self.assertTrue(
                any(
                    "picture outline must be none" in item["message"]
                    for item in report["blockers"]
                )
            )
            picture.line.fill.background()
            ppt.save(deck)
            self.assertTrue(
                self.subject.audit_deconstruction_pptx(
                    deck,
                    specs,
                    reviewed(),
                    builder_backend="windows_com_v584",
                )["ok"]
            )

    def test_bitmap_audit_rejects_header_and_hash_chain_tampering(self):
        with tempfile.TemporaryDirectory() as directory:
            project, contract, asset = self.bitmap_contract(directory)
            ppt, slide, deck = self.deck(directory); add_picture(slide, asset, "EL_S01_BODY_BITMAP_1", .56, 1.1725, 11.13, 4.98); ppt.save(deck)
            bad_header = copy.deepcopy(contract); bad_header["schema_version"] = "5.9"
            self.assertFalse(self.subject.audit_bitmap_pptx(deck, bad_header, project)["ok"])
            blueprint = project / contract["pages"]["S01"]["source_blueprint"]
            Image.new("RGB", (400, 200), "red").save(blueprint)
            self.assertFalse(self.subject.audit_bitmap_pptx(deck, contract, project)["ok"])
            project, contract, asset = self.bitmap_contract(directory)
            ppt, slide, deck = self.deck(directory); add_picture(slide, asset, "EL_S01_BODY_BITMAP_1", .56, 1.1725, 11.13, 4.98); ppt.save(deck)
            asset.write_bytes(b"tampered crop")
            self.assertFalse(self.subject.audit_bitmap_pptx(deck, contract, project)["ok"])
            project, contract, asset = self.bitmap_contract(directory)
            embedded = Path(directory) / "embedded.png"; Image.new("RGB", (380, 170), "red").save(embedded)
            ppt, slide, deck = self.deck(directory); add_picture(slide, embedded, "EL_S01_BODY_BITMAP_1", .56, 1.1725, 11.13, 4.98); ppt.save(deck)
            self.assertFalse(self.subject.audit_bitmap_pptx(deck, contract, project)["ok"])
            escaped = copy.deepcopy(contract); escaped["pages"]["S01"]["asset_path"] = "../escaped.png"
            self.assertFalse(self.subject.audit_bitmap_pptx(deck, escaped, project)["ok"])
            extra = copy.deepcopy(contract); extra["pages"]["S02"] = copy.deepcopy(contract["pages"]["S01"])
            self.assertFalse(self.subject.audit_bitmap_pptx(deck, extra, project)["ok"])

    def test_bitmap_audit_proves_asset_is_the_declared_source_crop(self):
        with tempfile.TemporaryDirectory() as directory:
            project, contract, asset = self.bitmap_contract(directory)
            blueprint = project / contract["pages"]["S01"]["source_blueprint"]
            with Image.open(blueprint) as image:
                image.crop((0, 0, 380, 170)).save(asset)
            contract["pages"]["S01"]["asset_sha256"] = hashlib.sha256(asset.read_bytes()).hexdigest()
            ppt, slide, deck = self.deck(directory); add_picture(slide, asset, "EL_S01_BODY_BITMAP_1", .56, 1.1725, 11.13, 4.98); ppt.save(deck)
            report = self.subject.audit_bitmap_pptx(deck, contract, project)
            self.assertFalse(report["ok"])
            self.assertTrue(any("declared source crop" in item["message"] for item in report["blockers"]))
            invalid = copy.deepcopy(contract); invalid["pages"]["S01"]["source_px"] = [-1, 0, 380, 170]
            self.assertFalse(self.subject.audit_bitmap_pptx(deck, invalid, project)["ok"])


if __name__ == "__main__":
    unittest.main()
