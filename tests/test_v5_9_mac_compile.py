from __future__ import annotations

import importlib.util
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.chart import Chart


SKILL = Path(__file__).resolve().parents[1]


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def portable_font_path() -> Path:
    candidates = (
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
        Path("/Library/Fonts/Arial.ttf"),
    )
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    raise unittest.SkipTest("Arial-compatible test font is unavailable")


class V59MacCompileTests(unittest.TestCase):
    def _build_project(self, elements: list[dict]):
        compiler = load_module(
            f"v59_mac_compiler_{id(elements)}",
            SKILL / "scripts" / "project_compiler_mac.py",
        )
        temporary = tempfile.TemporaryDirectory(prefix="v59_mac_compile_")
        self.addCleanup(temporary.cleanup)
        project = Path(temporary.name)
        build = project / ".build"
        build.mkdir()
        source = project / "source.txt"
        source.write_text("source", encoding="utf-8")
        brief = {
            "schema_version": "5.9",
            "pipeline_revision": "5.9.0",
            "requested_page_count": 1,
            "production_mode": "fast",
            "platform_target": "auto",
            "source_files": [str(source)],
        }
        slides = [{
            "slide_id": "S01",
            "chapter": "第一章",
            "title": "Mac 本地生成",
            "core_points": ["python-pptx 创建可编辑对象"],
            "source": "资料来源：测试",
        }]
        (project / "project_brief.json").write_text(
            json.dumps(brief, ensure_ascii=False), encoding="utf-8"
        )
        (build / "slides.json").write_text(
            json.dumps(slides, ensure_ascii=False), encoding="utf-8"
        )
        (build / "page_specs.json").write_text(
            json.dumps({"S01": {"elements": elements}}, ensure_ascii=False),
            encoding="utf-8",
        )
        (build / "visual_manifest.json").write_text(
            json.dumps({"schema_version": "5.9", "pages": {}}), encoding="utf-8"
        )
        generator = compiler.compile_project(project)
        runtime = load_module(f"generated_mac_deck_{id(elements)}", generator)
        font_path = portable_font_path()
        output = runtime.build_deck(
            project / "output" / "report.pptx",
            font_catalog={"Microsoft YaHei": (font_path, 0)},
        )
        return project, output

    def test_compiled_generator_builds_editable_basic_elements(self):
        project, output = self._build_project([
            {"type": "text", "text": "正文", "box": [0.2, 0.2, 2.0, 0.5]},
            {"type": "rect", "box": [2.4, 0.2, 1.0, 0.5], "fill": "#EDEDED"},
        ])
        prs = Presentation(output)
        report = json.loads(
            (project / ".build" / "font_report.json").read_text(encoding="utf-8")
        )
        self.assertEqual(1, len(prs.slides))
        self.assertEqual("pillow_font_metrics", report["measurement_backend"])
        texts = [shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text")]
        self.assertTrue(any("Mac 本地生成" in text for text in texts))
        self.assertTrue(any("正文" in text for text in texts))

    def test_template_dispatches_every_allowed_element_type(self):
        template = (
            SKILL / "assets" / "python_pptx_generator_template.py"
        ).read_text(encoding="utf-8")
        required = {
            "section_header", "text", "rect", "oval", "line", "arrow",
            "text_card", "metric_strip", "hbar_chart", "column_chart",
            "line_chart", "combo_chart", "donut_chart", "grouped_hbar_chart",
            "flow", "matrix", "asset",
        }
        for element_type in required:
            self.assertIn(f'"{element_type}"', template, element_type)

    def test_mac_chart_build_is_editable(self):
        _, output = self._build_project([{
            "type": "column_chart",
            "box": [0.2, 0.2, 5.0, 2.5],
            "data": [{"label": "2025", "value": 10}, {"label": "2026", "value": 12}],
        }])
        prs = Presentation(output)
        charts = [
            shape.chart
            for shape in prs.slides[0].shapes
            if getattr(shape, "has_chart", False)
        ]
        self.assertTrue(charts)
        self.assertIsInstance(charts[0], Chart)

    def test_asset_is_inserted_from_project_build_directory(self):
        temporary = tempfile.TemporaryDirectory(prefix="v59_mac_asset_")
        self.addCleanup(temporary.cleanup)
        asset_source = Path(temporary.name) / "asset.png"
        Image.new("RGB", (80, 40), "blue").save(asset_source)
        project, _ = self._build_project([])
        asset_dir = project / ".build" / "assets" / "S01"
        asset_dir.mkdir(parents=True)
        (asset_dir / "A01.png").write_bytes(asset_source.read_bytes())
        specs = {"S01": {"elements": [
            {"type": "asset", "asset_id": "A01", "box": [0.2, 0.2, 2.0, 1.0]}
        ]}}
        (project / ".build" / "page_specs.json").write_text(
            json.dumps(specs), encoding="utf-8"
        )
        compiler = load_module(
            "v59_mac_compiler_asset",
            SKILL / "scripts" / "project_compiler_mac.py",
        )
        runtime = load_module("generated_mac_asset", compiler.compile_project(project))
        font_path = portable_font_path()
        output = runtime.build_deck(
            project / "output" / "asset.pptx",
            font_catalog={"Microsoft YaHei": (font_path, 0)},
        )
        prs = Presentation(output)
        self.assertTrue(any(shape.name == "ASSET_A01" for shape in prs.slides[0].shapes))


if __name__ == "__main__":
    unittest.main()
