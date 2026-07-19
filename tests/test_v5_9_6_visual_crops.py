from __future__ import annotations

import importlib.util
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches


SKILL = Path(__file__).resolve().parents[1]
HASH = "a" * 64
TILE_IDS = ["Q1", "Q2", "Q3", "Q4"]


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def review_record(*visual_ids: str) -> dict:
    tile_subjects = {tile_id: [] for tile_id in TILE_IDS}
    for index, visual_id in enumerate(visual_ids):
        tile_subjects[TILE_IDS[index % len(TILE_IDS)]].append(visual_id)
    return {
        "full_page_reviewed": True,
        "blueprint_sha256": HASH,
        "tile_manifest_sha256": HASH,
        "reviewed_tile_ids": TILE_IDS,
        "tile_subjects": tile_subjects,
    }


def page(*visuals: dict, **overrides) -> dict:
    payload = {
        "pipeline_revision": "5.9.6",
        "visual_reviewed": True,
        "visual_review": "reviewed_inventory",
        "visual_census_result": "reviewed_inventory",
        "page_graphics_grade": "G1",
        "design_draft_sha256": HASH,
        "observed_candidate_count": len(visuals),
        "candidate_count": len(visuals),
        "visual_review_tiles": review_record(
            *(str(item["visual_id"]) for item in visuals)
        ),
        "visuals": list(visuals),
    }
    payload.update(overrides)
    return payload


def visual(
    *,
    visual_id: str = "S01_V01",
    kind: str = "pictogram",
    treatment: str = "crop",
) -> dict:
    payload = {
        "visual_id": visual_id,
        "kind": kind,
        "description": "reviewed visual subject",
        "retention_grade": "B",
        "treatment": treatment,
        "source_px": [10, 10, 90, 90],
        "target_box_in": [1.0, 1.0, 0.8, 0.8],
        "review_tile_ids": ["Q1"],
    }
    if treatment == "crop":
        payload["asset_id"] = visual_id
    elif treatment == "native":
        payload["element_id"] = visual_id
        payload["rebuild_recipe"] = "basic_shape"
    elif treatment == "omit":
        payload["omit_reason"] = "non_evidence_decoration"
    return payload


def write_blueprint(project: Path) -> None:
    blueprints = project / "blueprints"
    blueprints.mkdir()
    image = Image.new("RGB", (400, 240), "white")
    draw = ImageDraw.Draw(image)
    for index in range(8):
        left = 12 + index * 46
        draw.rectangle((left, 80, left + 28, 108), fill=(30, 100, 220))
    path = blueprints / "S01.png"
    image.save(path)
    path.with_suffix(".composition.json").write_text(
        json.dumps({"body_roi": [0, 0, 400, 240]}),
        encoding="utf-8",
    )


def write_generator(project: Path, *, missing_last_crop: bool = False) -> Path:
    visuals = []
    crops = {}
    complex_visuals = []
    for index in range(8):
        asset_id = f"S01_V{index + 1:02d}"
        left = 10 + index * 46
        item = {
            "visual_id": asset_id,
            "asset_id": asset_id,
            "kind": "pictogram",
            "description": f"subject {index + 1}",
            "retention_grade": "B",
            "treatment": "crop",
            "source_px": [left, 78, left + 32, 112],
            "target_box_in": [0.5 + index, 1.0, 0.5, 0.5],
            "review_tile_ids": ["Q1" if index < 4 else "Q2"],
        }
        visuals.append(item)
        complex_visuals.append(
            {
                "asset_id": asset_id,
                "kind": "pictogram",
                "description": item["description"],
            }
        )
        if not missing_last_crop or index < 7:
            crops[asset_id] = {
                "slide_id": "S01",
                "kind": "pictogram",
                "source_px": item["source_px"],
                "target_box_in": item["target_box_in"],
            }
    generator = project / "generate_deck.py"
    generator.write_text(
        "\n".join(
            [
                "DECK_META = " + repr({"schema_version": "5.9"}),
                "SLIDES = "
                + repr(
                    [
                        {
                            "slide_id": "S01",
                            "visual_review": "reviewed_inventory",
                            "visual_inventory": visuals,
                            "complex_visuals": complex_visuals,
                        }
                    ]
                ),
                "BLUEPRINTS = "
                + repr({"S01": {"path": "blueprints/S01.png"}}),
                "DESIGN_DRAFTS = "
                + repr({"S01": {"path": "blueprints/S01.png"}}),
                "ASSET_CROPS = " + repr(crops),
            ]
        )
        + "\n",
        encoding="utf-8",
    )
    return generator


class V596VisualCropTests(unittest.TestCase):
    def setUp(self):
        self.contracts = load_module(
            "v596_reconstruction_contract",
            SKILL / "scripts" / "v591_reconstruction_contract.py",
        )

    def test_pipeline_accepts_v596(self):
        pipeline = load_module(
            "v596_pipeline",
            SKILL / "scripts" / "project_pipeline.py",
        )
        brief = {
            "schema_version": "5.9",
            "pipeline_revision": "5.9.6",
            "requested_page_count": 1,
            "production_mode": "blueprint",
            "blueprint_engine": "builtin_imagegen",
            "platform_target": "auto",
            "source_files": ["C:/fixture/source.docx"],
        }
        self.assertEqual([], pipeline.validate_brief(brief))

    def test_mandatory_crop_kinds_reject_native_and_omit(self):
        for treatment in ("native", "omit"):
            with self.subTest(treatment=treatment):
                subject = visual(treatment=treatment)
                report = self.contracts.validate_visual_page(
                    "S01",
                    page(subject),
                )
                self.assertIn(
                    "VISUAL_MANDATORY_CROP_REQUIRED",
                    {item["code"] for item in report["blockers"]},
                )

    def test_basic_geometry_can_remain_native(self):
        subject = visual(kind="oval", treatment="native")
        report = self.contracts.validate_visual_page("S01", page(subject))
        self.assertEqual([], report["blockers"], report)

    def test_all_review_tiles_are_required(self):
        subject = visual()
        review = review_record(subject["visual_id"])
        review["reviewed_tile_ids"] = ["Q1", "Q2", "Q3"]
        report = self.contracts.validate_visual_page(
            "S01",
            page(subject, visual_review_tiles=review),
        )
        self.assertIn(
            "VISUAL_REVIEW_TILES_INCOMPLETE",
            {item["code"] for item in report["blockers"]},
        )

    def test_every_visual_requires_location_and_tile_membership(self):
        subject = visual()
        subject.pop("source_px")
        subject.pop("review_tile_ids")
        report = self.contracts.validate_visual_page("S01", page(subject))
        codes = {item["code"] for item in report["blockers"]}
        self.assertIn("VISUAL_SOURCE_LOCATION_REQUIRED", codes)
        self.assertIn("VISUAL_TILE_MEMBERSHIP_REQUIRED", codes)

    def test_tile_subject_index_must_equal_visual_inventory(self):
        subject = visual()
        review = review_record(subject["visual_id"])
        review["tile_subjects"]["Q4"] = ["UNKNOWN"]
        report = self.contracts.validate_visual_page(
            "S01",
            page(subject, visual_review_tiles=review),
        )
        self.assertIn(
            "VISUAL_TILE_SUBJECT_MISMATCH",
            {item["code"] for item in report["blockers"]},
        )

    def test_review_tile_generator_writes_hash_bound_quadrants(self):
        module_path = SKILL / "scripts" / "v596_visual_review.py"
        self.assertTrue(module_path.is_file())
        generator = load_module("v596_visual_review", module_path)
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            blueprints = project / "blueprints"
            blueprints.mkdir()
            blueprint = blueprints / "S01.png"
            Image.new("RGB", (200, 100), "white").save(blueprint)
            manifest = generator.generate_review_tiles(project)
            page_record = manifest["pages"]["S01"]
            self.assertEqual(TILE_IDS, [item["tile_id"] for item in page_record["tiles"]])
            self.assertEqual([200, 100], page_record["pixel_size"])
            self.assertEqual(
                generator.sha256_file(blueprint),
                page_record["blueprint_sha256"],
            )
            for tile in page_record["tiles"]:
                self.assertTrue((project / tile["path"]).is_file())
            stored = json.loads(
                (project / ".build" / "visual_review_tiles.json").read_text(
                    encoding="utf-8"
                )
            )
            self.assertEqual(manifest, stored)

    def test_review_tile_manifest_must_match_locked_blueprint(self):
        generator = load_module(
            "v596_visual_review_binding",
            SKILL / "scripts" / "v596_visual_review.py",
        )
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_blueprint(project)
            manifest = generator.generate_review_tiles(project)
            alignment = {
                "pages": {
                    "S01": {
                        "design_draft_sha256": manifest["pages"]["S01"][
                            "blueprint_sha256"
                        ],
                        "visual_review_tiles": {
                            "full_page_reviewed": True,
                            "blueprint_sha256": manifest["pages"]["S01"][
                                "blueprint_sha256"
                            ],
                            "tile_manifest_sha256": generator.sha256_file(
                                project / ".build" / "visual_review_tiles.json"
                            ),
                            "reviewed_tile_ids": TILE_IDS,
                            "tile_subjects": {tile_id: [] for tile_id in TILE_IDS},
                        },
                    }
                }
            }
            self.assertEqual(
                [],
                generator.validate_review_tiles(project, alignment),
            )
            (project / "blueprints" / "S01.png").write_bytes(b"changed")
            self.assertTrue(generator.validate_review_tiles(project, alignment))

    def test_pipeline_review_tile_gate_blocks_stale_manifest(self):
        generator = load_module(
            "v596_visual_review_pipeline_fixture",
            SKILL / "scripts" / "v596_visual_review.py",
        )
        pipeline = load_module(
            "v596_pipeline_tile_gate",
            SKILL / "scripts" / "project_pipeline.py",
        )
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_blueprint(project)
            manifest = generator.generate_review_tiles(project)
            alignment = {
                "pages": {
                    "S01": {
                        "design_draft_sha256": manifest["pages"]["S01"][
                            "blueprint_sha256"
                        ],
                        "visual_review_tiles": {
                            "full_page_reviewed": True,
                            "blueprint_sha256": manifest["pages"]["S01"][
                                "blueprint_sha256"
                            ],
                            "tile_manifest_sha256": "b" * 64,
                            "reviewed_tile_ids": TILE_IDS,
                            "tile_subjects": {tile_id: [] for tile_id in TILE_IDS},
                        },
                    }
                }
            }
            report = pipeline._v596_review_tile_report(
                project,
                {
                    "schema_version": "5.9",
                    "pipeline_revision": "5.9.6",
                },
                alignment,
            )
            self.assertFalse(report["ok"])
            self.assertIn(
                "VISUAL_REVIEW_TILE_CONTRACT",
                {item["code"] for item in report["blockers"]},
            )

    def test_v596_extracts_all_eight_declared_crops(self):
        extractor = load_module(
            "v596_asset_extractor",
            SKILL / "scripts" / "extract_direct_assets.py",
        )
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_blueprint(project)
            (project / "project_brief.json").write_text(
                json.dumps(
                    {
                        "schema_version": "5.9",
                        "pipeline_revision": "5.9.6",
                        "production_mode": "blueprint",
                    }
                ),
                encoding="utf-8",
            )
            report = extractor.extract_direct_assets(
                write_generator(project),
                project,
            )
            self.assertTrue(report["ok"], report)
            self.assertEqual(8, report["mandatory_crop_count"])
            self.assertEqual(8, report["declared_assets"])
            self.assertEqual(8, report["extracted_assets"])
            self.assertEqual([], report["missing_crop_ids"])

    def test_v596_missing_declared_crop_blocks_before_builder(self):
        extractor = load_module(
            "v596_asset_extractor_missing",
            SKILL / "scripts" / "extract_direct_assets.py",
        )
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_blueprint(project)
            (project / "project_brief.json").write_text(
                json.dumps(
                    {
                        "schema_version": "5.9",
                        "pipeline_revision": "5.9.6",
                        "production_mode": "blueprint",
                    }
                ),
                encoding="utf-8",
            )
            with self.assertRaises(ValueError):
                extractor.extract_direct_assets(
                    write_generator(project, missing_last_crop=True),
                    project,
                )
            report = json.loads(
                (project / ".build" / "direct_asset_report.json").read_text(
                    encoding="utf-8"
                )
            )
            self.assertFalse(report["ok"])
            self.assertEqual(["S01_V08"], report["missing_crop_ids"])

    def test_v596_ppt_asset_audit_requires_all_eight_insertions(self):
        audit = load_module(
            "v596_ppt_asset_audit",
            SKILL / "scripts" / "ppt_asset_audit.py",
        )
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            build = project / ".build"
            build.mkdir()
            (project / "project_brief.json").write_text(
                json.dumps(
                    {
                        "schema_version": "5.9",
                        "pipeline_revision": "5.9.6",
                        "production_mode": "blueprint",
                    }
                ),
                encoding="utf-8",
            )
            generator_path = write_generator(project)
            inventory = []
            assets = []
            presentation = Presentation()
            presentation.slide_width = Inches(13.333333)
            presentation.slide_height = Inches(7.5)
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            image_path = project / "asset.png"
            Image.new("RGB", (40, 40), (30, 100, 220)).save(image_path)
            for index in range(8):
                asset_id = f"S01_V{index + 1:02d}"
                inventory.append({"asset_id": asset_id, "treatment": "crop"})
                assets.append({"asset_id": asset_id, "aspect_ratio": 1.0})
                picture = slide.shapes.add_picture(
                    str(image_path),
                    Inches(0.5 + index),
                    Inches(1.0),
                    Inches(0.5),
                    Inches(0.5),
                )
                picture.name = f"ASSET_{asset_id}"
            pptx = project / "report.pptx"
            presentation.save(pptx)
            (build / "direct_asset_report.json").write_text(
                json.dumps({"assets": assets}),
                encoding="utf-8",
            )
            (build / "visual_manifest.json").write_text(
                json.dumps({"pages": {"S01": {"visuals": inventory}}}),
                encoding="utf-8",
            )
            report = audit.audit_pptx(pptx, generator_path, project)
            self.assertTrue(report["ok"], report)
            self.assertEqual(8, report["declared_assets"])
            self.assertEqual(8, report["inserted_assets"])
            self.assertEqual(8, report["census_crop_assets"])


if __name__ == "__main__":
    unittest.main()
