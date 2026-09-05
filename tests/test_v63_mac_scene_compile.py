from __future__ import annotations

import importlib.util
import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from tests.test_v63_deconstruction import TEMPLATE, write_project
from tests.test_v63_windows_scene_renderer import graph


ROOT = Path(__file__).resolve().parents[1]


def load(name: str):
    path = ROOT / "scripts" / f"{name}.py"
    spec = importlib.util.spec_from_file_location(f"{name}_mac_v63_test", path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def load_path(path: Path):
    spec = importlib.util.spec_from_file_location(f"v63_mac_runtime_{id(path)}", path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


class V63MacSceneCompileTests(unittest.TestCase):
    def test_mac_master_core_also_receives_plain_text(self):
        values = load("v63_mac_scene_renderer")._skeleton_values(
            {"core_points": ["■ 已带符号", "未带符号"]}, 1
        )

        self.assertEqual("已带符号\n未带符号", values["core_judgment"])

    def test_mac_normalizer_consumes_the_shared_scene_without_redesign(self):
        scene = graph()
        scene["deconstruction_runtime_revision"] = "6.3.1"
        scene["color_authority"] = "blueprint_body"
        normalized, report = load("v6_mac_spec").normalize_v63_scene_graph(scene)

        self.assertEqual(scene, normalized)
        self.assertTrue(report["ok"], report["blockers"])
        self.assertTrue(report["mac_native_render_unverified"])

    def test_mac_compiler_builds_shared_atoms_and_preserves_five_placeholders(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_project(project)
            original = json.loads(
                (project / ".build" / "v63_scene_graph.json").read_text(encoding="utf-8")
            )
            page = graph()["pages"]["S01"]
            census_path = project / ".build" / "v63_visual_census.json"
            census = json.loads(census_path.read_text(encoding="utf-8"))
            target_body = census["pages"]["S01"]["body_roi_px"]
            source_body = page["body_roi_px"]

            def point(value):
                x, y = value
                return [
                    target_body[0] + (x - source_body[0]) / source_body[2] * target_body[2],
                    target_body[1] + (y - source_body[1]) / source_body[3] * target_body[3],
                ]

            def box(value):
                first = point(value[:2])
                second = point(value[2:])
                return first + second

            census["pages"]["S01"]["candidates"][0]["bbox_px"] = [
                target_body[0],
                target_body[1],
                target_body[0] + target_body[2],
                target_body[1] + target_body[3],
            ]
            census_path.write_text(json.dumps(census), encoding="utf-8")
            for element in page["elements"]:
                element["bbox_px"] = box(element["bbox_px"])
                if isinstance(element.get("points_px"), list):
                    element["points_px"] = [point(item) for item in element["points_px"]]
            original["pages"]["S01"]["body_roi_px"] = target_body
            original["pages"]["S01"]["elements"] = page["elements"]
            original["pages"]["S01"]["candidate_resolutions"]["C1"] = {
                "mode": "editable",
                "element_ids": [
                    item["element_id"]
                    for item in page["elements"]
                    if item["type"] != "group"
                ],
            }
            (project / ".build" / "v63_scene_graph.json").write_text(
                json.dumps(original), encoding="utf-8"
            )
            precheck = load("v63_deconstruction").prepare_deconstruction(
                project, backend="mac_python_pptx_v2", template_path=TEMPLATE
            )
            self.assertTrue(precheck["ok"], precheck["blockers"])
            generator = load("project_compiler_mac_v2").compile_project(project)
            runtime = load_path(generator)
            output = runtime.build_deck(project / "output" / "mac-v63.pptx")
            presentation = Presentation(output)
            names = {shape.name for shape in presentation.slides[0].shapes}
            placeholders = [
                shape for shape in presentation.slides[0].shapes if shape.is_placeholder
            ]
            compile_report = json.loads(
                (project / ".build" / "compile_report.json").read_text(encoding="utf-8")
            )

        self.assertEqual(5, len(placeholders))
        self.assertTrue(
            {"V63_RECT", "V63_ROUND", "V63_ELLIPSE", "V63_LINE", "V63_ARROW", "V63_FREE", "V63_TEXT"}.issubset(names)
        )
        self.assertTrue(compile_report["mac_native_render_unverified"])


if __name__ == "__main__":
    unittest.main()
