from __future__ import annotations

import hashlib
import importlib.util
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


ROOT = Path(__file__).resolve().parents[1]
V1_COMPILER_SHA256 = "b7960b4abf7d37d82e1ed1db9f72a8c532ad6d24e9f91ae13f4d4b7a3c479453"
V1_TEMPLATE_SHA256 = "7e2605fced46c30ee78b5ce72d5ebb0ed8c391f10675e60dd1171dbc1b3982ec"


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def sha256_file(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def portable_font_path() -> Path:
    candidates = (
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
        Path("/Library/Fonts/Arial.ttf"),
    )
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    raise unittest.SkipTest("Arial-compatible test font is unavailable")


def inches(value) -> float:
    return float(value) / 914400.0


class V6MacCompileTests(unittest.TestCase):
    def _project(self, mode: str) -> Path:
        temporary = tempfile.TemporaryDirectory(prefix=f"v6_mac_{mode}_")
        self.addCleanup(temporary.cleanup)
        project = Path(temporary.name)
        build = project / ".build"
        build.mkdir()
        (project / "blueprints").mkdir()
        (build / "design_drafts").mkdir()
        Image.new("RGB", (1600, 900), "#ddeeff").save(
            project / "blueprints" / "S01.png"
        )
        (build / "design_drafts" / "S01.png").write_bytes(
            (project / "blueprints" / "S01.png").read_bytes()
        )
        blueprint_hash = sha256_file(project / "blueprints" / "S01.png")
        source = project / "source.txt"
        source.write_text("source", encoding="utf-8")
        brief = {
            "schema_version": "6.0",
            "pipeline_revision": "6.0.0",
            "requested_page_count": 1,
            "production_mode": "blueprint",
            "construction_mode": mode,
            "blueprint_engine": "builtin_imagegen",
            "platform_target": "auto",
            "source_files": [str(source)],
        }
        slides = [
            {
                "slide_id": "S01",
                "chapter": "Chapter",
                "title": "Mac V2",
                "core_points": ["Editable reconstruction"],
                "source": "Source: test",
            }
        ]
        (project / "project_brief.json").write_text(
            json.dumps(brief), encoding="utf-8"
        )
        (build / "slides.json").write_text(
            json.dumps(slides), encoding="utf-8"
        )
        (build / "visual_manifest.json").write_text(
            json.dumps(
                {
                    "schema_version": "6.0",
                    "pipeline_revision": "6.0.0",
                    "construction_mode": mode,
                    "pages": {
                        "S01": {
                            "design_draft_path": ".build/design_drafts/S01.png",
                            "design_draft_sha256": blueprint_hash,
                            "formal_blueprint_path": "blueprints/S01.png",
                            "formal_blueprint_sha256": blueprint_hash,
                            "visuals": [],
                        }
                    },
                }
            ),
            encoding="utf-8",
        )
        for name, payload in (
            ("authoring_bundle.json", {"schema_version": "6.0"}),
            (
                "blueprint_alignment.json",
                {"schema_version": "6.0", "pages": {"S01": {}}},
            ),
            (
                "formal_blueprint_manifest.json",
                {
                    "schema_version": "6.0",
                    "pipeline_revision": "6.0.0",
                    "construction_mode": mode,
                    "pages": {
                        "S01": {
                            "design_draft_path": ".build/design_drafts/S01.png",
                            "design_draft_sha256": blueprint_hash,
                            "formal_blueprint_path": "blueprints/S01.png",
                            "formal_blueprint_sha256": blueprint_hash,
                        }
                    },
                },
            ),
        ):
            (build / name).write_text(json.dumps(payload), encoding="utf-8")
        return project

    def _compile_and_build(self, project: Path):
        compiler = load_module(
            f"v6_mac_compiler_{id(project)}",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        generator = compiler.compile_project(project)
        runtime = load_module(f"v6_mac_runtime_{id(project)}", generator)
        output = runtime.build_deck(
            project / "output" / "report.pptx",
            font_catalog={"Microsoft YaHei": (portable_font_path(), 0)},
        )
        return generator, output

    def _materialize_deconstruct(self, project: Path, elements: list[dict]) -> None:
        (project / ".build" / "page_specs.json").write_text(
            json.dumps({"S01": {"elements": elements}}),
            encoding="utf-8",
        )
        load_module(
            f"v6_mac_spec_materialize_{id(project)}",
            ROOT / "scripts" / "v6_mac_spec.py",
        ).materialize_mac_page_specs(project)

    def test_v1_files_remain_byte_identical(self):
        self.assertEqual(
            V1_COMPILER_SHA256,
            sha256_file(ROOT / "scripts" / "project_compiler_mac.py"),
        )
        self.assertEqual(
            V1_TEMPLATE_SHA256,
            sha256_file(ROOT / "assets" / "python_pptx_generator_template.py"),
        )

    def test_deconstruct_builds_native_editable_elements_with_stable_names(self):
        project = self._project("deconstruct")
        asset_dir = project / ".build" / "assets" / "S01"
        asset_dir.mkdir(parents=True)
        with Image.open(project / "blueprints" / "S01.png") as blueprint:
            blueprint.crop((0, 0, 400, 100)).save(asset_dir / "A01.png")
        elements = [
            {
                "type": "text",
                "element_id": "E_TEXT",
                "module_id": "claim",
                "text": "Editable body",
                "box": [0.1, 0.1, 2.0, 0.5],
                "fill": "#F0F1F2",
                "line": "#112233",
                "align": 2,
                "valign": 3,
                "margin_left": 0.08,
                "margin_right": 0.06,
                "margin_top": 0.03,
                "margin_bottom": 0.02,
            },
            {
                "type": "oval",
                "element_id": "E_OVAL",
                "module_id": "node",
                "text": "Editable oval",
                "box": [11.0, 0.1, 1.0, 0.7],
                "fill": "#FFFFFF",
                "line": "#7391B3",
                "color": "#1E386B",
                "align": 2,
                "valign": 3,
            },
            {
                "type": "matrix",
                "element_id": "E_TABLE",
                "module_id": "table",
                "headers": ["A", "B"],
                "rows": [["x", "y"]],
                "box": [0.1, 0.8, 2.4, 1.0],
            },
            {
                "type": "column_chart",
                "element_id": "E_CHART",
                "module_id": "chart",
                "data": [
                    {"label": "2025", "value": 10},
                    {"label": "2026", "value": 12},
                ],
                "box": [2.8, 0.1, 2.5, 1.7],
            },
            {
                "type": "combo_chart",
                "element_id": "E_COMBO",
                "module_id": "combo",
                "data": [
                    {"label": "2025", "value": 10, "line_value": 20},
                    {"label": "2026", "value": 12, "line_value": 18},
                ],
                "box": [5.5, 0.1, 2.8, 1.7],
            },
            {
                "type": "flow",
                "element_id": "E_FLOW",
                "module_id": "flow",
                "steps": [
                    {
                        "title": "First",
                        "body": "Evidence",
                        "detail": "Caveat",
                    },
                    {"label": "Second", "detail": "Decision"},
                    {"title": "Third", "body": "Action"},
                ],
                "box": [0.1, 2.1, 5.5, 0.8],
            },
            {
                "type": "asset",
                "element_id": "E_MAP",
                "module_id": "map",
                "asset_id": "A01",
                "box": [9.0, 2.0, 2.0, 2.0],
                "coord_space": "absolute",
                "fit": "contain",
            },
        ]
        (project / ".build" / "page_specs.json").write_text(
            json.dumps({"S01": {"elements": elements}}), encoding="utf-8"
        )
        manifest = json.loads(
            (project / ".build" / "visual_manifest.json").read_text(
                encoding="utf-8"
            )
        )
        manifest["pages"]["S01"]["visuals"] = [
            {
                "treatment": "crop",
                "asset_id": "A01",
                "kind": "map",
                "source_px": [0, 0, 400, 100],
                "target_box_in": [9.0, 2.0, 2.0, 2.0],
                "target_coord_space": "absolute",
            }
        ]
        (project / ".build" / "visual_manifest.json").write_text(
            json.dumps(manifest), encoding="utf-8"
        )
        spec_module = load_module(
            f"v6_mac_spec_{id(project)}", ROOT / "scripts" / "v6_mac_spec.py"
        )
        spec_module.materialize_mac_page_specs(project)
        generator, output = self._compile_and_build(project)

        prs = Presentation(output)
        slide = prs.slides[0]
        names = [shape.name for shape in slide.shapes]
        self.assertTrue(any(name.startswith("EL_E_TEXT_") for name in names))
        text_shape = next(
            shape
            for shape in slide.shapes
            if shape.name.startswith("EL_E_TEXT_")
        )
        self.assertEqual("Editable body", text_shape.text)
        self.assertAlmostEqual(0.08, inches(text_shape.text_frame.margin_left), places=2)
        self.assertEqual(3, int(text_shape.text_frame.vertical_anchor))
        oval = next(
            shape
            for shape in slide.shapes
            if shape.name.startswith("EL_E_OVAL_")
        )
        self.assertIn("Editable oval", oval.text)
        self.assertTrue(
            any(
                shape.name.startswith("EL_E_TABLE_") and shape.has_table
                for shape in slide.shapes
            )
        )
        self.assertTrue(
            any(
                shape.name.startswith("EL_E_CHART_") and shape.has_chart
                for shape in slide.shapes
            )
        )
        combo = [
            shape
            for shape in slide.shapes
            if shape.name.startswith("EL_E_COMBO_") and shape.has_chart
        ]
        self.assertEqual(2, len(combo))
        flow = [
            shape for shape in slide.shapes if shape.name.startswith("EL_E_FLOW_")
        ]
        self.assertGreaterEqual(len(flow), 5)
        self.assertEqual(MSO_SHAPE_TYPE.LINE, flow[0].shape_type)
        self.assertTrue(
            any(shape.shape_type != MSO_SHAPE_TYPE.LINE for shape in flow[1:])
        )
        flow_text = "\n".join(
            shape.text for shape in flow if getattr(shape, "has_text_frame", False)
        )
        self.assertIn("Evidence", flow_text)
        self.assertIn("Caveat", flow_text)
        self.assertIn("Decision", flow_text)
        picture = next(
            shape
            for shape in slide.shapes
            if shape.name.startswith("EL_E_MAP_")
        )
        self.assertAlmostEqual(2.0, inches(picture.width), places=2)
        self.assertAlmostEqual(0.5, inches(picture.height), places=2)
        self.assertAlmostEqual(2.75, inches(picture.top), places=2)
        self.assertEqual(0, picture.crop_left)
        self.assertEqual(0, picture.crop_right)
        generated_module = load_module(
            f"v6_generated_assert_{id(project)}", generator
        )
        self.assertTrue(
            all(
                element["type"] != "body_asset"
                for element in generated_module.PAGE_SPECS["S01"]["elements"]
            )
        )
        probe = Presentation()
        probe_slide = probe.slides.add_slide(probe.slide_layouts[6])
        with self.assertRaisesRegex(
            ValueError, "MAC_RECONSTRUCTION_UNSUPPORTED"
        ):
            generated_module.render_page_spec(
                probe_slide,
                {
                    "elements": [
                        {
                            "type": "column_chart",
                            "element_id": "BAD_RUNTIME_CHART",
                            "box": [0, 0, 2, 1],
                            "coord_space": "absolute",
                            "data": ["not-a-mapping"],
                        }
                    ]
                },
                {"x": 0, "y": 0, "w": 10, "h": 5},
                project,
                "S01",
            )
        generated_module.ASSET_REGISTRY["A01"]["asset_path"] = "../../outside.png"
        with self.assertRaisesRegex(ValueError, "MAC_ASSET_CONTRACT_MISMATCH"):
            generated_module._asset_path(project, "S01", "A01")
        compile_report = json.loads(
            (project / ".build" / "compile_report.json").read_text(
                encoding="utf-8"
            )
        )
        self.assertEqual("6.0", compile_report["schema_version"])
        self.assertEqual("deconstruct", compile_report["construction_mode"])
        self.assertEqual("mac_python_pptx_v2", compile_report["builder_backend"])

    def test_supported_styles_are_applied_and_flow_preserves_body_and_detail(self):
        project = self._project("deconstruct")
        self._materialize_deconstruct(
            project,
            [
                {
                    "type": "text_card",
                    "element_id": "CARD",
                    "box": [0.1, 0.1, 3.0, 1.2],
                    "title": "Styled title",
                    "body": "Styled body",
                    "title_fill": "#112233",
                    "body_fill": "#F1F2F3",
                    "title_color": "#FF0000",
                    "body_color": "#008000",
                },
                {
                    "type": "metric_strip",
                    "element_id": "METRIC",
                    "box": [3.3, 0.1, 2.5, 0.9],
                    "metrics": [
                        {
                            "label": "Styled label",
                            "value": "42%",
                            "value_color": "#123456",
                            "label_color": "#654321",
                        }
                    ],
                },
                {
                    "type": "flow",
                    "element_id": "FLOW",
                    "box": [0.1, 1.6, 6.0, 1.0],
                    "steps": [
                        {
                            "title": "Discover",
                            "body": "Evidence",
                            "detail": "Caveat",
                            "title_color": "#ABCDEF",
                            "body_color": "#234567",
                        },
                        {"title": "Decide", "body": "Action"},
                    ],
                },
                {
                    "type": "column_chart",
                    "element_id": "CHART",
                    "box": [6.3, 0.1, 3.0, 2.4],
                    "style": 12,
                    "data": [{"label": "A", "value": 1}],
                },
            ],
        )
        _, output = self._compile_and_build(project)
        slide = Presentation(output).slides[0]

        def text_shape(text: str):
            return next(
                shape
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text == text
            )

        self.assertEqual(
            "FF0000",
            str(
                text_shape("Styled title")
                .text_frame.paragraphs[0]
                .runs[0]
                .font.color.rgb
            ),
        )
        self.assertEqual(
            "008000",
            str(
                text_shape("Styled body")
                .text_frame.paragraphs[0]
                .runs[0]
                .font.color.rgb
            ),
        )
        self.assertEqual(
            "123456",
            str(
                text_shape("42%")
                .text_frame.paragraphs[0]
                .runs[0]
                .font.color.rgb
            ),
        )
        self.assertEqual(
            "654321",
            str(
                text_shape("Styled label")
                .text_frame.paragraphs[0]
                .runs[0]
                .font.color.rgb
            ),
        )
        detail_shape = text_shape("Evidence\nCaveat")
        self.assertEqual(
            "234567",
            str(
                detail_shape.text_frame.paragraphs[0].runs[0].font.color.rgb
            ),
        )
        self.assertEqual(
            "ABCDEF",
            str(
                text_shape("Discover")
                .text_frame.paragraphs[0]
                .runs[0]
                .font.color.rgb
            ),
        )
        chart = next(shape.chart for shape in slide.shapes if shape.has_chart)
        self.assertEqual(12, chart.chart_style)

    def test_bitmap_build_uses_exact_runtime_body_contain_and_editable_skeleton(self):
        project = self._project("bitmap")
        asset_id = "S01_BODY_BITMAP"
        asset_dir = project / ".build" / "assets" / "S01"
        asset_dir.mkdir(parents=True)
        asset = asset_dir / f"{asset_id}.png"
        Image.new("RGB", (1000, 400), "#225588").save(asset)
        page_specs = {
            "S01": {
                "elements": [
                    {
                        "type": "body_asset",
                        "element_id": asset_id,
                        "asset_id": asset_id,
                        "fit": "contain",
                        "target": "runtime_body_box",
                        "outline": "none",
                    }
                ]
            }
        }
        contract = {
            "schema_version": "6.0",
            "pipeline_revision": "6.0.0",
            "construction_mode": "bitmap",
            "pages": {
                "S01": {
                    "asset_id": asset_id,
                    "asset_path": f".build/assets/S01/{asset_id}.png",
                    "asset_sha256": sha256_file(asset),
                    "source_blueprint": "blueprints/S01.png",
                    "source_blueprint_sha256": sha256_file(
                        project / "blueprints" / "S01.png"
                    ),
                    "fit": "contain",
                    "target": "runtime_body_box",
                    "outline": "none",
                }
            },
        }
        (project / ".build" / "bitmap_page_specs.json").write_text(
            json.dumps(page_specs), encoding="utf-8"
        )
        (project / ".build" / "bitmap_contract.json").write_text(
            json.dumps(contract), encoding="utf-8"
        )
        _, output = self._compile_and_build(project)

        prs = Presentation(output)
        slide = prs.slides[0]
        body = next(
            shape
            for shape in slide.shapes
            if shape.name.startswith(f"EL_{asset_id}_")
        )
        core = next(shape for shape in slide.shapes if shape.name == "SKEL_CORE")
        source = next(shape for shape in slide.shapes if shape.name == "SKEL_SOURCE")
        body_left = inches(core.left)
        body_top = inches(core.top + core.height) + 0.12
        body_right = body_left + inches(core.width)
        body_bottom = inches(source.top) - 0.195
        box_width = body_right - body_left
        box_height = body_bottom - body_top
        expected_width = min(box_width, box_height * 2.5)
        expected_height = expected_width / 2.5
        self.assertAlmostEqual(expected_width, inches(body.width), places=2)
        self.assertAlmostEqual(expected_height, inches(body.height), places=2)
        self.assertAlmostEqual(
            body_left + (box_width - expected_width) / 2,
            inches(body.left),
            places=2,
        )
        self.assertAlmostEqual(
            body_top + (box_height - expected_height) / 2,
            inches(body.top),
            places=2,
        )
        self.assertEqual(0, body.crop_left)
        self.assertEqual(0, body.crop_bottom)
        line = body._element.spPr.find(qn("a:ln"))
        self.assertIsNotNone(line)
        self.assertIsNotNone(line.find(qn("a:noFill")))
        self.assertTrue(all(
            any(shape.name == skeleton and shape.has_text_frame for shape in slide.shapes)
            for skeleton in (
                "SKEL_CHAPTER",
                "SKEL_TITLE",
                "SKEL_CORE",
                "SKEL_SOURCE",
                "SKEL_PAGE_NUMBER",
            )
        ))
        pictures = [
            shape for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        self.assertEqual(1, len(pictures))
        asset_report = json.loads(
            (project / ".build" / "mac_asset_report.json").read_text(
                encoding="utf-8"
            )
        )
        self.assertEqual("6.0", asset_report["schema_version"])
        self.assertEqual("bitmap", asset_report["construction_mode"])
        self.assertEqual(sha256_file(asset), asset_report["assets"][0]["source_sha256"])

    def test_core_judgment_normalizes_legacy_leading_bullets(self):
        project = self._project("deconstruct")
        slides_path = project / ".build" / "slides.json"
        slides = json.loads(slides_path.read_text(encoding="utf-8"))
        slides[0]["core_points"] = ["■ ■ Editable reconstruction"]
        slides_path.write_text(json.dumps(slides), encoding="utf-8")
        self._materialize_deconstruct(project, [])

        _, output = self._compile_and_build(project)

        presentation = Presentation(output)
        core = next(
            shape
            for shape in presentation.slides[0].shapes
            if shape.name == "SKEL_CORE"
        )
        self.assertEqual("■ Editable reconstruction", core.text)

    def test_compiler_rejects_non_v6_projects(self):
        project = self._project("deconstruct")
        brief = json.loads((project / "project_brief.json").read_text(encoding="utf-8"))
        brief["schema_version"] = "5.9"
        (project / "project_brief.json").write_text(
            json.dumps(brief), encoding="utf-8"
        )
        compiler = load_module(
            "v6_mac_compiler_invalid",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        with self.assertRaisesRegex(ValueError, "6.0.0"):
            compiler.compile_project(project)

    def test_compiler_requires_locked_deconstruction_contracts(self):
        compiler = load_module(
            "v6_mac_compiler_locked_contracts",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        for missing in (
            "formal_blueprint_manifest.json",
            "visual_manifest.json",
            "blueprint_alignment.json",
            "authoring_bundle.json",
        ):
            project = self._project("deconstruct")
            specs = {
                "S01": {
                    "elements": [
                        {
                            "type": "text",
                            "element_id": "TEXT",
                            "text": "Locked",
                            "box": [0, 0, 2, 1],
                        }
                    ]
                }
            }
            (project / ".build" / "page_specs.json").write_text(
                json.dumps(specs), encoding="utf-8"
            )
            load_module(
                f"v6_mac_spec_missing_{missing}",
                ROOT / "scripts" / "v6_mac_spec.py",
            ).materialize_mac_page_specs(project)
            (project / ".build" / missing).unlink()
            with self.subTest(missing=missing):
                with self.assertRaisesRegex(
                    (ValueError, FileNotFoundError), "MAC_V6_CONTRACT_INVALID"
                ):
                    compiler.compile_project(project)

        project = self._project("deconstruct")
        (project / ".build" / "page_specs.json").write_text(
            json.dumps(
                {
                    "S01": {
                        "elements": [
                            {
                                "type": "text",
                                "element_id": "TEXT",
                                "text": "Locked",
                                "box": [0, 0, 2, 1],
                            }
                        ]
                    }
                }
            ),
            encoding="utf-8",
        )
        load_module(
            "v6_mac_spec_tamper",
            ROOT / "scripts" / "v6_mac_spec.py",
        ).materialize_mac_page_specs(project)
        Image.new("RGB", (1600, 900), "red").save(
            project / "blueprints" / "S01.png"
        )
        with self.assertRaisesRegex(ValueError, "MAC_BLUEPRINT_HASH_MISMATCH"):
            compiler.compile_project(project)

    def test_compiler_binds_formal_blueprint_to_original_design_draft(self):
        compiler = load_module(
            "v6_mac_compiler_blueprint_provenance",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        project = self._project("deconstruct")
        self._materialize_deconstruct(
            project,
            [
                {
                    "type": "text",
                    "element_id": "TEXT",
                    "text": "Locked",
                    "box": [0, 0, 2, 1],
                }
            ],
        )
        Image.new("RGB", (1600, 900), "red").save(
            project / "blueprints" / "S01.png"
        )
        replacement_hash = sha256_file(project / "blueprints" / "S01.png")
        for name in (
            "formal_blueprint_manifest.json",
            "visual_manifest.json",
        ):
            path = project / ".build" / name
            payload = json.loads(path.read_text(encoding="utf-8"))
            page = payload["pages"]["S01"]
            page["design_draft_sha256"] = replacement_hash
            page["formal_blueprint_sha256"] = replacement_hash
            path.write_text(json.dumps(payload), encoding="utf-8")
        with self.assertRaisesRegex(ValueError, "MAC_BLUEPRINT_HASH_MISMATCH"):
            compiler.compile_project(project)

    def test_compiler_validates_formal_manifest_revision_mode_and_paths(self):
        compiler = load_module(
            "v6_mac_compiler_formal_manifest_contract",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        for field, value in (
            ("pipeline_revision", "6.0.1"),
            ("construction_mode", "bitmap"),
        ):
            project = self._project("deconstruct")
            self._materialize_deconstruct(
                project,
                [
                    {
                        "type": "text",
                        "element_id": "TEXT",
                        "text": "Locked",
                        "box": [0, 0, 2, 1],
                    }
                ],
            )
            path = project / ".build" / "formal_blueprint_manifest.json"
            payload = json.loads(path.read_text(encoding="utf-8"))
            payload[field] = value
            path.write_text(json.dumps(payload), encoding="utf-8")
            with self.subTest(field=field):
                with self.assertRaisesRegex(
                    ValueError, "MAC_V6_CONTRACT_INVALID"
                ):
                    compiler.compile_project(project)

        project = self._project("deconstruct")
        self._materialize_deconstruct(
            project,
            [
                {
                    "type": "text",
                    "element_id": "TEXT",
                    "text": "Locked",
                    "box": [0, 0, 2, 1],
                }
            ],
        )
        path = project / ".build" / "formal_blueprint_manifest.json"
        payload = json.loads(path.read_text(encoding="utf-8"))
        payload["pages"]["S01"]["design_draft_path"] = "blueprints/S01.png"
        path.write_text(json.dumps(payload), encoding="utf-8")
        with self.assertRaisesRegex(ValueError, "MAC_BLUEPRINT_HASH_MISMATCH"):
            compiler.compile_project(project)

    def test_compiler_validates_visual_manifest_revision_and_mode(self):
        compiler = load_module(
            "v6_mac_compiler_visual_manifest_contract",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        for field, value in (
            ("schema_version", "5.9"),
            ("pipeline_revision", "6.0.1"),
            ("construction_mode", "bitmap"),
        ):
            project = self._project("deconstruct")
            self._materialize_deconstruct(
                project,
                [
                    {
                        "type": "text",
                        "element_id": "TEXT",
                        "text": "Locked",
                        "box": [0, 0, 2, 1],
                    }
                ],
            )
            path = project / ".build" / "visual_manifest.json"
            payload = json.loads(path.read_text(encoding="utf-8"))
            payload[field] = value
            path.write_text(json.dumps(payload), encoding="utf-8")
            with self.subTest(field=field):
                with self.assertRaisesRegex(
                    ValueError, "MAC_V6_CONTRACT_INVALID"
                ):
                    compiler.compile_project(project)

    def test_deconstruct_crop_must_equal_reviewed_locked_blueprint_pixels(self):
        compiler = load_module(
            "v6_mac_compiler_crop_pixels",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        for mutation in ("replace", "missing"):
            project = self._project("deconstruct")
            asset_dir = project / ".build" / "assets" / "S01"
            asset_dir.mkdir(parents=True)
            asset = asset_dir / "A01.png"
            with Image.open(project / "blueprints" / "S01.png") as blueprint:
                blueprint.crop((10, 20, 410, 120)).save(asset)
            self._materialize_deconstruct(
                project,
                [
                    {
                        "type": "asset",
                        "element_id": "MAP",
                        "module_id": "map",
                        "asset_id": "A01",
                        "box": [0, 0, 2, 1],
                        "fit": "contain",
                    }
                ],
            )
            visual_path = project / ".build" / "visual_manifest.json"
            visual = json.loads(visual_path.read_text(encoding="utf-8"))
            visual["pages"]["S01"]["visuals"] = [
                {
                    "treatment": "crop",
                    "asset_id": "A01",
                    "kind": "map",
                    "source_px": [10, 20, 410, 120],
                    "target_box_in": [0, 0, 2, 1],
                    "target_coord_space": "body",
                }
            ]
            visual_path.write_text(json.dumps(visual), encoding="utf-8")
            if mutation == "replace":
                Image.new("RGB", (400, 100), "red").save(asset)
            else:
                asset.unlink()
            with self.subTest(mutation=mutation):
                with self.assertRaisesRegex(
                    (ValueError, FileNotFoundError),
                    "^MAC_ASSET_CONTRACT_MISMATCH:",
                ):
                    compiler.compile_project(project)

    def test_deconstruct_rejects_crop_with_complete_perimeter_frame(self):
        compiler = load_module(
            "v6_mac_compiler_crop_frame",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        project = self._project("deconstruct")
        blueprint_path = project / "blueprints" / "S01.png"
        draft_path = project / ".build" / "design_drafts" / "S01.png"
        image = Image.new("RGB", (1600, 900), "#ddeeff")
        ImageDraw.Draw(image).rectangle(
            (100, 100, 499, 299),
            outline="#111111",
            width=3,
        )
        image.save(blueprint_path)
        image.save(draft_path)
        digest = sha256_file(blueprint_path)
        for manifest_name in (
            "visual_manifest.json",
            "formal_blueprint_manifest.json",
        ):
            manifest_path = project / ".build" / manifest_name
            manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            page = manifest["pages"]["S01"]
            page["design_draft_sha256"] = digest
            page["formal_blueprint_sha256"] = digest
            manifest_path.write_text(json.dumps(manifest), encoding="utf-8")

        asset_dir = project / ".build" / "assets" / "S01"
        asset_dir.mkdir(parents=True)
        with Image.open(blueprint_path) as blueprint:
            blueprint.crop((100, 100, 500, 300)).save(asset_dir / "A01.png")
        self._materialize_deconstruct(
            project,
            [
                {
                    "type": "asset",
                    "element_id": "FRAME",
                    "module_id": "visual",
                    "asset_id": "A01",
                    "box": [0, 0, 3, 1.5],
                    "fit": "contain",
                }
            ],
        )
        visual_path = project / ".build" / "visual_manifest.json"
        visual = json.loads(visual_path.read_text(encoding="utf-8"))
        visual["pages"]["S01"]["visuals"] = [
            {
                "treatment": "crop",
                "asset_id": "A01",
                "kind": "illustration",
                "source_px": [100, 100, 500, 300],
                "target_box_in": [0, 0, 3, 1.5],
                "target_coord_space": "body",
            }
        ]
        visual_path.write_text(json.dumps(visual), encoding="utf-8")

        with self.assertRaisesRegex(
            ValueError,
            "^MAC_DECONSTRUCTION_BODY_FRAME_INCLUDED:",
        ):
            compiler.compile_project(project)

    def test_compiler_rederives_mac_spec_and_checks_materialization_report(self):
        compiler = load_module(
            "v6_mac_compiler_spec_derivation",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        project = self._project("deconstruct")
        self._materialize_deconstruct(
            project,
            [
                {
                    "type": "text",
                    "element_id": "TEXT",
                    "text": "Original",
                    "box": [0, 0, 2, 1],
                }
            ],
        )
        mac_path = project / ".build" / "mac_page_specs.json"
        forged = json.loads(mac_path.read_text(encoding="utf-8"))
        forged["S01"]["elements"][0]["text"] = "Forged but valid"
        mac_path.write_text(json.dumps(forged), encoding="utf-8")
        report_path = project / ".build" / "mac_spec_report.json"
        report = json.loads(report_path.read_text(encoding="utf-8"))
        report["normalized_sha256"] = sha256_file(mac_path)
        report_path.write_text(json.dumps(report), encoding="utf-8")
        with self.assertRaisesRegex(ValueError, "MAC_V6_CONTRACT_INVALID"):
            compiler.compile_project(project)

    def test_bitmap_contract_rejects_unsafe_paths_and_missing_lock(self):
        compiler = load_module(
            "v6_mac_compiler_bitmap_security",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        for relative in (
            "../../outside.png",
            "C:/absolute.png",
            ".build/assets/S01/wrong.png",
        ):
            project = self._project("bitmap")
            asset_id = "S01_BODY_BITMAP"
            asset_dir = project / ".build" / "assets" / "S01"
            asset_dir.mkdir(parents=True)
            asset = asset_dir / f"{asset_id}.png"
            Image.new("RGB", (400, 200), "blue").save(asset)
            (project / ".build" / "bitmap_page_specs.json").write_text(
                json.dumps(
                    {
                        "S01": {
                            "elements": [
                                {
                                    "type": "body_asset",
                                    "element_id": asset_id,
                                    "asset_id": asset_id,
                                    "fit": "contain",
                                    "target": "runtime_body_box",
                                    "outline": "none",
                                }
                            ]
                        }
                    }
                ),
                encoding="utf-8",
            )
            contract = {
                "schema_version": "6.0",
                "pipeline_revision": "6.0.0",
                "construction_mode": "bitmap",
                "pages": {
                    "S01": {
                        "asset_id": asset_id,
                        "asset_path": relative,
                        "asset_sha256": sha256_file(asset),
                        "source_blueprint": "blueprints/S01.png",
                        "source_blueprint_sha256": sha256_file(
                            project / "blueprints" / "S01.png"
                        ),
                        "fit": "contain",
                        "target": "runtime_body_box",
                        "outline": "none",
                    }
                },
            }
            (project / ".build" / "bitmap_contract.json").write_text(
                json.dumps(contract), encoding="utf-8"
            )
            with self.subTest(relative=relative):
                with self.assertRaisesRegex(
                    ValueError, "MAC_ASSET_CONTRACT_MISMATCH"
                ):
                    compiler.compile_project(project)

        project = self._project("bitmap")
        (project / ".build" / "bitmap_page_specs.json").write_text(
            json.dumps(
                {
                    "S01": {
                        "elements": [
                            {
                                "type": "body_asset",
                                "element_id": "S01_BODY_BITMAP",
                                "asset_id": "S01_BODY_BITMAP",
                                "fit": "contain",
                                "target": "runtime_body_box",
                                "outline": "none",
                            }
                        ]
                    }
                }
            ),
            encoding="utf-8",
        )
        (project / ".build" / "formal_blueprint_manifest.json").unlink()
        with self.assertRaisesRegex(
            (ValueError, FileNotFoundError), "MAC_V6_CONTRACT_INVALID"
        ):
            compiler.compile_project(project)

        project = self._project("bitmap")
        (project / ".build" / "bitmap_page_specs.json").write_text(
            json.dumps(
                {
                    "S01": {
                        "elements": [
                            {
                                "type": "body_asset",
                                "element_id": "S01_BODY_BITMAP",
                                "asset_id": "S01_BODY_BITMAP",
                                "fit": "contain",
                                "target": "runtime_body_box",
                                "outline": "none",
                            }
                        ]
                    }
                }
            ),
            encoding="utf-8",
        )
        with self.assertRaisesRegex(
            (ValueError, FileNotFoundError), "MAC_V6_CONTRACT_INVALID"
        ):
            compiler.compile_project(project)

    def test_bitmap_contract_source_blueprint_hash_is_locked(self):
        compiler = load_module(
            "v6_mac_compiler_bitmap_source_lock",
            ROOT / "scripts" / "project_compiler_mac_v2.py",
        )
        project = self._project("bitmap")
        asset_id = "S01_BODY_BITMAP"
        asset_dir = project / ".build" / "assets" / "S01"
        asset_dir.mkdir(parents=True)
        asset = asset_dir / f"{asset_id}.png"
        Image.new("RGB", (400, 200), "blue").save(asset)
        (project / ".build" / "bitmap_page_specs.json").write_text(
            json.dumps(
                {
                    "S01": {
                        "elements": [
                            {
                                "type": "body_asset",
                                "element_id": asset_id,
                                "asset_id": asset_id,
                                "fit": "contain",
                                "target": "runtime_body_box",
                                "outline": "none",
                            }
                        ]
                    }
                }
            ),
            encoding="utf-8",
        )
        contract = {
            "schema_version": "6.0",
            "pipeline_revision": "6.0.0",
            "construction_mode": "bitmap",
            "pages": {
                "S01": {
                    "asset_id": asset_id,
                    "asset_path": f".build/assets/S01/{asset_id}.png",
                    "asset_sha256": sha256_file(asset),
                    "source_blueprint": "blueprints/S01.png",
                    "source_blueprint_sha256": "0" * 64,
                    "fit": "contain",
                    "target": "runtime_body_box",
                    "outline": "none",
                }
            },
        }
        (project / ".build" / "bitmap_contract.json").write_text(
            json.dumps(contract), encoding="utf-8"
        )
        with self.assertRaisesRegex(ValueError, "MAC_BLUEPRINT_HASH_MISMATCH"):
            compiler.compile_project(project)


if __name__ == "__main__":
    unittest.main()
