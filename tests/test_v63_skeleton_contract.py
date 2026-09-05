from __future__ import annotations

import importlib.util
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "assets" / "company_template.pptx"


def load_module():
    path = ROOT / "scripts" / "v63_skeleton_contract.py"
    spec = importlib.util.spec_from_file_location("v63_skeleton_contract_test", path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


class V63SkeletonContractTests(unittest.TestCase):
    def test_template_resolves_exactly_five_master_owned_roles(self):
        contract = load_module().read_template_contract(TEMPLATE)

        self.assertEqual(
            ["chapter", "page_title", "core_judgment", "source", "page_number"],
            list(contract["roles"]),
        )
        self.assertEqual(5, len({item["shape_id"] for item in contract["roles"].values()}))
        self.assertGreater(contract["body_roi_in"][2], 10.0)
        self.assertGreater(contract["body_roi_in"][3], 4.5)

    def test_text_update_preserves_skeleton_identity_geometry_and_style(self):
        module = load_module()
        values = {
            "chapter": "一、发展形势与背景意义",
            "page_title": "中国市场增速显著高于全球",
            "core_judgment": "■ 中国玩具市场预计年均增长9.5%。",
            "source": "资料来源：研究报告",
            "page_number": "3",
        }
        with tempfile.TemporaryDirectory() as directory:
            output = Path(directory) / "filled.pptx"
            module.update_pptx_skeleton(TEMPLATE, output, values)

            audit = module.audit_pptx_skeleton(TEMPLATE, output)
            presentation = Presentation(output)
            resolved = module.resolve_python_pptx_shapes(presentation.slides[0])

        self.assertTrue(audit["ok"], audit["errors"])
        self.assertEqual(values, {role: shape.text for role, shape in resolved.items()})

    def test_audit_rejects_recreated_or_restyled_master_object(self):
        module = load_module()
        with tempfile.TemporaryDirectory() as directory:
            candidate = Path(directory) / "changed.pptx"
            presentation = Presentation(TEMPLATE)
            shape = module.resolve_python_pptx_shapes(presentation.slides[0])["page_title"]
            shape.left += 91440
            presentation.save(candidate)

            audit = module.audit_pptx_skeleton(TEMPLATE, candidate)

        self.assertFalse(audit["ok"])
        self.assertTrue(any("page_title" in error for error in audit["errors"]))


if __name__ == "__main__":
    unittest.main()
