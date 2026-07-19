from __future__ import annotations

import importlib.util
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


SKILL = Path(__file__).resolve().parents[1]


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


class V59MacRenderTests(unittest.TestCase):
    def test_no_renderer_is_structurally_valid_unrendered(self):
        renderer = load_module(
            "v59_mac_render", SKILL / "scripts" / "mac_render_slides.py"
        )
        quality = load_module(
            "v59_mac_quality", SKILL / "scripts" / "mac_quality.py"
        )
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            pptx = root / "report.pptx"
            Presentation().save(pptx)
            render_result = renderer.render_project(
                pptx, root, expected_page_count=0, detector=lambda: None
            )
            report = quality.audit_mac_pptx(
                pptx,
                expected_page_count=0,
                render_result=render_result,
                font_fallbacks=[],
            )
        self.assertFalse(render_result["visual_verification"])
        self.assertEqual("structurally_valid_unrendered", report["status"])

    def test_font_fallback_prevents_clean_pass(self):
        quality = load_module(
            "v59_mac_quality_fallback", SKILL / "scripts" / "mac_quality.py"
        )
        with tempfile.TemporaryDirectory() as directory:
            pptx = Path(directory) / "report.pptx"
            Presentation().save(pptx)
            report = quality.audit_mac_pptx(
                pptx,
                expected_page_count=0,
                render_result={
                    "ok": True,
                    "visual_verification": True,
                    "renderer": "powerpoint_mac",
                },
                font_fallbacks=[{
                    "requested": "Microsoft YaHei",
                    "resolved": "PingFang SC",
                }],
            )
        self.assertEqual("pass_with_warnings", report["status"])

    def test_off_slide_shape_blocks_structural_pass(self):
        quality = load_module(
            "v59_mac_quality_geometry", SKILL / "scripts" / "mac_quality.py"
        )
        with tempfile.TemporaryDirectory() as directory:
            pptx = Path(directory) / "report.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(13.2),
                Inches(1),
                Inches(1),
                Inches(1),
            )
            prs.save(pptx)
            report = quality.audit_mac_pptx(
                pptx,
                expected_page_count=1,
                render_result={
                    "ok": True,
                    "visual_verification": False,
                    "renderer": None,
                },
                font_fallbacks=[],
            )
        self.assertEqual("blocked", report["status"])
        self.assertTrue(any("off slide" in error for error in report["errors"]))


if __name__ == "__main__":
    unittest.main()
