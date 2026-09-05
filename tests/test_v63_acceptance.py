from __future__ import annotations

import importlib.util
import json
import os
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from tests.test_v63_windows_scene_renderer import graph


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "assets" / "company_template.pptx"


def load(name: str):
    path = ROOT / "scripts" / f"{name}.py"
    spec = importlib.util.spec_from_file_location(f"{name}_acceptance_test", path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def build(project: Path) -> Path:
    (project / ".build").mkdir()
    (project / ".build" / "v63_scene_graph.json").write_text(json.dumps(graph()), encoding="utf-8")
    (project / ".build" / "v63_asset_ledger.json").write_text(json.dumps({"assets": []}), encoding="utf-8")
    (project / ".build" / "slides.json").write_text(json.dumps([{"slide_id": "S01", "chapter": "一、章节", "title": "页标题", "core_points": ["核心判断"], "source": "资料来源：测试"}], ensure_ascii=False), encoding="utf-8")
    output = project / "output" / "scene.pptx"
    load("v63_windows_scene_renderer").build_deck(project, output, template_path=TEMPLATE)
    return output


@unittest.skipUnless(os.name == "nt", "requires Windows PowerPoint COM")
class V63AcceptanceTests(unittest.TestCase):
    def test_audit_proves_every_scene_atom_and_master_placeholder(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            output = build(project)
            audit = load("v63_acceptance").audit_v63_pptx(output, project, template_path=TEMPLATE)

        self.assertTrue(audit["ok"], audit["blockers"])
        self.assertEqual(7, audit["editable_body_count"])
        self.assertEqual(0, audit["image_count"])

    def test_audit_blocks_missing_scene_atom_and_changed_skeleton(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            output = build(project)
            presentation = Presentation(output)
            slide = presentation.slides[0]
            for shape in list(slide.shapes):
                if shape.name == "V63_TEXT":
                    shape._element.getparent().remove(shape._element)
            placeholder = next(shape for shape in slide.shapes if shape.is_placeholder)
            placeholder.left += 91440
            presentation.save(output)
            audit = load("v63_acceptance").audit_v63_pptx(output, project, template_path=TEMPLATE)

        codes = {item["code"] for item in audit["blockers"]}
        self.assertIn("V63_PPTX_ELEMENT_MISSING", codes)
        self.assertIn("V63_PPTX_SKELETON_CHANGED", codes)

    def test_audit_blocks_incorrect_master_placeholder_text(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            output = build(project)
            presentation = Presentation(output)
            chapter = load("v63_skeleton_contract").resolve_python_pptx_shapes(
                presentation.slides[0]
            )["chapter"]
            chapter.text = "wrong chapter"
            presentation.save(output)
            audit = load("v63_acceptance").audit_v63_pptx(
                output, project, template_path=TEMPLATE
            )

        self.assertIn(
            "V63_PPTX_SKELETON_TEXT_MISMATCH",
            {item["code"] for item in audit["blockers"]},
        )

    def test_acceptance_requires_valid_render_and_binds_pptx_hash(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            output = build(project)
            module = load("v63_acceptance")
            audit = module.audit_v63_pptx(output, project, template_path=TEMPLATE)
            rejected = module.evaluate_v63_acceptance(project, output, audit, {"ok": False, "visual_verification": False})
            accepted = module.evaluate_v63_acceptance(project, output, audit, {"ok": True, "status": "pass", "visual_verification": True})

        self.assertFalse(rejected["accepted"])
        self.assertTrue(accepted["accepted"])
        self.assertEqual(64, len(accepted["pptx_sha256"]))

    def test_locked_acceptance_reuses_only_the_same_pptx_bytes(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            output = build(project)
            module = load("v63_acceptance")
            audit = module.audit_v63_pptx(output, project, template_path=TEMPLATE)
            module.evaluate_v63_acceptance(
                project,
                output,
                audit,
                {"ok": True, "status": "pass", "visual_verification": True},
            )
            locked = module.locked_v63_acceptance(project, output)
            output.write_bytes(output.read_bytes() + b"tampered")
            stale = module.locked_v63_acceptance(project, output)

        self.assertIsNotNone(locked)
        self.assertIsNone(stale)


if __name__ == "__main__":
    unittest.main()
