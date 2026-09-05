from __future__ import annotations

import importlib.util
import json
import os
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]


def load_module():
    path = ROOT / "scripts" / "v63_windows_scene_renderer.py"
    spec = importlib.util.spec_from_file_location("v63_windows_renderer_test", path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def graph() -> dict:
    common = {"source_candidate_ids": ["C1"]}
    return {
        "schema_version": "6.3",
        "pages": {
            "S01": {
                "body_roi_px": [100, 150, 1000, 400],
                "elements": [
                    {"element_id": "G1", "type": "group", "bbox_px": [100, 150, 1100, 550], "z_order": 0, "style": {}, "source_candidate_ids": []},
                    {"element_id": "RECT", "type": "rect", "bbox_px": [100, 150, 400, 250], "z_order": 1, "style": {"fill": "#244A82", "line": "#8CAAD0"}, **common},
                    {"element_id": "ROUND", "type": "round_rect", "bbox_px": [420, 150, 620, 250], "z_order": 2, "style": {"fill": "#EEF3F9", "line": "none"}, **common},
                    {"element_id": "ELLIPSE", "type": "ellipse", "bbox_px": [650, 150, 750, 250], "z_order": 3, "style": {"fill": "#FFFFFF", "line": "#244A82"}, **common},
                    {"element_id": "LINE", "type": "line", "bbox_px": [150, 300, 450, 300.1], "points_px": [[150, 300], [450, 300]], "z_order": 4, "style": {"line": "#2F66A8", "line_width": 1.2}, **common},
                    {"element_id": "ARROW", "type": "arrow", "bbox_px": [150, 340, 450, 340.1], "points_px": [[150, 340], [450, 340]], "z_order": 5, "style": {"line": "#2F66A8", "line_width": 2.0}, **common},
                    {"element_id": "FREE", "type": "freeform", "bbox_px": [500, 300, 700, 420], "points_px": [[500, 420], [600, 300], [700, 420], [500, 420]], "z_order": 6, "style": {"fill": "#6E96C6", "line": "none"}, **common},
                    {"element_id": "TEXT", "type": "text", "bbox_px": [150, 430, 650, 500], "z_order": 7, "text": "Editable text", "style": {"font_size": 12, "font_name": "Microsoft YaHei", "color": "#111111", "bold": True, "align": "left", "valign": "middle"}, **common},
                ],
            }
        },
    }


class V63WindowsSceneRendererTests(unittest.TestCase):
    def test_master_core_receives_plain_text_without_a_second_bullet(self):
        values = load_module()._skeleton_values(
            {"core_points": ["■ 已带符号", "未带符号"]}, 1
        )

        self.assertEqual("已带符号\r未带符号", values["core_judgment"])

    def test_render_plan_keeps_atomic_z_order_and_coordinate_mapping(self):
        plan = load_module().render_plan(
            graph()["pages"]["S01"], [0.564, 1.782, 12.205, 5.356]
        )

        self.assertEqual(["RECT", "ROUND", "ELLIPSE", "LINE", "ARROW", "FREE", "TEXT"], [item["element_id"] for item in plan])
        self.assertEqual("rect", plan[0]["type"])
        self.assertAlmostEqual(0.564, plan[0]["bbox_in"][0], places=3)

    @unittest.skipUnless(os.name == "nt", "requires Windows PowerPoint COM")
    def test_real_powerpoint_build_preserves_placeholders_and_adds_editable_atoms(self):
        module = load_module()
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            (project / ".build").mkdir()
            (project / ".build" / "v63_scene_graph.json").write_text(json.dumps(graph()), encoding="utf-8")
            (project / ".build" / "v63_asset_ledger.json").write_text(json.dumps({"assets": []}), encoding="utf-8")
            slides = [{"slide_id": "S01", "chapter": "一、章节", "title": "页标题", "core_points": ["核心判断"], "source": "资料来源：测试"}]
            (project / ".build" / "slides.json").write_text(json.dumps(slides, ensure_ascii=False), encoding="utf-8")
            output = project / "output" / "scene.pptx"

            report = module.build_deck(
                project,
                output,
                template_path=ROOT / "assets" / "company_template.pptx",
            )
            render_report = module.render_deck(
                output, project, expected_page_count=1
            )
            presentation = Presentation(output)
            slide = presentation.slides[0]
            names = {shape.name for shape in slide.shapes}
            placeholders = [shape for shape in slide.shapes if shape.is_placeholder]
            render_exists = (
                project / ".build" / "rendered" / "current" / "S01.png"
            ).is_file()

        self.assertEqual(5, len(placeholders))
        self.assertTrue({"V63_RECT", "V63_ROUND", "V63_ELLIPSE", "V63_LINE", "V63_ARROW", "V63_FREE", "V63_TEXT"}.issubset(names))
        self.assertEqual(7, report["editable_body_count"])
        self.assertEqual(0, report["image_count"])
        self.assertTrue(render_report["ok"])
        self.assertTrue(render_exists)


if __name__ == "__main__":
    unittest.main()
