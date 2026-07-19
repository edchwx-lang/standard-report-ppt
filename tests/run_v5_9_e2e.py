from __future__ import annotations

import argparse
import importlib.util
import json
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation


SKILL = Path(__file__).resolve().parents[1]


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def portable_font_path() -> Path:
    for candidate in (
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
        Path("/Library/Fonts/Arial.ttf"),
    ):
        if candidate.is_file():
            return candidate
    raise RuntimeError("Arial-compatible E2E font is unavailable")


def main(page_count: int, mode: str) -> None:
    with tempfile.TemporaryDirectory(prefix="v59_e2e_") as directory:
        project = Path(directory) / "含空格 Mac 项目"
        project.mkdir()
        build = project / ".build"
        build.mkdir()
        source = project / "source.txt"
        source.write_text("跨平台蓝图先于构建器", encoding="utf-8")
        brief = {
            "schema_version": "5.9",
            "pipeline_revision": "5.9.4",
            "requested_page_count": page_count,
            "production_mode": mode,
            "platform_target": "auto",
            "source_files": [str(source)],
        }
        if mode == "blueprint":
            brief["blueprint_engine"] = "builtin_imagegen"
        (project / "project_brief.json").write_text(
            json.dumps(brief, ensure_ascii=False), encoding="utf-8"
        )
        slides = [{
            "slide_id": f"S{index:02d}",
            "chapter": "第一章",
            "title": f"V5.9 Mac E2E {index}",
            "core_points": ["本地生成可编辑 PPTX"],
            "source": "资料来源：测试",
        } for index in range(1, page_count + 1)]
        specs = {
            f"S{index:02d}": {"elements": [
                {
                    "type": "column_chart",
                    "box": [0.2, 0.2, 5.0, 2.5],
                    "data": [
                        {"label": "Windows", "value": 10},
                        {"label": "macOS", "value": 10},
                    ],
                },
                {
                    "type": "text_card",
                    "box": [5.4, 0.2, 4.0, 2.5],
                    "title": "结论",
                    "body": "共享规格，双后端",
                },
            ]}
            for index in range(1, page_count + 1)
        }
        (build / "slides.json").write_text(
            json.dumps(slides, ensure_ascii=False), encoding="utf-8"
        )
        (build / "page_specs.json").write_text(
            json.dumps(specs, ensure_ascii=False), encoding="utf-8"
        )
        (build / "visual_manifest.json").write_text(
            json.dumps({"schema_version": "5.9", "pages": {}}),
            encoding="utf-8",
        )
        if mode == "blueprint":
            gate = load_module(
                "v59_e2e_gate", SKILL / "scripts" / "v59_blueprint_gate.py"
            )
            alignment_pages = {}
            for index in range(1, page_count + 1):
                slide_id = f"S{index:02d}"
                returned = project / f"{slide_id}-imagegen-result.png"
                Image.new("RGB", (1600, 900), "white").save(returned)
                record = gate.record_artifact(
                    project, slide_id, returned, transport_attempt_count=1
                )
                alignment_pages[slide_id] = {
                    "reviewed": True,
                    "design_draft_sha256": record["artifact_sha256"],
                    "resolved_page_spec": specs[slide_id],
                }
            (build / "blueprint_alignment.json").write_text(
                json.dumps({
                    "schema_version": "5.9",
                    "skill_version": "5.9.0",
                    "pages": alignment_pages,
                }, ensure_ascii=False),
                encoding="utf-8",
            )
        compiler = load_module(
            "v59_e2e_compiler", SKILL / "scripts" / "project_compiler_mac.py"
        )
        generator = compiler.compile_project(project)
        runtime = load_module("v59_e2e_generator", generator)
        output = runtime.build_deck(
            project / "output" / "report.pptx",
            font_catalog={"Microsoft YaHei": (portable_font_path(), 0)},
        )
        presentation = Presentation(output)
        assert len(presentation.slides) == page_count
        assert all(
            any(getattr(shape, "has_chart", False) for shape in slide.shapes)
            for slide in presentation.slides
        )
        print(json.dumps({
            "ok": True,
            "pptx": str(output),
            "slides": page_count,
            "mode": mode,
        }, ensure_ascii=False))


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--pages", type=int, choices=(1, 3, 6), default=1)
    parser.add_argument("--mode", choices=("fast", "blueprint"), default="fast")
    args = parser.parse_args()
    main(args.pages, args.mode)
