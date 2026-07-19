from __future__ import annotations

import importlib.util
import json
import tempfile
import unittest
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches


SKILL = Path(__file__).resolve().parents[1]


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


def write_generator(
    project: Path,
    *,
    crops: dict,
    visual_inventory: list[dict],
) -> Path:
    generator = project / "generate_deck.py"
    generator.write_text(
        "\n".join(
            [
                "DECK_META = " + repr({"schema_version": "5.9"}),
                "SLIDES = " + repr(
                    [
                        {
                            "slide_id": "S01",
                            "visual_review": "reviewed_inventory",
                            "visual_inventory": visual_inventory,
                            "complex_visuals": [
                                {
                                    "asset_id": item["asset_id"],
                                    "kind": item.get("kind", "pictogram"),
                                    "description": item.get("description", ""),
                                }
                                for item in visual_inventory
                                if item.get("treatment") == "crop"
                            ],
                        }
                    ]
                ),
                "DESIGN_DRAFTS = "
                + repr({"S01": {"path": "blueprints/S01.png"}}),
                "ASSET_CROPS = " + repr(crops),
                "",
            ]
        ),
        encoding="utf-8",
    )
    return generator


def write_brief(project: Path, revision: str) -> None:
    (project / "project_brief.json").write_text(
        json.dumps(
            {
                "schema_version": "5.9",
                "pipeline_revision": revision,
                "production_mode": "blueprint",
            }
        ),
        encoding="utf-8",
    )


def crop_visual(asset_id: str = "S01_V01") -> dict:
    return {
        "visual_id": asset_id,
        "asset_id": asset_id,
        "kind": "pictogram",
        "description": "blue square",
        "treatment": "crop",
        "source_px": [30, 30, 170, 170],
        "target_box_in": [2.0, 2.0, 2.0, 2.0],
    }


def crop_spec(asset_id: str = "S01_V01") -> dict:
    return {
        asset_id: {
            "slide_id": "S01",
            "source_px": [30, 30, 170, 170],
            "target_box_in": [2.0, 2.0, 2.0, 2.0],
        }
    }


def write_blueprint(project: Path) -> None:
    blueprints = project / "blueprints"
    blueprints.mkdir()
    image = Image.new("RGB", (200, 200), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((55, 55, 145, 145), fill=(30, 100, 220))
    path = blueprints / "S01.png"
    image.save(path)
    path.with_suffix(".composition.json").write_text(
        json.dumps({"body_roi": [0, 0, 200, 200]}),
        encoding="utf-8",
    )


def write_asset_pptx(
    project: Path,
    *,
    asset_count: int = 1,
    outside_target: bool = False,
) -> Path:
    image_path = project / "asset.png"
    Image.new("RGB", (80, 80), (30, 100, 220)).save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    core = slide.shapes.add_shape(
        1,
        Inches(0.5),
        Inches(0.5),
        Inches(1.0),
        Inches(0.5),
    )
    core.name = "SKEL_CORE"
    source = slide.shapes.add_shape(
        1,
        Inches(0.5),
        Inches(6.8),
        Inches(1.0),
        Inches(0.2),
    )
    source.name = "SKEL_SOURCE"
    left = Inches(5.0 if outside_target else 2.5)
    for _ in range(asset_count):
        picture = slide.shapes.add_picture(
            str(image_path),
            left,
            Inches(2.5),
            Inches(1.0),
            Inches(1.0),
        )
        picture.name = "ASSET_S01_V01"
    path = project / "report.pptx"
    presentation.save(path)
    return path


class V594PostBlueprintCropTests(unittest.TestCase):
    def setUp(self):
        self.extractor = load_module(
            "v594_crop_extractor",
            SKILL / "scripts" / "extract_direct_assets.py",
        )
        self.contracts = load_module(
            "v594_reconstruction_contract",
            SKILL / "scripts" / "v591_reconstruction_contract.py",
        )
        self.revision_contracts = load_module(
            "v594_revision_contracts",
            SKILL / "scripts" / "v591_contracts.py",
        )
        self.audit = load_module(
            "v594_asset_audit",
            SKILL / "scripts" / "ppt_asset_audit.py",
        )

    def test_v594_extracts_every_explicit_reviewed_crop(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_brief(project, "5.9.4")
            write_blueprint(project)
            generator = write_generator(
                project,
                crops=crop_spec(),
                visual_inventory=[crop_visual()],
            )
            report = self.extractor.extract_direct_assets(generator, project)
            self.assertTrue(report["ok"])
            self.assertEqual(["S01_V01"], report["requested_crop_ids"])
            self.assertEqual(["S01_V01"], report["extracted_crop_ids"])
            self.assertEqual([], report["missing_crop_ids"])
            asset = report["assets"][0]
            self.assertGreater(asset["pixel_size"][0], 0)
            self.assertTrue((project / asset["path"]).is_file())

    def test_v594_missing_blueprint_blocks_after_writing_report(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_brief(project, "5.9.4")
            generator = write_generator(
                project,
                crops=crop_spec(),
                visual_inventory=[crop_visual()],
            )
            with self.assertRaises(ValueError):
                self.extractor.extract_direct_assets(generator, project)
            report = json.loads(
                (project / ".build" / "direct_asset_report.json").read_text(
                    encoding="utf-8"
                )
            )
            self.assertFalse(report["ok"])
            self.assertEqual(["S01_V01"], report["missing_crop_ids"])

    def test_v592_missing_blueprint_remains_advisory(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_brief(project, "5.9.2")
            generator = write_generator(
                project,
                crops=crop_spec(),
                visual_inventory=[crop_visual()],
            )
            report = self.extractor.extract_direct_assets(generator, project)
            self.assertTrue(report["ok"])
            self.assertTrue(report["warnings"])

    def test_v594_zero_crop_page_is_valid(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            write_brief(project, "5.9.4")
            generator = write_generator(
                project,
                crops={},
                visual_inventory=[
                    {
                        "visual_id": "S01_ARROW",
                        "kind": "arrow",
                        "description": "editable arrow",
                        "treatment": "native",
                        "element_id": "S01_ARROW",
                        "rebuild_recipe": "line_arrow",
                    },
                    {
                        "visual_id": "S01_DECORATION",
                        "kind": "decoration",
                        "description": "optional decoration",
                        "treatment": "omit",
                        "omit_reason": "non_evidence_decoration",
                    },
                ],
            )
            report = self.extractor.extract_direct_assets(generator, project)
            self.assertTrue(report["ok"])
            self.assertEqual([], report["requested_crop_ids"])
            self.assertEqual([], report["extracted_crop_ids"])

    def test_duplicate_crop_asset_id_is_rejected_by_review_contract(self):
        visual = crop_visual()
        duplicate = {
            **crop_visual(),
            "visual_id": "S01_V02",
            "description": "second subject",
        }
        report = self.contracts.validate_visual_page(
            "S01",
            {
                "visual_reviewed": True,
                "observed_candidate_count": 2,
                "candidate_count": 2,
                "visuals": [visual, duplicate],
            },
        )
        self.assertIn(
            "VISUAL_ASSET_ID_INVALID",
            {item["code"] for item in report["blockers"]},
        )

    def test_v594_python_pptx_reader_finds_named_crop_without_com(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            pptx = write_asset_pptx(project)
            manifest = self.audit.extract_manifest_python_pptx(pptx)
            self.assertEqual(
                1,
                len(manifest["pages"][0]["assets"]["S01_V01"]),
            )

    def test_v594_asset_audit_binds_crop_to_current_pptx(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            build = project / ".build"
            build.mkdir()
            write_brief(project, "5.9.4")
            pptx = write_asset_pptx(project)
            generator = write_generator(
                project,
                crops={
                    "S01_V01": {
                        "slide_id": "S01",
                        "source_px": [30, 30, 170, 170],
                        "target_box_in": [2.0, 2.0, 2.0, 2.0],
                    }
                },
                visual_inventory=[crop_visual()],
            )
            (build / "direct_asset_report.json").write_text(
                json.dumps(
                    {
                        "assets": [
                            {
                                "asset_id": "S01_V01",
                                "aspect_ratio": 1.0,
                            }
                        ]
                    }
                ),
                encoding="utf-8",
            )
            (build / "visual_manifest.json").write_text(
                json.dumps(
                    {
                        "pages": {
                            "S01": {
                                "visuals": [
                                    {
                                        "asset_id": "S01_V01",
                                        "treatment": "crop",
                                    }
                                ]
                            }
                        }
                    }
                ),
                encoding="utf-8",
            )
            report = self.audit.audit_pptx(
                pptx,
                generator,
                project,
            )
            self.assertTrue(report["ok"], report)
            self.assertEqual(
                self.extractor.sha256_file(pptx),
                report["pptx_sha256"],
            )

    def test_v594_duplicate_inserted_crop_is_blocked(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            pptx = write_asset_pptx(project, asset_count=2)
            manifest = self.audit.extract_manifest_python_pptx(pptx)
            report = self.audit.audit_manifest(
                manifest,
                {
                    "S01_V01": {
                        "slide_id": "S01",
                        "target_box_in": [2.0, 2.0, 2.0, 2.0],
                    }
                },
                {
                    "assets": [
                        {
                            "asset_id": "S01_V01",
                            "aspect_ratio": 1.0,
                        }
                    ]
                },
                census_crop_ids={"S01_V01"},
            )
            self.assertFalse(report["ok"])
            self.assertIn("found 2", " ".join(report["errors"]))

    def test_v594_crop_outside_target_is_blocked(self):
        with tempfile.TemporaryDirectory() as directory:
            project = Path(directory)
            pptx = write_asset_pptx(project, outside_target=True)
            manifest = self.audit.extract_manifest_python_pptx(pptx)
            report = self.audit.audit_manifest(
                manifest,
                {
                    "S01_V01": {
                        "slide_id": "S01",
                        "target_box_in": [2.0, 2.0, 2.0, 2.0],
                    }
                },
                {
                    "assets": [
                        {
                            "asset_id": "S01_V01",
                            "aspect_ratio": 1.0,
                        }
                    ]
                },
                census_crop_ids={"S01_V01"},
            )
            self.assertFalse(report["ok"])
            self.assertIn("outside its target box", " ".join(report["errors"]))

    def test_only_v594_makes_asset_audit_blocking(self):
        base = {
            "schema_version": "5.9",
            "production_mode": "blueprint",
        }
        self.assertEqual(
            "blocker",
            self.revision_contracts.audit_policy(
                {**base, "pipeline_revision": "5.9.4"},
                "ppt_asset_audit",
            ),
        )
        self.assertEqual(
            "warning",
            self.revision_contracts.audit_policy(
                {**base, "pipeline_revision": "5.9.2"},
                "ppt_asset_audit",
            ),
        )


if __name__ == "__main__":
    unittest.main()
