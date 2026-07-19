from __future__ import annotations

import importlib.util
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


SKILL = Path(__file__).resolve().parents[1]


def load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


class V59MacOoxmlTests(unittest.TestCase):
    def test_adapter_sets_east_asian_font_and_dash(self):
        adapter = load_module("v59_ooxml", SKILL / "scripts" / "mac_pptx_ooxml.py")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "中文"
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(2), Inches(1))
        adapter.set_east_asian_font(run, "PingFang SC")
        adapter.set_line_dash(line, "dash")
        self.assertIn('typeface="PingFang SC"', slide.element.xml)
        self.assertIn('val="dash"', slide.element.xml)

    def test_package_validation_and_slide_removal(self):
        adapter = load_module("v59_ooxml_validate", SKILL / "scripts" / "mac_pptx_ooxml.py")
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.slides.add_slide(prs.slide_layouts[6])
        adapter.remove_last_slide(prs)
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "sample.pptx"
            prs.save(path)
            result = adapter.validate_pptx_package(path)
        self.assertEqual(1, result["slides"])
