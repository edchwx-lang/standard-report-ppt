"""Regression tests for the actual one-page release failures (protocol fixtures)."""
import hashlib
import importlib.util
import io
import json
import os
import subprocess
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch
from zipfile import ZipFile

from PIL import Image
from tests import test_v6_delivery as legacy_delivery

ROOT = Path(__file__).resolve().parents[1]
def load(name):
    spec = importlib.util.spec_from_file_location('release_'+name, ROOT/'scripts'/f'{name}.py')
    module = importlib.util.module_from_spec(spec); spec.loader.exec_module(module)
    return module
def sha(path): return hashlib.sha256(path.read_bytes()).hexdigest()
def write(path, value): path.write_text(json.dumps(value), encoding='utf-8')
def read(path): return json.loads(path.read_text(encoding='utf-8'))

def v63_project(project):
    pptx, generator = legacy_delivery.V6DeliveryTests().make_project(project, 'deconstruct')
    build=project/'.build'
    # Protocol-only fixture, not a claim of visual acceptance of a real deck.
    write(build/'v63_scene_graph.json', {'pages':{'S01':{'elements':[]}}})
    write(build/'v63_visual_census.json', {'pages':{'S01':{'candidates':[]}}})
    asset=build/'assets/S01/map.png';asset.parent.mkdir(parents=True)
    Image.new('RGB',(4,4),'blue').save(asset)
    write(build/'v63_asset_ledger.json', {'assets':[{'asset_id':'MAP','asset_path':'.build/assets/S01/map.png','asset_sha256':sha(asset)}]})
    compile=read(build/'compile_report.json')
    compile.update(schema_version='6.3',deconstruction_runtime_revision='6.3.1',scene_graph_sha256=sha(build/'v63_scene_graph.json'),asset_ledger_sha256=sha(build/'v63_asset_ledger.json'))
    write(build/'compile_report.json',compile)
    acceptance=load('v63_acceptance')
    bindings=acceptance.visual_bindings(project,pptx)
    write(build/'v63_visual_review.json',{'reviewed':True,'bindings':bindings,'pages':[{'slide_id':'S01','object_checks':[],'differences':[]}]})
    write(build/'deconstruction_acceptance.json',{'schema_version':'6.3','deconstruction_runtime_revision':'6.3.1','accepted':True,'status':'pass','pptx_sha256':sha(pptx),'bindings':bindings,'visual_review_sha256':sha(build/'v63_visual_review.json')})
    return pptx,generator

class ReleaseTests(unittest.TestCase):
    def test_packaging_failure_is_recorded_and_retry_does_not_build(self):
        from types import SimpleNamespace
        with tempfile.TemporaryDirectory() as d:
            p=Path(d);pptx,gen=v63_project(p);pipeline=load('project_pipeline')
            result=read(p/'.build/pipeline_result.json')
            result.update(pptx_sha256=sha(pptx),build_attempt=2,status='success')
            original=pptx.read_bytes()
            with patch.object(pipeline,'_load_module',return_value=SimpleNamespace(package_v6_delivery=lambda *a: (_ for _ in ()).throw(ValueError('disk-full')))):
                with self.assertRaisesRegex(ValueError,'disk-full'):
                    pipeline._package_v631_result(p,pptx,result)
            failed=read(p/'.build/pipeline_result.json')
            self.assertFalse(failed['ok'])
            self.assertTrue(failed['pptx_accepted'])
            self.assertEqual('failed',failed['delivery_status'])
            pipeline._package_v631_result(p,pptx,failed)
            saved=read(p/'.build/pipeline_result.json')
            self.assertTrue(saved['ok'])
            self.assertEqual('packaged',saved['delivery_status'])
            self.assertEqual(2,saved['build_attempt'])
            self.assertEqual(original,pptx.read_bytes())

    @unittest.skipUnless(os.name=='nt','native Windows rendering')
    def test_bitmap_preview_failure_recovers_existing_pptx_without_build(self):
        from pptx import Presentation
        from pptx.util import Inches
        with tempfile.TemporaryDirectory() as d:
            p=Path(d);pptx=p/'existing.pptx'
            deck=Presentation();s=deck.slides.add_slide(deck.slide_layouts[6])
            s.shapes.add_textbox(Inches(1),Inches(1),Inches(3),Inches(1)).text='Render only'
            deck.save(pptx);before=sha(pptx)
            pipeline=load('project_pipeline')
            with patch.object(pipeline,'_run',side_effect=subprocess.CalledProcessError(2,['renderer'],stderr='RUNTIME_NODE unavailable')):
                result=pipeline._render_windows_post_lock(p,pptx,1,'bitmap')
            self.assertTrue(result['ok'])
            self.assertEqual('powerpoint_windows_com',result['renderer'])
            self.assertEqual(before,sha(pptx))
            self.assertTrue((p/'.build/rendered/current/S01.png').is_file())
            self.assertFalse((p/'.build/v6_build_attempt.json').exists())

    def test_v63_package_accepts_current_receipt_and_contains_actual_assets(self):
        with tempfile.TemporaryDirectory() as d:
            p=Path(d);pptx,gen=v63_project(p)
            out=load('pack_delivery').package_v6_delivery(p,pptx,gen,p/'delivery.zip')
            with ZipFile(out) as z:
                self.assertEqual({'report.pptx','py.zip','blueprints.zip'},set(z.namelist()))
                with ZipFile(io.BytesIO(z.read('blueprints.zip'))) as inner:
                    self.assertIn('assets/S01/map.png',inner.namelist())
                    self.assertEqual((p/'.build/assets/S01/map.png').read_bytes(),inner.read('assets/S01/map.png'))

    def test_v63_package_still_rejects_changed_scene_or_generator(self):
        for filename in ['v63_scene_graph.json','generate_deck.py']:
            with self.subTest(filename=filename),tempfile.TemporaryDirectory() as d:
                p=Path(d);pptx,gen=v63_project(p)
                path=gen if filename.endswith('.py') else p/'.build'/filename
                path.write_text('{}',encoding='utf-8')
                with self.assertRaises(ValueError):
                    load('pack_delivery').package_v6_delivery(p,pptx,gen,p/'bad.zip')
                self.assertFalse((p/'bad.zip').exists())

    def test_system_python_discovers_installed_node_without_global_env_change(self):
        with tempfile.TemporaryDirectory() as d:
            home=Path(d);dep=home/'.cache/codex-runtimes/codex-primary-runtime/dependencies'
            node=dep/'node/bin'/('node.exe' if os.name=='nt' else 'node')
            node.parent.mkdir(parents=True);node.write_bytes(b'fixture')
            modules=dep/'node/node_modules';modules.mkdir()
            original={'PATH':'unchanged'}
            with patch.object(Path,'home',return_value=home):
                env=load('project_pipeline')._bitmap_render_environment(original,executable=home/'system/python.exe')
            self.assertEqual(str(node),env.get('RUNTIME_NODE'))
            self.assertEqual(str(modules),env.get('RUNTIME_NODE_MODULES'))
            self.assertEqual({'PATH':'unchanged'},original)

    def test_measured_rounding_has_same_small_radius_on_both_backends(self):
        page={'body_roi_px':[0,0,1000,500],'coordinate_mode':'source_pixels_contain','elements':[{
            'element_id':'P','type':'round_rect','bbox_px':[10,10,510,410],'z_order':0,
            'style':{'corner_radius_px':8,'fill':'#FFFFFF'}}]}
        for backend in ['v63_windows_scene_renderer','v63_mac_scene_renderer']:
            with self.subTest(backend=backend):
                command=load(backend).render_plan(page,[0,0,10,5])[0]
                self.assertEqual('freeform',command['type'])
                self.assertTrue(command['closed'])
                self.assertEqual([5.02,.1], [round(v,2) for v in command['points_in'][0]])
                self.assertEqual([.1,.1,5.,4.],command['bbox_in'])

    def test_invalid_measured_radius_is_rejected(self):
        for value in [-1,300,float('nan')]:
            with self.subTest(value=value),self.assertRaises(ValueError):
                load('v63_scene_graph').normalize_element_geometry({'type':'round_rect','bbox_px':[0,0,100,50],'style':{'corner_radius_px':value}})

if __name__=='__main__':unittest.main()
