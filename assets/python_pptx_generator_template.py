"""V5.9 local macOS generator template compiled into one generate_deck.py."""

from __future__ import annotations

import argparse
import hashlib
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DECK_META = {
    "schema_version": "__PROJECT_SCHEMA_VERSION__",
    "production_mode": "__PRODUCTION_MODE__",
    "builder_backend": "mac_python_pptx_v1",
    "template_path": "__COMPANY_TEMPLATE_PATH__",
    "template_sha256": "__COMPANY_TEMPLATE_SHA256__",
    "page_count": 0,  # __PAGE_COUNT__
    "contract_hashes": {},  # __CONTRACT_HASHES__
}
SLIDES = []
DESIGN_DRAFTS = {}
ASSET_CROPS = {}
PAGE_SPECS = {}
TEXT_METRICS_SOURCE = ""
OOXML_ADAPTER_SOURCE = ""
FONTTOOLS_VENDOR_PATH = "__FONTTOOLS_VENDOR_PATH__"
if FONTTOOLS_VENDOR_PATH and FONTTOOLS_VENDOR_PATH not in sys.path:
    sys.path.insert(0, FONTTOOLS_VENDOR_PATH)
exec(TEXT_METRICS_SOURCE, globals())
exec(OOXML_ADAPTER_SOURCE, globals())


FONT_PREFERENCES = ("Microsoft YaHei", "PingFang SC", "Noto Sans CJK SC")
FONT_NAME = FONT_PREFERENCES[0]
FONT_PATH = None
FONT_FACE_INDEX = 0
TEXT_MEASUREMENTS = []
ASSET_INSERTIONS = []
NAVY = "#1E386B"
MID_BLUE = "#3F628F"
BLUE = "#7391B3"
LIGHT_BLUE = "#E6EBF1"
LIGHT_GRAY = "#D9D9D9"
WHITE = "#FFFFFF"
BLACK = "#000000"


def sha256_file(path: str | Path) -> str:
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(str(value).lstrip("#").upper())


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        remove_shape(shape)


def _absolute_box(element: dict, body: dict[str, float]) -> list[float]:
    box = element.get("box")
    if not isinstance(box, list) or len(box) != 4:
        raise ValueError("page element requires box=[x,y,w,h]")
    x, y, width, height = [float(value) for value in box]
    if element.get("coord_space", "body") == "body":
        x += body["x"]
        y += body["y"]
    if width <= 0 or height <= 0:
        raise ValueError("page element width and height must be positive")
    return [x, y, width, height]


def add_textbox(
    slide,
    text: str,
    box: list[float],
    *,
    size: float,
    color: str,
    bold: bool = False,
    name: str | None = None,
    align=PP_ALIGN.LEFT,
):
    x, y, width, height = box
    shape = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.text_frame.clear()
    shape.text_frame.margin_left = Inches(0.04)
    shape.text_frame.margin_right = Inches(0.04)
    shape.text_frame.margin_top = 0
    shape.text_frame.margin_bottom = 0
    shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    set_east_asian_font(run, FONT_NAME)
    if name:
        set_shape_name(shape, name)
    if FONT_PATH is not None and text:
        measurement = measure_text_box(
            str(text),
            font_path=FONT_PATH,
            font_index=FONT_FACE_INDEX,
            resolved_font_name=FONT_NAME,
            font_size_pt=size,
            max_width_pt=max(1.0, (width - 0.08) * 72),
            safety_margin_ratio=1.10,
        )
        TEXT_MEASUREMENTS.append({
            "shape_name": name or shape.name,
            "text": str(text),
            "allocated_height_pt": measurement.allocated_height_pt,
            "box_height_pt": height * 72,
            "fits": measurement.allocated_height_pt <= height * 72,
        })
    return shape


def _add_shape(slide, shape_type, box, *, fill=WHITE, line=LIGHT_GRAY, name=None):
    x, y, width, height = box
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
    if name:
        set_shape_name(shape, name)
    return shape


def add_skeleton(slide, spec: dict, page_number: int) -> dict[str, float]:
    add_textbox(
        slide, spec["chapter"], [0.56, 0.1575, 12.20, 0.47],
        size=20, color=NAVY, bold=True, name="SKEL_CHAPTER",
    )
    _add_shape(
        slide, MSO_SHAPE.RECTANGLE, [0.56, 0.5906, 12.20, 0.394],
        fill=NAVY, line=None, name="SKEL_TITLE_BAR",
    )
    add_textbox(
        slide, spec["title"], [0.60, 0.60, 12.12, 0.34],
        size=16, color=WHITE, bold=True, name="SKEL_TITLE",
    )
    core_text = "\n".join(f"■ {point}" for point in spec.get("core_points", []))
    if FONT_PATH is None:
        raise RuntimeError("generation font was not resolved")
    measurement = measure_text_box(
        core_text,
        font_path=FONT_PATH,
        font_index=FONT_FACE_INDEX,
        resolved_font_name=FONT_NAME,
        font_size_pt=12,
        max_width_pt=12.04 * 72,
        margin_top_pt=4,
        margin_bottom_pt=4,
        safety_margin_ratio=1.10,
    )
    core_height = max(0.64, measurement.allocated_height_pt / 72)
    core = _add_shape(
        slide, MSO_SHAPE.RECTANGLE, [0.56, 1.063, 12.20, core_height],
        fill=WHITE, line=BLACK, name="SKEL_CORE_BORDER",
    )
    core.line.width = Pt(1)
    set_line_dash(core, "dash")
    add_textbox(
        slide, core_text, [0.64, 1.09, 12.04, core_height - 0.05],
        size=12, color=BLACK, name="SKEL_CORE",
    )
    add_textbox(
        slide, spec.get("source", ""), [0.56, 7.218, 10.56, 0.282],
        size=7, color="#666666", name="SKEL_SOURCE",
    )
    add_textbox(
        slide, str(page_number), [12.267, 7.218, 0.50, 0.282],
        size=8, color="#666666", name="SKEL_PAGE_NUMBER", align=PP_ALIGN.RIGHT,
    )
    body_y = 1.063 + core_height + 0.15
    return {"x": 0.56, "y": body_y, "w": 12.20, "h": 7.05 - body_y}


def configure_generation_font(project_dir: Path, font_catalog=None):
    global FONT_NAME, FONT_PATH, FONT_FACE_INDEX
    catalog = font_catalog if font_catalog is not None else build_font_catalog()
    resolved = resolve_font(FONT_PREFERENCES, catalog=catalog)
    FONT_NAME = resolved.name
    FONT_PATH = resolved.path
    FONT_FACE_INDEX = resolved.face_index
    payload = {
        "schema_version": "5.9",
        "measurement_backend": "pillow_font_metrics",
        "requested_order": list(FONT_PREFERENCES),
        "resolved_font": resolved.name,
        "font_path": str(resolved.path),
        "face_index": resolved.face_index,
        "fallbacks": ([{
            "requested": FONT_PREFERENCES[0],
            "resolved": resolved.name,
            "reason": "preferred font unavailable",
        }] if resolved.fallback_used else []),
    }
    report = project_dir / ".build" / "font_report.json"
    report.parent.mkdir(parents=True, exist_ok=True)
    report.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return resolved


def add_section_header_mac(slide, element: dict, box: list[float]):
    return add_textbox(
        slide, str(element.get("text", element.get("title", ""))), box,
        size=float(element.get("font_size", 10)),
        color=str(element.get("color", NAVY)), bold=True,
    )


def add_metric_strip_mac(slide, element: dict, box: list[float]):
    metrics = element.get("metrics", element.get("data", []))
    x, y, width, height = box
    cell_width = width / max(1, len(metrics))
    for index, metric in enumerate(metrics):
        left = x + index * cell_width
        add_textbox(
            slide, str(metric.get("value", "")),
            [left, y, cell_width, height * 0.58],
            size=15, color=str(metric.get("color", NAVY)), bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, str(metric.get("label", "")),
            [left, y + height * 0.58, cell_width, height * 0.42],
            size=8, color="#666666", align=PP_ALIGN.CENTER,
        )


def add_flow_mac(slide, element: dict, box: list[float]):
    steps = element.get("steps", [])
    x, y, width, height = box
    gap = 0.18
    step_width = (width - gap * max(0, len(steps) - 1)) / max(1, len(steps))
    for index, raw_step in enumerate(steps):
        step = raw_step if isinstance(raw_step, dict) else {"label": raw_step}
        left = x + index * (step_width + gap)
        _add_shape(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            [left, y, step_width, height],
            fill=str(step.get("fill", "#F5F5F5")),
            line=str(step.get("line_color", BLUE)),
        )
        add_textbox(
            slide, str(step.get("title", step.get("label", ""))),
            [left + 0.05, y + 0.05, max(0.05, step_width - 0.1), max(0.05, height - 0.1)],
            size=9, color=BLACK, bold=True, align=PP_ALIGN.CENTER,
        )
        if index:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(left - gap), Inches(y + height / 2),
                Inches(left), Inches(y + height / 2),
            )
            set_arrow_end(connector)


def _chart_data(element: dict) -> CategoryChartData:
    rows = element.get("data", [])
    data = CategoryChartData()
    data.categories = [str(row.get("label", "")) for row in rows]
    if rows and isinstance(rows[0].get("values"), list):
        series_defs = element.get("series", [])
        count = max(len(row.get("values", [])) for row in rows)
        for index in range(count):
            name = (
                str(series_defs[index].get("name", f"Series {index + 1}"))
                if index < len(series_defs) and isinstance(series_defs[index], dict)
                else f"Series {index + 1}"
            )
            data.add_series(
                name,
                [
                    (row.get("values", []) + [None] * count)[index]
                    for row in rows
                ],
            )
    else:
        data.add_series(
            str(element.get("series_name", "Value")),
            [row.get("value", 0) for row in rows],
        )
    return data


def add_category_chart(slide, kind: str, element: dict, box: list[float]):
    chart_types = {
        "hbar_chart": XL_CHART_TYPE.BAR_CLUSTERED,
        "column_chart": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line_chart": XL_CHART_TYPE.LINE_MARKERS,
        "donut_chart": XL_CHART_TYPE.DOUGHNUT,
        "grouped_hbar_chart": XL_CHART_TYPE.BAR_CLUSTERED,
    }
    x, y, width, height = box
    frame = slide.shapes.add_chart(
        chart_types[kind],
        Inches(x), Inches(y), Inches(width), Inches(height),
        _chart_data(element),
    )
    chart = frame.chart
    chart.has_title = False
    chart.chart_style = 10
    chart.has_legend = bool(element.get("show_legend", len(chart.series) > 1))
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    set_shape_name(frame, f"CHART_{kind.upper()}")
    return frame


def add_combo_chart(slide, element: dict, box: list[float]):
    columns = dict(element)
    columns["data"] = [
        {"label": row.get("label", ""), "value": row.get("column", row.get("value", 0))}
        for row in element.get("data", [])
    ]
    lines = dict(element)
    lines["data"] = [
        {"label": row.get("label", ""), "value": row.get("line", row.get("value2", 0))}
        for row in element.get("data", [])
    ]
    column_frame = add_category_chart(slide, "column_chart", columns, box)
    inset = [box[0] + box[2] * 0.52, box[1] + 0.08, box[2] * 0.46, box[3] * 0.42]
    line_frame = add_category_chart(slide, "line_chart", lines, inset)
    set_shape_name(column_frame, "CHART_COMBO_COLUMN")
    set_shape_name(line_frame, "CHART_COMBO_LINE")
    return column_frame, line_frame


def _format_table_text(table) -> None:
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = FONT_NAME
                    run.font.size = Pt(8)
                    set_east_asian_font(run, FONT_NAME)


def render_page_spec(
    slide, page_spec: dict, body: dict[str, float], project_dir: Path, slide_id: str
) -> None:
    elements = page_spec.get("elements")
    if not isinstance(elements, list):
        raise ValueError(f"{slide_id}: page spec elements must be a list")
    for index, element in enumerate(elements):
        if not isinstance(element, dict):
            raise ValueError(f"{slide_id}/element[{index}]: must be a mapping")
        kind = element.get("type")
        absolute = _absolute_box(element, body)
        x, y, width, height = absolute
        if kind == "section_header":
            add_section_header_mac(slide, element, absolute)
        elif kind == "text":
            add_textbox(
                slide, str(element.get("text", "")), absolute,
                size=float(element.get("font_size", 10)),
                color=str(element.get("color", BLACK)),
                bold=bool(element.get("bold", False)),
            )
        elif kind in {"rect", "oval"}:
            _add_shape(
                slide,
                MSO_SHAPE.RECTANGLE if kind == "rect" else MSO_SHAPE.OVAL,
                absolute,
                fill=str(element.get("fill", WHITE)),
                line=str(element.get("line", element.get("line_color", NAVY))),
            )
        elif kind in {"line", "arrow"}:
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x), Inches(y), Inches(x + width), Inches(y + height),
            )
            line.line.color.rgb = rgb(str(element.get("color", NAVY)))
            line.line.width = Pt(float(element.get("weight", 1.0)))
            if kind == "arrow":
                set_arrow_end(line)
        elif kind == "asset":
            asset_id = str(element["asset_id"])
            record = ASSET_CROPS.get(asset_id, {})
            path = (
                project_dir / ".build" / "assets"
                / str(record.get("slide_id", slide_id)) / f"{asset_id}.png"
            )
            if not path.is_file():
                raise FileNotFoundError(
                    f"{slide_id}: missing extracted asset {asset_id}: {path}"
                )
            picture = slide.shapes.add_picture(
                str(path), Inches(x), Inches(y), Inches(width), Inches(height)
            )
            shape_name = f"ASSET_{asset_id}"
            set_shape_name(picture, shape_name)
            ASSET_INSERTIONS.append({
                "slide_id": slide_id,
                "asset_id": asset_id,
                "shape_name": shape_name,
                "source_path": str(path),
                "source_sha256": sha256_file(path),
                "target_box_in": absolute,
            })
        elif kind == "text_card":
            _add_shape(
                slide, MSO_SHAPE.RECTANGLE, absolute,
                fill=str(element.get("body_fill", WHITE)), line=LIGHT_GRAY,
            )
            add_textbox(
                slide, str(element.get("title", "")),
                [x + 0.08, y + 0.05, width - 0.16, min(0.28, height * 0.30)],
                size=10, color=NAVY, bold=True,
            )
            add_textbox(
                slide, str(element.get("body", "")),
                [x + 0.08, y + min(0.34, height * 0.35), width - 0.16, max(0.05, height - 0.39)],
                size=9, color=BLACK,
            )
        elif kind == "metric_strip":
            add_metric_strip_mac(slide, element, absolute)
        elif kind in {
            "hbar_chart", "column_chart", "line_chart",
            "donut_chart", "grouped_hbar_chart",
        }:
            add_category_chart(slide, kind, element, absolute)
        elif kind == "combo_chart":
            add_combo_chart(slide, element, absolute)
        elif kind == "flow":
            add_flow_mac(slide, element, absolute)
        elif kind == "matrix":
            headers = element.get("headers", [])
            rows = element.get("rows", [])
            if not headers:
                raise ValueError(f"{slide_id}: matrix requires headers")
            table_shape = slide.shapes.add_table(
                len(rows) + 1, len(headers),
                Inches(x), Inches(y), Inches(width), Inches(height),
            )
            table = table_shape.table
            for column, value in enumerate(headers):
                table.cell(0, column).text = str(value)
            for row_index, row in enumerate(rows, start=1):
                for column, value in enumerate(row[:len(headers)]):
                    table.cell(row_index, column).text = str(value)
            _format_table_text(table)
        else:
            raise ValueError(
                f"{slide_id}/element[{index}]: unsupported Mac element type {kind!r}"
            )


# __PAGE_BUILDERS__


def validate_embedded_contract() -> None:
    expected = [f"S{index:02d}" for index in range(1, DECK_META["page_count"] + 1)]
    if [slide.get("slide_id") for slide in SLIDES] != expected:
        raise ValueError("embedded slides are not in canonical order")
    if sorted(PAGE_SPECS) != expected or sorted(PAGE_BUILDERS) != expected:
        raise ValueError("embedded page specs/builders do not cover every slide")
    if DECK_META["production_mode"] == "blueprint" and sorted(DESIGN_DRAFTS) != expected:
        raise ValueError("blueprint mode requires one locked design draft per slide")


def build_deck(output_path: str | Path | None = None, *, font_catalog=None) -> Path:
    TEXT_MEASUREMENTS.clear()
    ASSET_INSERTIONS.clear()
    validate_embedded_contract()
    project_dir = Path(__file__).resolve().parent
    output = Path(output_path) if output_path else project_dir / "output" / "report.pptx"
    output = output.resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    template = Path(DECK_META["template_path"])
    if not template.is_file() or sha256_file(template) != DECK_META["template_sha256"]:
        raise ValueError("company template file/hash mismatch")
    configure_generation_font(project_dir, font_catalog)
    presentation = Presentation(template)
    if not presentation.slides:
        raise ValueError("company template must contain its initial slide")
    first_layout = presentation.slides[0].slide_layout
    while len(presentation.slides) < DECK_META["page_count"]:
        presentation.slides.add_slide(first_layout)
    while len(presentation.slides) > DECK_META["page_count"]:
        remove_last_slide(presentation)
    for page_number, spec in enumerate(SLIDES, start=1):
        slide = presentation.slides[page_number - 1]
        clear_slide(slide)
        body = add_skeleton(slide, spec, page_number)
        PAGE_BUILDERS[spec["slide_id"]](slide, spec, body, project_dir)
        clear_forbidden_effects(slide.element)
    for master in presentation.slide_masters:
        clear_forbidden_effects(master.element)
        for layout in master.slide_layouts:
            clear_forbidden_effects(layout.element)
    temporary = output.with_name(f".{output.stem}.building.pptx")
    if temporary.exists():
        temporary.unlink()
    presentation.save(temporary)
    validate_pptx_package(temporary)
    temporary.replace(output)
    metrics_report = project_dir / ".build" / "mac_text_metrics_report.json"
    metrics_report.write_text(
        json.dumps({
            "schema_version": "5.9",
            "measurement_backend": "pillow_font_metrics",
            "items": TEXT_MEASUREMENTS,
            "ok": all(item["fits"] for item in TEXT_MEASUREMENTS),
        }, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    asset_report = project_dir / ".build" / "mac_asset_report.json"
    asset_report.write_text(
        json.dumps({
            "schema_version": "5.9",
            "assets": ASSET_INSERTIONS,
        }, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    return output


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Build a V5.9 deck locally with python-pptx."
    )
    parser.add_argument("--output", type=Path)
    args = parser.parse_args()
    print(build_deck(args.output))


if __name__ == "__main__":
    main()
